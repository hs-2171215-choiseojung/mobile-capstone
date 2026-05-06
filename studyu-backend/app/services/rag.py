"""
RAG 파이프라인 (Supabase 영속 저장, pdfplumber + OpenAI)

문서 청크는 Supabase document_chunks 테이블에 저장되어
서버 재시작 / 배포 후에도 데이터가 유지됩니다.

임베딩 기반 유사도 검색:
  - 인덱싱 시 text-embedding-3-small로 청크 임베딩 생성 후 JSONB 컬럼에 저장
  - 쿼리 시 질문 임베딩과 코사인 유사도로 top-k 청크만 선택
  - document_chunks 테이블에 embedding JSONB 컬럼이 있어야 합니다:
      ALTER TABLE document_chunks ADD COLUMN IF NOT EXISTS embedding JSONB;
"""

import io
import json
import re
from typing import Any
from concurrent.futures import ThreadPoolExecutor, as_completed
import numpy as np
import pdfplumber
from openai import OpenAI
from app.core.config import settings
from app.core.supabase import supabase_admin

CHUNK_SIZE = 900
CHUNK_OVERLAP = 100
TOP_K = 5  # 문서당 검색할 최대 청크 수

# 이미지를 직접 가리키는 질문 패턴 — Python 레벨 사전 판단용
_IMAGE_DIRECT_REF_RE = re.compile(
    r'(이|저|그)\s+[가-힣]+'                      # 이 개, 저 꽃, 그 사람
    r'|이\s*(사진|이미지|그림|첨부|파일)'           # 이 사진, 이 이미지
    r'|사진\s*(속|에서|의|에|안)'                  # 사진 속, 사진에서
    r'|이미지\s*(속|에서|의|에|안)'
    r'|여기\s*(서|에|에서|나오는|에\s*있는)'
    r'|이것|저것|그것'
    r'|털\s*색|색상|색깔|무슨\s*색|어떤\s*색'      # 시각적 속성
    r'|외형|생김새|외모|모습|모양새',
    re.UNICODE,
)

_UNRELATED_IMG_PREFIX = "📌 이 질문은 업로드된 이미지와 직접적인 관련은 없지만, 알고 계시면 도움이 될 것 같아 답변드립니다."

_MEDIA_META_RE = re.compile(r"\[MEDIA_META\](.*?)\[/MEDIA_META\]", re.DOTALL)
_MEDIA_SUMMARY_RE = re.compile(r"\[MEDIA_SUMMARY\](.*?)\[/MEDIA_SUMMARY\]", re.DOTALL)
_MEDIA_TRANSCRIPT_HEADER_RE = re.compile(r"^\[(?:음성 전사 내용 - .*?|MEDIA_TRANSCRIPT)\]\s*", re.MULTILINE)


# ─────────────────────────────────────────────
# 멀티 LLM 헬퍼
# ─────────────────────────────────────────────

def _is_claude(model: str) -> bool:
    return model.startswith("claude-")

def _convert_messages_for_anthropic(messages: list[dict]) -> tuple[str, list[dict]]:
    """OpenAI 형식 messages → Anthropic 형식 (system 분리, 이미지 포맷 변환)"""
    system = ""
    anthropic_msgs: list[dict] = []
    for msg in messages:
        if msg["role"] == "system":
            system = msg["content"] if isinstance(msg["content"], str) else ""
            continue
        content = msg["content"]
        if isinstance(content, list):
            new_parts: list[dict] = []
            for part in content:
                if part.get("type") == "text":
                    new_parts.append({"type": "text", "text": part["text"]})
                elif part.get("type") == "image_url":
                    url = part["image_url"]["url"]
                    if url.startswith("data:"):
                        header, b64data = url.split(",", 1)
                        media_type = header.split(":")[1].split(";")[0]
                        new_parts.append({
                            "type": "image",
                            "source": {"type": "base64", "media_type": media_type, "data": b64data},
                        })
            anthropic_msgs.append({"role": msg["role"], "content": new_parts})
        else:
            anthropic_msgs.append({"role": msg["role"], "content": content})
    return system, anthropic_msgs


# ─────────────────────────────────────────────
# 통합 LLM 호출 유틸리티
# ─────────────────────────────────────────────

_DEFAULT_MODEL = "claude-haiku-4-5-20251001"

# sentence-transformers 모델 싱글톤 (첫 호출 시 다운로드)
import threading as _threading
_embedding_model = None
_embedding_lock = _threading.Lock()

def _get_embedding_model():
    global _embedding_model
    if _embedding_model is None:
        with _embedding_lock:
            if _embedding_model is None:  # double-checked locking
                from sentence_transformers import SentenceTransformer
                _embedding_model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")
    return _embedding_model


def _call_llm(
    messages: list[dict],
    model: str = _DEFAULT_MODEL,
    temperature: float = 0.7,
    max_tokens: int = 2048,
    json_mode: bool = False,
) -> str:
    """OpenAI/Claude 자동 라우팅 LLM 호출. model prefix로 분기."""
    if _is_claude(model):
        from anthropic import Anthropic as _Anthropic
        client = _Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        system, anthropic_msgs = _convert_messages_for_anthropic(messages)
        if json_mode:
            system = (system + "\n\n반드시 유효한 JSON 객체만 응답하세요. 마크다운 코드블록 없이 JSON만 출력하세요.").strip()
        kwargs: dict = {
            "model": model,
            "max_tokens": max_tokens,
            "temperature": temperature,
            "messages": anthropic_msgs,
        }
        if system:
            kwargs["system"] = system
        resp = client.messages.create(**kwargs)
        return resp.content[0].text
    else:
        oai = OpenAI(api_key=settings.OPENAI_API_KEY)
        create_kwargs: dict = {
            "model": model,
            "messages": messages,
            "temperature": temperature,
            "max_tokens": max_tokens,
        }
        if json_mode:
            create_kwargs["response_format"] = {"type": "json_object"}
        resp = oai.chat.completions.create(**create_kwargs)
        return resp.choices[0].message.content or ""


def _call_llm_vision(
    image_b64: str,
    mime_type: str,
    prompt: str,
    model: str = "claude-sonnet-4-6",
    json_mode: bool = False,
    max_tokens: int = 2000,
) -> str:
    """이미지 포함 LLM Vision 호출. OpenAI/Claude 자동 라우팅."""
    if _is_claude(model):
        from anthropic import Anthropic as _Anthropic
        client = _Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        content = [
            {"type": "image", "source": {"type": "base64", "media_type": mime_type, "data": image_b64}},
            {"type": "text", "text": prompt},
        ]
        sys_prompt = "반드시 유효한 JSON 객체만 응답하세요. 마크다운 코드블록 없이 JSON만 출력하세요." if json_mode else None
        kwargs: dict = {
            "model": model,
            "max_tokens": max_tokens,
            "messages": [{"role": "user", "content": content}],
        }
        if sys_prompt:
            kwargs["system"] = sys_prompt
        resp = client.messages.create(**kwargs)
        return resp.content[0].text
    else:
        safe_vision_model = model if "gpt-4" in model else "gpt-4o"
        oai = OpenAI(api_key=settings.OPENAI_API_KEY)
        create_kwargs: dict = {
            "model": safe_vision_model,
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_b64}", "detail": "high"}},
                ],
            }],
            "max_tokens": max_tokens,
        }
        if json_mode:
            create_kwargs["response_format"] = {"type": "json_object"}
        resp = oai.chat.completions.create(**create_kwargs)
        return resp.choices[0].message.content or ""


def _elevenlabs_tts_simple(text: str, voice_name: str = "sarah") -> bytes:
    """동기 ElevenLabs TTS (오디오 오버뷰용 간단 호출)."""
    import httpx as _httpx
    from app.services.tts import ELEVENLABS_VOICES
    voice_id = ELEVENLABS_VOICES.get(voice_name, ELEVENLABS_VOICES["sarah"])
    resp = _httpx.post(
        f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}",
        json={
            "text": text,
            "model_id": "eleven_multilingual_v2",
            "voice_settings": {"stability": 0.45, "similarity_boost": 0.80},
        },
        headers={"xi-api-key": settings.ELEVENLABS_API_KEY, "Content-Type": "application/json"},
        timeout=60,
    )
    resp.raise_for_status()
    return resp.content


# ─────────────────────────────────────────────
# 텍스트 정제
# ─────────────────────────────────────────────

def _sanitize(text: str) -> str:
    text = text.translate({0: None})
    text = re.sub(r'[\x01-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    text = text.encode('utf-8', errors='ignore').decode('utf-8')
    text = text.translate({0: None})
    return text


def _hard_sanitize(text: str) -> str:
    return "".join(ch for ch in text if ord(ch) != 0)


def _format_media_timestamp(total_seconds: int | float) -> str:
    safe = max(0, int(total_seconds or 0))
    hours = safe // 3600
    minutes = (safe % 3600) // 60
    seconds = safe % 60
    if hours > 0:
        return f"{hours}:{minutes:02d}:{seconds:02d}"
    return f"{minutes}:{seconds:02d}"


_WHISPER_SUPPORTED_MEDIA_EXTENSIONS = {"flac", "m4a", "mp3", "mp4", "mpeg", "mpga", "oga", "ogg", "wav", "webm"}


def _find_ffmpeg_binary() -> str | None:
    import shutil

    return shutil.which("ffmpeg")


def _prepare_media_file_for_transcription(input_path: str, original_ext: str) -> tuple[str, bool]:
    import os
    import subprocess
    import tempfile

    normalized_ext = (original_ext or "").lower().lstrip(".")
    if normalized_ext in _WHISPER_SUPPORTED_MEDIA_EXTENSIONS:
        return input_path, False

    ffmpeg_path = _find_ffmpeg_binary()
    if not ffmpeg_path:
        raise ValueError(
            f".{normalized_ext} 형식은 Whisper가 직접 지원하지 않습니다. "
            "서버에 ffmpeg를 설치하면 자동으로 wav로 변환해서 처리할 수 있습니다."
        )

    converted_fd, converted_path = tempfile.mkstemp(suffix=".wav")
    os.close(converted_fd)

    try:
        completed = subprocess.run(
            [
                ffmpeg_path,
                "-y",
                "-i",
                input_path,
                "-vn",
                "-ac",
                "1",
                "-ar",
                "16000",
                converted_path,
            ],
            check=False,
            capture_output=True,
            text=True,
        )
        if completed.returncode != 0:
            stderr = (completed.stderr or "").strip()
            raise ValueError(
                f".{normalized_ext} 파일을 전사용 wav로 변환하지 못했습니다."
                + (f" ffmpeg 오류: {stderr[:240]}" if stderr else "")
            )
        return converted_path, True
    except Exception:
        if os.path.exists(converted_path):
            os.unlink(converted_path)
        raise


def prepare_media_for_browser_playback(file_bytes: bytes, filename: str) -> tuple[bytes, str, str]:
    import os
    import subprocess
    import tempfile

    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext not in {"avi", "mkv", "mov"}:
        return file_bytes, ext, ""

    ffmpeg_path = _find_ffmpeg_binary()
    if not ffmpeg_path:
        return file_bytes, ext, ""

    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as src_tmp:
        src_tmp.write(file_bytes)
        src_path = src_tmp.name

    out_fd, out_path = tempfile.mkstemp(suffix=".mp4")
    os.close(out_fd)

    try:
        completed = subprocess.run(
            [
                ffmpeg_path,
                "-y",
                "-i",
                src_path,
                "-c:v",
                "libx264",
                "-preset",
                "veryfast",
                "-pix_fmt",
                "yuv420p",
                "-movflags",
                "+faststart",
                "-c:a",
                "aac",
                "-b:a",
                "128k",
                out_path,
            ],
            check=False,
            capture_output=True,
            text=True,
        )
        if completed.returncode != 0:
            return file_bytes, ext, ""
        converted_bytes = open(out_path, "rb").read()
        if not converted_bytes:
            return file_bytes, ext, ""
        return converted_bytes, "mp4", "video/mp4"
    finally:
        if os.path.exists(src_path):
            os.unlink(src_path)
        if os.path.exists(out_path):
            os.unlink(out_path)


def _looks_like_music_filename(filename: str) -> bool:
    lower = (filename or "").lower()
    keywords = ["가사", "lyric", "lyrics", "music", "뮤직", "ost", "mv", "official audio", "audio"]
    return any(keyword in lower for keyword in keywords)


def _is_transcript_credit_noise(text: str) -> bool:
    normalized = " ".join((text or "").strip().lower().split())
    if not normalized:
        return True
    credit_patterns = [
        "한글자막 by",
        "subtitle by",
        "subtitles by",
        "lyrics by",
        "lyric by",
    ]
    return any(pattern in normalized for pattern in credit_patterns)


def _extract_words_from_segment(segment: Any) -> list[dict[str, Any]]:
    raw_words = getattr(segment, "words", None)
    if not raw_words:
        return []

    words: list[dict[str, Any]] = []
    for word in raw_words:
        text = _hard_sanitize(_sanitize(getattr(word, "word", "") or "")).strip()
        if not text:
            continue
        start_sec = getattr(word, "start", None)
        end_sec = getattr(word, "end", None)
        if not isinstance(start_sec, (int, float)) or not isinstance(end_sec, (int, float)):
            continue
        words.append({
            "word": text,
            "start_sec": float(start_sec),
            "end_sec": float(max(start_sec, end_sec)),
        })
    return words


def _refine_media_segments(segments: list[dict[str, Any]], *, is_music: bool) -> list[dict[str, Any]]:
    refined: list[dict[str, Any]] = []
    for segment in segments:
        text = _hard_sanitize(_sanitize(segment.get("text", ""))).strip()
        if not text or _is_transcript_credit_noise(text):
            continue

        start_sec = float(segment.get("start_sec", 0) or 0)
        end_sec = float(max(start_sec, segment.get("end_sec", start_sec) or start_sec))
        words = segment.get("words") if isinstance(segment.get("words"), list) else []

        valid_words = [word for word in words if isinstance(word, dict)]
        if valid_words:
            start_candidates = [float(word["start_sec"]) for word in valid_words if isinstance(word.get("start_sec"), (int, float))]
            end_candidates = [float(word["end_sec"]) for word in valid_words if isinstance(word.get("end_sec"), (int, float))]
            if start_candidates:
                start_sec = min(start_candidates)
            if end_candidates:
                end_sec = max(end_candidates)

        refined.append({
            "start_sec": start_sec,
            "end_sec": max(start_sec, end_sec),
            "text": text,
        })

    if not refined:
        return []

    if is_music:
        if refined[0]["start_sec"] < 1.0:
            refined[0]["start_sec"] = min(max(refined[0]["start_sec"], 0.0) + 1.0, refined[0]["end_sec"])

        for idx in range(1, len(refined)):
            prev = refined[idx - 1]
            current = refined[idx]
            # Keep segments monotonic and reduce overlaps that make music lyrics look early.
            if current["start_sec"] < prev["end_sec"]:
                midpoint = max(prev["start_sec"], min(prev["end_sec"], (prev["end_sec"] + current["start_sec"]) / 2))
                prev["end_sec"] = midpoint
                current["start_sec"] = max(midpoint, current["start_sec"])
            current["start_sec"] = max(current["start_sec"], prev["start_sec"])
            current["end_sec"] = max(current["start_sec"], current["end_sec"])

    return refined


def _parse_media_metadata(content: str) -> dict[str, Any]:
    match = _MEDIA_META_RE.search(content or "")
    if not match:
        return {}
    try:
        parsed = json.loads(match.group(1))
        return parsed if isinstance(parsed, dict) else {}
    except Exception:
        return {}


def _normalize_media_summary_payload(payload: Any) -> dict[str, Any]:
    if isinstance(payload, list):
        sections = [item for item in payload if isinstance(item, dict)]
        return {
            "title": "전체 내용 요약",
            "overview": "",
            "sections": sections,
        }

    if isinstance(payload, dict):
        title = str(payload.get("title", "")).strip() or "전체 내용 요약"
        overview = str(payload.get("overview", "")).strip()
        raw_sections = payload.get("sections", [])
        sections = [item for item in raw_sections if isinstance(item, dict)] if isinstance(raw_sections, list) else []
        return {
            "title": title,
            "overview": overview,
            "sections": sections,
        }

    return {"title": "전체 내용 요약", "overview": "", "sections": []}


def _parse_media_summary(content: str) -> dict[str, Any]:
    match = _MEDIA_SUMMARY_RE.search(content or "")
    if not match:
        return {"title": "전체 내용 요약", "overview": "", "sections": []}
    try:
        parsed = json.loads(match.group(1))
        return _normalize_media_summary_payload(parsed)
    except Exception:
        return {"title": "전체 내용 요약", "overview": "", "sections": []}


def _strip_media_metadata(content: str) -> str:
    cleaned = _MEDIA_META_RE.sub("", content or "")
    cleaned = _MEDIA_SUMMARY_RE.sub("", cleaned)
    cleaned = _MEDIA_TRANSCRIPT_HEADER_RE.sub("", cleaned)
    return cleaned.strip()


def _inject_media_summary(chunks: list[str], summary: dict[str, Any]) -> list[str]:
    if not chunks or not summary or not summary.get("sections"):
        return chunks
    encoded = json.dumps(summary, ensure_ascii=False)
    cleaned_first = _MEDIA_SUMMARY_RE.sub("", chunks[0]).strip()
    chunks[0] = f"[MEDIA_SUMMARY]{encoded}[/MEDIA_SUMMARY]\n{cleaned_first}".strip()
    return chunks


def _build_media_timeline_from_chunks_v2(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    timeline: list[dict[str, Any]] = []
    for row in rows or []:
        content = row.get("content", "")
        text = _strip_media_metadata(content)
        if not text:
            continue
        meta = _parse_media_metadata(content)
        start_sec = meta.get("start_sec")
        if not isinstance(start_sec, (int, float)):
            start_sec = 0
        end_sec = meta.get("end_sec")
        entry: dict[str, Any] = {
            "time_sec": int(max(0, start_sec)),
            "text": text,
        }
        if isinstance(end_sec, (int, float)):
            entry["end_sec"] = int(max(entry["time_sec"], end_sec))
        timeline.append(entry)
    return timeline


_MEDIA_ALIGNMENT_TOKEN_RE = re.compile(r"[\uac00-\ud7a3]+|[a-z0-9]{2,}")


def _tokenize_media_alignment_text(text: str) -> list[str]:
    normalized = re.sub(r"\s+", " ", (text or "").lower()).strip()
    if not normalized:
        return []
    return _MEDIA_ALIGNMENT_TOKEN_RE.findall(normalized)


def _build_media_alignment_bigrams(tokens: list[str]) -> set[str]:
    if len(tokens) < 2:
        return set()
    return {f"{tokens[idx]} {tokens[idx + 1]}" for idx in range(len(tokens) - 1)}


def _extract_media_lead_query(title: str, summary: str) -> str:
    cleaned_summary = re.sub(r"\r\n?", "\n", summary or "").strip()
    summary_lines = [line.strip() for line in cleaned_summary.split("\n") if line.strip()]

    lead_part = ""
    if summary_lines:
        first_line = re.sub(r"^\s*(?:[-*+]|\d+\.)\s*", "", summary_lines[0]).strip()
        sentence_parts = re.split(r"(?<=[.!?])\s+|(?<=다)\s+|(?<=요)\s+", first_line)
        lead_part = next((part.strip() for part in sentence_parts if part.strip()), first_line)

    query_parts = [part for part in [title.strip(), lead_part] if part]
    return "\n".join(query_parts).strip()


def _score_media_alignment(
    section_tokens: list[str],
    section_bigrams: set[str],
    candidate_tokens: list[str],
    candidate_bigrams: set[str],
) -> float:
    if not section_tokens or not candidate_tokens:
        return 0.0

    section_counts: dict[str, int] = {}
    for token in section_tokens:
        section_counts[token] = section_counts.get(token, 0) + 1

    candidate_counts: dict[str, int] = {}
    for token in candidate_tokens:
        candidate_counts[token] = candidate_counts.get(token, 0) + 1

    overlap = 0
    for token, count in section_counts.items():
        overlap += min(count, candidate_counts.get(token, 0))

    token_recall = overlap / max(len(section_tokens), 1)
    token_precision = overlap / max(len(candidate_tokens), 1)
    bigram_overlap = len(section_bigrams & candidate_bigrams)
    bigram_recall = bigram_overlap / max(len(section_bigrams), 1) if section_bigrams else 0.0

    return (token_recall * 0.72) + (token_precision * 0.18) + (bigram_recall * 0.35)


def align_media_summary_to_timeline(
    summary: dict[str, Any],
    timeline: list[dict[str, Any]],
) -> dict[str, Any]:
    sections = summary.get("sections")
    if not isinstance(sections, list) or not sections or not timeline:
        return summary

    aligned_sections: list[dict[str, Any]] = []
    cursor = 0
    total_sections = len(sections)
    min_gap_seconds = 3
    top_k_candidates = 4

    timeline_texts = [str(item.get("text", "")).strip() for item in timeline]
    timeline_tokens = [_tokenize_media_alignment_text(text) for text in timeline_texts]
    timeline_bigrams = [_build_media_alignment_bigrams(tokens) for tokens in timeline_tokens]

    def _normalize_alignment_text(value: str) -> str:
        normalized = re.sub(r"\s+", " ", (value or "").strip().lower())
        normalized = re.sub(r"[^\w\s가-힣]", "", normalized)
        return normalized

    query_texts: list[str] = []
    for section in sections:
        if not isinstance(section, dict):
            query_texts.append("요약 섹션")
            continue
        title = str(section.get("title", "")).strip()
        section_summary = str(section.get("summary", "")).strip()
        anchor_text = str(section.get("anchor_text", "")).strip()
        lead_query_text = anchor_text or _extract_media_lead_query(title, section_summary)
        query_texts.append(lead_query_text or f"{title}\n{section_summary}".strip() or "요약 섹션")

    query_embeddings: list[list[float]] | None = None
    timeline_embeddings: list[list[float]] | None = None
    try:
        if query_texts and timeline_texts:
            query_embeddings = _embed_batch(query_texts)
            timeline_embeddings = _embed_batch([text or "전사 청크" for text in timeline_texts])
    except Exception:
        query_embeddings = None
        timeline_embeddings = None

    for index, section in enumerate(sections):
        if not isinstance(section, dict):
            continue

        title = str(section.get("title", "")).strip()
        section_summary = str(section.get("summary", "")).strip()
        anchor_text = str(section.get("anchor_text", "")).strip()
        lead_query_text = query_texts[index]
        section_tokens = _tokenize_media_alignment_text(lead_query_text)
        section_bigrams = _build_media_alignment_bigrams(section_tokens)

        original_start = section.get("start_sec")
        try:
            original_start_sec = max(0, int(float(original_start)))
        except Exception:
            original_start_sec = None

        best_idx: int | None = None
        best_score = -1.0
        remaining_sections = total_sections - index - 1
        latest_start_idx = max(cursor, len(timeline) - remaining_sections - 1)
        min_window_idx = int((len(timeline) * index) / max(total_sections, 1))
        max_window_idx = int((len(timeline) * (index + 1)) / max(total_sections, 1))
        window_padding = max(2, len(timeline) // max(total_sections * 2, 1))
        search_start = max(cursor, min_window_idx - window_padding)
        search_end = min(latest_start_idx + 1, max_window_idx + window_padding + 1)
        if search_start >= search_end:
            search_start = cursor
            search_end = latest_start_idx + 1

        previous_start_sec = None
        if aligned_sections:
            previous_start_sec = int(aligned_sections[-1].get("start_sec", 0) or 0)

        candidate_scores: list[tuple[int, float]] = []
        embedding_scores: np.ndarray | None = None
        if (
            query_embeddings is not None
            and timeline_embeddings is not None
            and index < len(query_embeddings)
            and len(timeline_embeddings) == len(timeline)
        ):
            embedding_scores = _cosine_similarity(query_embeddings[index], timeline_embeddings)

        exact_anchor_idx: int | None = None
        normalized_anchor = _normalize_alignment_text(anchor_text)
        if len(normalized_anchor) >= 8:
            for chunk_idx in range(search_start, search_end):
                normalized_timeline = _normalize_alignment_text(timeline_texts[chunk_idx])
                if normalized_anchor and normalized_anchor in normalized_timeline:
                    exact_anchor_idx = chunk_idx
                    break

        for chunk_idx in range(search_start, search_end):
            score = _score_media_alignment(
                section_tokens,
                section_bigrams,
                timeline_tokens[chunk_idx],
                timeline_bigrams[chunk_idx],
            ) * 0.45

            if exact_anchor_idx is not None:
                if chunk_idx == exact_anchor_idx:
                    score += 1.5
                else:
                    score -= 0.4

            if embedding_scores is not None and chunk_idx < len(embedding_scores):
                score += float(embedding_scores[chunk_idx]) * 0.55

            candidate_time = int(timeline[chunk_idx].get("time_sec", 0) or 0)
            if original_start_sec is not None:
                drift = abs(candidate_time - original_start_sec)
                score += max(0.0, 1.0 - (drift / 90.0)) * 0.05

            expected_center_idx = int(((index + 0.5) / max(total_sections, 1)) * len(timeline))
            index_distance = abs(chunk_idx - expected_center_idx)
            section_span = max(len(timeline) / max(total_sections, 1), 1)
            score += max(0.0, 1.0 - (index_distance / section_span)) * 0.06

            if previous_start_sec is not None:
                if candidate_time < previous_start_sec:
                    score -= 1.1
                elif candidate_time - previous_start_sec < min_gap_seconds:
                    score -= 0.35

            if chunk_idx < min_window_idx:
                score -= min(0.35, (min_window_idx - chunk_idx) * 0.03)
            elif chunk_idx > max_window_idx:
                score -= min(0.35, (chunk_idx - max_window_idx) * 0.03)

            candidate_scores.append((chunk_idx, score))
            if score > best_score:
                best_score = score
                best_idx = chunk_idx

        if candidate_scores:
            candidate_scores.sort(key=lambda item: item[1], reverse=True)
            score_threshold = best_score - 0.08
            top_candidates = [
                idx for idx, score in candidate_scores[:top_k_candidates]
                if score >= score_threshold
            ]
            if top_candidates:
                best_idx = min(top_candidates, key=lambda idx: int(timeline[idx].get("time_sec", 0) or 0))

        if exact_anchor_idx is not None:
            best_idx = exact_anchor_idx
            best_score = max(best_score, 1.5)

        resolved_start = original_start_sec if original_start_sec is not None else 0
        if best_idx is not None:
            matched_start = int(timeline[best_idx].get("time_sec", resolved_start) or resolved_start)
            if best_score >= 0.16 or original_start_sec is None:
                resolved_start = matched_start

        if aligned_sections:
            previous_start = int(aligned_sections[-1].get("start_sec", 0) or 0)
            if resolved_start <= previous_start:
                resolved_start = previous_start + min_gap_seconds
            elif resolved_start - previous_start < min_gap_seconds:
                resolved_start = previous_start + min_gap_seconds

        aligned_sections.append({
            **section,
            "title": title,
            "summary": section_summary,
            "start_sec": resolved_start,
        })

        if best_idx is not None:
            cursor = min(best_idx + 1, len(timeline) - 1)

    return {
        "title": str(summary.get("title", "")).strip() or "전체 내용 요약",
        "overview": str(summary.get("overview", "")).strip(),
        "sections": aligned_sections,
    }


def _build_balanced_media_transcript_excerpt(
    timeline: list[dict[str, Any]],
    *,
    max_chars: int = 12000,
) -> str:
    transcript_lines = [
        f"{_format_media_timestamp(item['time_sec'])} {item['text']}"
        for item in timeline
        if item.get("text")
    ]
    if not transcript_lines:
        return ""

    full_transcript = "\n".join(transcript_lines)
    if len(full_transcript) <= max_chars:
        return full_transcript

    avg_len = max(1, sum(len(line) for line in transcript_lines) // len(transcript_lines))
    target_count = max(8, min(len(transcript_lines), max_chars // max(avg_len, 40)))
    if target_count >= len(transcript_lines):
        return full_transcript[:max_chars]

    selected_indices = {0, len(transcript_lines) - 1}
    for idx in range(target_count):
        sampled_idx = round(idx * (len(transcript_lines) - 1) / max(target_count - 1, 1))
        selected_indices.add(sampled_idx)

    ordered_lines = [transcript_lines[idx] for idx in sorted(selected_indices)]
    excerpt = "\n".join(ordered_lines)
    if len(excerpt) <= max_chars:
        return excerpt

    trimmed_lines: list[str] = []
    current_len = 0
    for line in ordered_lines:
        next_len = current_len + (1 if trimmed_lines else 0) + len(line)
        if next_len > max_chars:
            break
        trimmed_lines.append(line)
        current_len = next_len
    return "\n".join(trimmed_lines)


def _normalize_media_fallback_text(text: str) -> str:
    normalized = re.sub(r"\s+", " ", (text or "")).strip()
    normalized = normalized.strip(" .,-")
    return normalized


def _truncate_media_fallback_text(text: str, max_len: int) -> str:
    normalized = _normalize_media_fallback_text(text)
    if len(normalized) <= max_len:
        return normalized

    cut = normalized[:max_len].rstrip()
    for delimiter in (". ", "! ", "? ", ", ", " "):
        idx = cut.rfind(delimiter)
        if idx >= max_len // 2:
            cut = cut[:idx].rstrip(" .,!?")
            break
    return cut.rstrip(" .,!?") + "..."


def _build_media_fallback_section_title(text: str, index: int) -> str:
    normalized = _normalize_media_fallback_text(text)
    if not normalized:
        return f"주요 내용 {index}"

    title = _truncate_media_fallback_text(normalized, 28)
    if title.endswith("..."):
        title = title[:-3].rstrip()
    return title or f"주요 내용 {index}"


def _build_media_fallback_section_summary(text: str) -> str:
    normalized = _normalize_media_fallback_text(text)
    if not normalized:
        return "이 구간의 주요 내용을 요약했습니다."

    summary = _truncate_media_fallback_text(normalized, 180)
    if not re.search(r"[.!?]$", summary):
        summary += "."
    return summary


def _extract_media_transcript(file_bytes: bytes, filename: str) -> tuple[list[dict[str, Any]], int]:
    import os
    import tempfile

    ext = filename.lower().rsplit(".", 1)[-1]
    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    prepared_path = tmp_path
    prepared_is_temp = False
    try:
        prepared_path, prepared_is_temp = _prepare_media_file_for_transcription(tmp_path, ext)
        with open(prepared_path, "rb") as f:
            response = client.audio.transcriptions.create(
                model="whisper-1",
                file=f,
                prompt=(
                    "Transcribe the audio exactly as spoken. Do not translate. "
                    "Keep Korean in Korean and English in English. "
                    "If multiple languages are spoken, preserve each part in its original language."
                ),
                response_format="verbose_json",
                timestamp_granularities=["segment", "word"],
            )

        segments = getattr(response, "segments", None)
        if not segments:
            text = getattr(response, "text", "") or ""
            if text.strip():
                return [{"start_sec": 0, "end_sec": 0, "text": text.strip()}], 0
            return [], 0

        normalized: list[dict[str, Any]] = []
        for seg in segments:
            raw_text = getattr(seg, "text", "") or ""
            text = _hard_sanitize(_sanitize(raw_text)).strip()
            if not text:
                continue
            start_sec = getattr(seg, "start", 0) or 0
            end_sec = getattr(seg, "end", start_sec) or start_sec
            normalized.append({
                "start_sec": float(start_sec),
                "end_sec": float(max(start_sec, end_sec)),
                "text": text,
                "words": _extract_words_from_segment(seg),
            })
        return _refine_media_segments(normalized, is_music=_looks_like_music_filename(filename)), 0
    finally:
        if prepared_is_temp and os.path.exists(prepared_path):
            os.unlink(prepared_path)
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)


def _build_media_chunks(segments: list[dict[str, Any]], filename: str) -> list[str]:
    chunks: list[str] = []
    for segment in segments:
        text = _hard_sanitize(_sanitize(segment.get("text", ""))).strip()
        if not text:
            continue
        meta = {
            "filename": filename,
            "start_sec": float(segment.get("start_sec", 0) or 0),
            "end_sec": float(segment.get("end_sec", 0) or 0),
        }
        chunks.append(
            f"[MEDIA_META]{json.dumps(meta, ensure_ascii=False)}[/MEDIA_META]\n[MEDIA_TRANSCRIPT]\n{text}"
        )
    return chunks


def build_media_summary_from_chunks(rows: list[dict[str, Any]], filename: str) -> dict[str, Any]:
    timeline = _build_media_timeline_from_chunks_v2(rows)
    if not timeline:
        return {"title": "전체 내용 요약", "overview": "", "sections": []}

    transcript = _build_balanced_media_transcript_excerpt(timeline, max_chars=12000)

    prompt = f"""다음은 업로드된 비디오/오디오 '{filename}'의 시간별 전사입니다.
전체 시간 흐름을 읽고, 짧은 메모가 아니라 README나 요약 보고서처럼 읽히는 문서형 요약을 만들어주세요.

중요 규칙:
- 전사 원문을 그대로 복사하지 말고, 내용을 이해한 뒤 자연스럽고 정리된 설명문으로 다시 작성하세요.
- 전체 요약은 "제목 + 간단한 도입 설명 + 주제별 본문" 구조의 문서를 만든다는 느낌으로 작성하세요.
- 맨 위의 title은 문서 제목처럼 구체적으로 쓰세요.
- overview는 제목 바로 아래에 들어가는 짧은 소개 문단처럼 2~4문장으로 작성하세요.
- sections는 시간 순서대로 정리하고, 각 section은 하나의 주제나 흐름 전환을 대표해야 합니다.
- 각 section의 title은 소제목처럼 자연스럽고 구체적으로 작성하세요. '주제 1', '핵심 내용 2' 같은 generic 이름은 쓰지 마세요.
- 각 section의 summary는 한 줄 메모가 아니라, 해당 구간에서 무엇을 설명하고 왜 중요한지 드러나는 2~5문장짜리 본문처럼 작성하세요.
- overview와 각 section의 summary는 기본적으로 자연스러운 본문 형태로 작성하되, 필요하면 Markdown 형식의 강조나 목록을 사용할 수 있습니다.
- 핵심 개념, 중요한 결론, 주의할 부분은 필요하면 **굵게** 표시해도 됩니다.
- 설명이 여러 포인트로 나뉘면 불릿 목록(-)이나 번호 목록을 사용해 읽기 쉽게 정리해도 됩니다.
- 화면에서는 소제목 오른쪽에 시간만 따로 붙여 보여줄 예정이므로, summary 안에 시간을 다시 반복해서 쓰지 마세요.
- 각 section에는 그 주제가 본격적으로 시작되는 대표 시작 시간만 start_sec에 넣으세요.
- Markdown 형식은 허용하지만, 불필요한 구분선은 넣지 말고 문서 본문처럼 매끄럽게 읽히는 내용을 작성하세요.
- 반드시 JSON 객체만 출력하세요.

출력 형식:
{{
  "title": "전체 내용 요약 제목",
  "overview": "영상/오디오 전체 내용을 소개하는 2~4문장 분량의 도입 설명",
  "sections": [
    {{
      "title": "소제목 1",
      "start_sec": 35,
      "summary": "이 구간에서 다루는 핵심 내용, 맥락, 중요한 포인트를 문서 본문처럼 2~5문장으로 자세히 설명"
    }}
  ]
}}

시간별 전사:
{transcript}
"""

    try:
        raw = _call_llm(
            messages=[{"role": "user", "content": prompt}],
            model=_DEFAULT_MODEL,
            temperature=0.2,
            max_tokens=900,
            json_mode=True,
        )
        parsed = json.loads(raw)
        payload = _normalize_media_summary_payload(parsed)
        sections = payload.get("sections", [])
        if isinstance(sections, list):
            normalized_sections: list[dict[str, Any]] = []
            for item in sections:
                if not isinstance(item, dict):
                    continue
                title = str(item.get("title", "")).strip()
                summary = str(item.get("summary", "")).strip()
                start_sec = item.get("start_sec", 0)
                if not title or not summary:
                    continue
                try:
                    start_value = int(float(start_sec))
                except Exception:
                    start_value = 0
                normalized_sections.append({
                    "title": title,
                    "start_sec": max(0, start_value),
                    "summary": summary,
                })
            if normalized_sections:
                return {
                    "title": str(payload.get("title", "")).strip() or "전체 내용 요약",
                    "overview": str(payload.get("overview", "")).strip(),
                    "sections": normalized_sections,
                }
    except Exception:
        pass

    fallback: list[dict[str, Any]] = []
    sample_points = timeline[: min(4, len(timeline))]
    for idx, item in enumerate(sample_points, start=1):
        fallback.append({
            "title": _build_media_fallback_section_title(str(item.get("text", "")), idx),
            "start_sec": item["time_sec"],
            "summary": _build_media_fallback_section_summary(str(item.get("text", ""))),
        })
    return {
        "title": "전체 내용 요약",
        "overview": f"{filename}의 핵심 내용을 시간 흐름에 따라 자연스럽게 정리한 요약입니다.",
        "sections": fallback,
    }


# ─────────────────────────────────────────────
# DB 안전 삽입
# ─────────────────────────────────────────────

def _db_safe_insert(table: str, records: list[dict]) -> None:
    """null byte가 절대 PostgREST에 전달되지 않도록 직접 HTTP로 삽입."""
    import httpx
    serialized = json.dumps(records, ensure_ascii=True)
    cleaned = serialized.replace('\\u0000', '')
    safe_records = json.loads(cleaned)

    url = f"{settings.SUPABASE_URL}/rest/v1/{table}"
    headers = {
        "apikey": settings.SUPABASE_SERVICE_ROLE_KEY,
        "Authorization": f"Bearer {settings.SUPABASE_SERVICE_ROLE_KEY}",
        "Content-Type": "application/json",
        "Prefer": "return=minimal",
    }
    batch_size = 50
    for i in range(0, len(safe_records), batch_size):
        batch_json = json.dumps(safe_records[i:i + batch_size], ensure_ascii=True).replace('\\u0000', '')
        print(f"[_db_safe_insert] batch {i//batch_size+1}, rows={min(batch_size, len(safe_records)-i)}")
        response = httpx.post(url, content=batch_json.encode('utf-8'), headers=headers, timeout=30)
        if response.status_code not in (200, 201):
            raise Exception(f"document_chunks insert 실패: {response.status_code} {response.text[:200]}")


# ─────────────────────────────────────────────
# 임베딩
# ─────────────────────────────────────────────

def _embed_batch(texts: list[str], batch_size: int = 100) -> list[list[float]]:
    """sentence-transformers 로컬 모델로 임베딩 생성 (API 키 불필요)."""
    model = _get_embedding_model()
    all_embeddings: list[list[float]] = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        vecs = model.encode(batch, normalize_embeddings=True)
        all_embeddings.extend(vecs.tolist())
    return all_embeddings


def _cosine_similarity(query_vec: list[float], embeddings: list[list[float]]) -> np.ndarray:
    """query_vec과 embeddings 행렬의 코사인 유사도 벡터 반환.
    차원 불일치 시(기존 DB 임베딩과 새 모델 차원이 다를 때) 균등 점수로 폴백."""
    q = np.array(query_vec, dtype=np.float32)
    m = np.array(embeddings, dtype=np.float32)
    if q.shape[0] != m.shape[1]:
        # 차원이 다르면 모든 청크를 동등하게 취급 (순서대로 선택)
        return np.ones(m.shape[0], dtype=np.float32)
    q_norm = np.linalg.norm(q)
    m_norms = np.linalg.norm(m, axis=1)
    m_norms = np.where(m_norms == 0, 1e-10, m_norms)
    return (m @ q) / (m_norms * (q_norm if q_norm != 0 else 1e-10))


# ─────────────────────────────────────────────
# 파일 형식별 텍스트 추출
# ─────────────────────────────────────────────

def _extract_text_from_pdf(file_bytes: bytes, vision_descriptions: dict | None = None) -> tuple[str, int]:
    """PDF에서 텍스트 추출. 페이지마다 [페이지 N] 마커를 붙여 반환.

    Args:
        vision_descriptions: {page_num(int): description} — _analyze_pdf_images에서 전달
    """
    parts: list[str] = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        page_count = len(pdf.pages)
        for page_idx, page in enumerate(pdf.pages, start=1):
            lines: list[str] = [f"[페이지 {page_idx}]"]

            # 텍스트 추출
            raw = page.extract_text() or ""
            text = _hard_sanitize(_sanitize(raw)).strip()
            if text:
                lines.append(text)

            # 표 추출 (마크다운 형식으로 변환)
            try:
                tables = page.extract_tables()
                for table in (tables or []):
                    if not table:
                        continue
                    rows_md: list[str] = []
                    for row_idx, row in enumerate(table):
                        cells = [str(cell or "").strip().replace("\n", " ") for cell in row]
                        rows_md.append("| " + " | ".join(cells) + " |")
                        if row_idx == 0:
                            rows_md.append("|" + "|".join(["---"] * len(cells)) + "|")
                    if rows_md:
                        lines.append("\n".join(rows_md))
            except Exception:
                pass

            # Vision AI 분석 결과 삽입
            if vision_descriptions and page_idx in vision_descriptions:
                lines.append(f"[시각 자료] {vision_descriptions[page_idx]}")

            if len(lines) > 1:  # [페이지 N] 마커 외에 내용이 있을 때만 포함
                parts.append("\n".join(lines))

    return "\n\n".join(parts), page_count


def _extract_text_from_docx(file_bytes: bytes) -> tuple[str, int]:
    """DOCX에서 텍스트 추출. [페이지 N] 마커 포함. (text, page_count) 반환."""
    from docx import Document
    from docx.oxml.ns import qn
    doc = Document(io.BytesIO(file_bytes))

    W_BR    = qn("w:br")
    W_TYPE  = qn("w:type")
    W_LRPB  = qn("w:lastRenderedPageBreak")

    lines: list[str] = []
    page_num = 1
    lines.append(f"[페이지 {page_num}]")

    def _has_page_break(para) -> bool:
        """단락 XML 안에 페이지 나누기(명시적/렌더링) 요소가 있으면 True."""
        for elem in para._p.iter():
            tag = elem.tag
            if tag == W_BR and elem.get(W_TYPE) == "page":
                return True
            if tag == W_LRPB:
                return True
        return False

    for para in doc.paragraphs:
        if _has_page_break(para):
            page_num += 1
            lines.append(f"[페이지 {page_num}]")
        text = para.text.strip()
        if text:
            lines.append(text)

    # 표 내용 — 표는 페이지 순서 유지를 위해 body element 순회로 재처리
    # (doc.paragraphs 는 표 안 단락을 포함하지 않으므로 별도 추가)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                lines.append(row_text)

    return "\n".join(lines), page_num


def _extract_text_from_pptx(file_bytes: bytes, vision_descriptions: dict | None = None) -> tuple[str, int]:
    """PPTX에서 텍스트 추출. vision_descriptions가 있으면 이미지/영상 설명도 포함.

    Args:
        vision_descriptions: {slide_num(int): description} — _generate_and_upload_slides에서 전달
    Returns:
        (text, slide_count)
    """
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    prs = Presentation(io.BytesIO(file_bytes))
    slides_text: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        lines: list[str] = [f"[슬라이드 {slide_idx}]"]

        # 텍스트 추출
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)

        # 이미지 포함 여부 감지 (PICTURE = 13)
        has_picture = any(shape.shape_type == 13 for shape in slide.shapes)

        # 영상 포함 여부 감지
        has_video = any(
            "video" in rel.reltype.lower()
            for rel in slide.part.rels.values()
        )

        # Vision AI 설명 추가
        if vision_descriptions and slide_idx in vision_descriptions:
            lines.append(f"[시각 자료 설명] {vision_descriptions[slide_idx]}")
        elif has_picture:
            lines.append("[이 슬라이드에는 이미지가 포함되어 있습니다]")

        if has_video:
            lines.append("[이 슬라이드에는 동영상이 포함되어 있습니다]")

        slides_text.append("\n".join(lines))

    return "\n\n".join(slides_text), len(prs.slides)


def _extract_text_from_hwpx(file_bytes: bytes) -> tuple[str, int]:
    """HWPX(zip 기반 XML)에서 텍스트 추출. (text, 0) 반환.

    HWPX는 ZIP 컨테이너 안에 Contents/section*.xml 파일들이 있고,
    실제 글자는 <hp:t> (로컬명 't') 요소에만 저장된다.
    모든 요소를 순회하면 메타데이터·속성값이 섞여 깨진 텍스트가 나오므로
    반드시 로컬명이 't'인 요소만 수집한다.
    """
    import zipfile
    from xml.etree import ElementTree as ET

    paragraphs: list[str] = []
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        # section*.xml 만 처리 (Contents/ 폴더)
        section_files = sorted(
            name for name in z.namelist()
            if name.startswith("Contents/") and name.endswith(".xml")
            and "section" in name.lower()
        )
        if not section_files:
            # fallback: Contents 폴더의 모든 xml
            section_files = sorted(
                name for name in z.namelist()
                if "Contents" in name and name.endswith(".xml")
            )

        for name in section_files:
            with z.open(name) as f:
                try:
                    root = ET.fromstring(f.read())
                except Exception:
                    continue

                # 단락(p) 단위로 묶어서 줄바꿈 보존
                # 로컬명 'p' 요소를 단락으로 간주
                p_tag_found = any(
                    (elem.tag.split("}")[-1] == "p")
                    for elem in root.iter()
                )

                if p_tag_found:
                    for elem in root.iter():
                        if elem.tag.split("}")[-1] != "p":
                            continue
                        # 이 단락 안의 모든 <t> 텍스트를 합침
                        run_texts = [
                            child.text
                            for child in elem.iter()
                            if child.tag.split("}")[-1] == "t" and child.text
                        ]
                        line = "".join(run_texts)
                        paragraphs.append(line)
                else:
                    # 단락 구조가 없으면 <t> 만 수집
                    for elem in root.iter():
                        if elem.tag.split("}")[-1] == "t" and elem.text:
                            paragraphs.append(elem.text)

    return "\n".join(paragraphs), 0


def _extract_text_from_hwp(file_bytes: bytes) -> tuple[str, int]:
    """HWP 5.x (OLE 바이너리 포맷)에서 텍스트 추출. (text, 0) 반환.

    HWP5 OLE 스트림 구조:
      BodyText/Section0000 (or Section0) → zlib 압융 해제 후 레코드 스트림
      레코드 포맷: 4바이트 헤더 (rec_type 10bit + 두패 + size 12bit)
                          size == 0xFFF 이면 다음 4바이트가 확장 크기

    HWPTAG_PARA_TEXT (rec_type 67) 파싱 규칙:
      - 0x0020 이상: 2바이트 UTF-16-LE 코드포인트 → 실제 문자
      - 0x0000~0x001F: 인라인 콘트롤 → 코드(2) + 헤더(8) = 10바이트 스킵
        깴떴지 않으면 헤더 바이트가 한자 등으로 잘못 디코딩됨
    """
    import zlib
    import struct
    import olefile

    texts: list[str] = []
    try:
        ole = olefile.OleFileIO(io.BytesIO(file_bytes))
        section_idx = 0
        while True:
            # HWP5 파일에 따라 Section0000 / Section0 두 가지 이름 형식이 존재
            section_name = f"BodyText/Section{section_idx:04d}"
            if not ole.exists(section_name):
                section_name = f"BodyText/Section{section_idx}"
                if not ole.exists(section_name):
                    break
            stream_data = ole.openstream(section_name).read()
            try:
                decompressed = zlib.decompress(stream_data, -15)
            except Exception:
                decompressed = stream_data

            pos = 0
            while pos < len(decompressed):
                if pos + 4 > len(decompressed):
                    break
                header = struct.unpack_from("<I", decompressed, pos)[0]
                rec_type = header & 0x3FF
                size = (header >> 20) & 0xFFF
                if size == 0xFFF:
                    if pos + 8 > len(decompressed):
                        break
                    size = struct.unpack_from("<I", decompressed, pos + 4)[0]
                    pos += 8
                else:
                    pos += 4
                data = decompressed[pos: pos + size]
                pos += size
                # Record type 67 = 단락 텍스트 (HWPTAG_PARA_TEXT)
                if rec_type == 67 and data:
                    try:
                        # HWP PARA_TEXT 포맷:
                        # - 일반 문자: 2바이트 UTF-16-LE 코드포인트 (0x0020 이상)
                        # - 인라인 컨트롤(0x0000~0x001F): 코드 2바이트 + 컨트롤 헤더 8바이트
                        #   → 건너뛰지 않으면 헤더 바이트가 한자 등으로 잘못 디코딩됨
                        text_chars: list[str] = []
                        i = 0
                        while i + 1 < len(data):
                            code = struct.unpack_from("<H", data, i)[0]
                            if code < 0x0020:
                                # 인라인 컨트롤: 코드(2) + 헤더(8) = 10바이트 스킵
                                i += 10
                            else:
                                ch = chr(code)
                                if ch.isprintable() or ch in "\n\t ":
                                    text_chars.append(ch)
                                i += 2
                        text = "".join(text_chars)
                        if text.strip():
                            texts.append(text.strip())
                    except Exception:
                        pass
            section_idx += 1
    except Exception as e:
        raise ValueError(f"HWP 파일을 읽을 수 없습니다: {e}")
    return "\n".join(texts), 0


def _compress_image_to_base64(img_bytes: bytes, max_width: int = 800, quality: int = 75) -> tuple[str, str]:
    """이미지 bytes를 리사이즈·압축 후 (base64_str, mime_type) 반환.

    브라우저에서 원본 크기 이미지(수 MB)를 Base64로 한꺼번에 로드하면
    파싱 부하로 페이지가 멈추는 문제를 방지하기 위해 축소·압축한다.

    처리 흐름:
      1. Pillow로 이미지 열기
      2. 투명 채널(RGBA/PA/LA) → 흰 배경 RGB 합성 (JPEG는 알파 미지원)
      3. 가로가 max_width(기본 800px) 초과 시 비율 유지 축소
      4. JPEG 품질 quality(기본 75)로 저장해 파일 크기 대폭 감소
      5. Base64 인코딩 후 (base64_str, "image/jpeg") 반환

    실패 시 Pillow 없이 원본을 그대로 Base64로 반환(PNG로 가정).
    """
    try:
        from PIL import Image as _PILImage
        import io as _io
        import base64 as _b64
        img = _PILImage.open(_io.BytesIO(img_bytes))
        # RGBA / 팔레트(P) / LA 모드는 흰 배경 RGB로 합성
        # (JPEG 포맷이 알파 채널을 지원하지 않으므로 변환 필요)
        if img.mode in ("RGBA", "P", "LA"):
            bg = _PILImage.new("RGB", img.size, (255, 255, 255))
            if img.mode == "P":
                img = img.convert("RGBA")
            bg.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
            img = bg
        elif img.mode != "RGB":
            img = img.convert("RGB")
        # 가로 max_width 초과 시 비율 유지 축소 (세로는 자동 계산)
        if img.width > max_width:
            ratio = max_width / img.width
            img = img.resize((max_width, int(img.height * ratio)), _PILImage.LANCZOS)
        buf = _io.BytesIO()
        img.save(buf, format="JPEG", quality=quality, optimize=True)
        b64 = _b64.b64encode(buf.getvalue()).decode("utf-8")
        return b64, "image/jpeg"
    except Exception:
        # Pillow 미설치 또는 지원하지 않는 포맷 → 원본 bytes를 그대로 Base64
        import base64 as _b64
        return _b64.b64encode(img_bytes).decode("utf-8"), "image/png"


def _extract_markdown_from_docx(file_bytes: bytes) -> str:
    """DOCX(Office Open XML) 파일을 마크다운으로 변환.

    DOCX 구조에서 다음을 추출한다:
      - 단락(w:p): 스타일 기반 헤딩(#~#####), 목록(- ), 인라인 볼드/이탤릭
      - 표(w:tbl): GFM 파이프 테이블 형식으로 변환, 첫 행에 헤더 구분선 삽입
      - 이미지(w:drawing → a:blip r:embed): 관계 ID로 blob 조회 후
        _compress_image_to_base64()로 압축해 data: URL Base64 마크다운 삽입

    주요 내부 함수:
      _get_image_md(rId)  : 관계 ID → ![이미지](data:...) 마크다운
      _is_on(rpr, tag)    : w:rPr 속성(w:val)으로 서식 on/off 판단
      _para_to_md(p_elem) : 단락 하나를 마크다운 라인으로 변환
      _tc_text(tc_elem)   : 셀(w:tc) 텍스트 추출 (중첩 표 포함)
      _table_to_md(tbl)   : 표 전체를 GFM 테이블 행 목록으로 변환

    반환: 마크다운 문자열 (연속 빈 줄 정리 완료)
    """
    import base64 as _b64
    import zipfile as _zf
    from docx import Document
    from docx.oxml.ns import qn

    # ── OOXML 태그 QName 상수 ─────────────────────────────────────────
    # python-docx의 qn()은 "prefix:local" → 완전한 Clark 표기({ns}local)로 변환
    W_P      = qn("w:p")       # 단락
    W_TBL    = qn("w:tbl")     # 표
    W_TR     = qn("w:tr")      # 표 행
    W_TC     = qn("w:tc")      # 표 셀
    W_R      = qn("w:r")       # 텍스트 런(run)
    W_T      = qn("w:t")       # 실제 텍스트 내용
    W_B      = qn("w:b")       # 굵게(bold)
    W_I      = qn("w:i")       # 기울임(italic)
    W_RPR    = qn("w:rPr")     # 런 속성(run properties)
    W_PPR    = qn("w:pPr")     # 단락 속성(paragraph properties)
    W_PSTYLE = qn("w:pStyle")  # 단락 스타일 참조
    W_VAL    = qn("w:val")     # 속성값 (w:val="0" 이면 서식 off 의미)
    W_DRAWING = qn("w:drawing") # 이미지/도형 컨테이너
    W_BLIP   = qn("a:blip")    # DrawingML 이미지 참조 (r:embed rId)
    R_EMBED  = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"  # 관계 ID 속성

    doc = Document(io.BytesIO(file_bytes))

    # 스타일 ID → 이름 맵 (예: "Heading1" → "Heading 1")
    # 헤딩 판별에 사용 — 한글 스타일명("제목 1" 등)도 포함
    style_map: dict[str, str] = {s.style_id: s.name for s in doc.styles if s.style_id}

    HEADING_MAP: dict[str, str] = {
        "Heading 1": "#", "Heading 2": "##", "Heading 3": "###",
        "Heading 4": "####", "Heading 5": "#####",
        "Title": "#", "Subtitle": "##",
        "제목 1": "#", "제목 2": "##", "제목 3": "###", "제목 4": "####",
        "제목": "#",
    }

    # 이미지 관계 ID → Base64 마크다운 캐시
    # 같은 이미지가 여러 단락에 반복 삽입된 경우 재압축 방지
    _img_cache: dict[str, str] = {}

    def _get_image_md(rId: str) -> str:
        """관계 ID(rId)로 DOCX 내부 이미지를 조회해 마크다운 Base64 이미지로 반환.

        DOCX에서 이미지는 w:drawing → a:blip r:embed=rId 경로로 참조된다.
        doc.part.related_parts[rId].blob 으로 raw bytes를 가져와
        _compress_image_to_base64()로 축소·압축 후 data: URL로 삽입한다.
        """
        if rId in _img_cache:
            return _img_cache[rId]
        try:
            part = doc.part.related_parts.get(rId)
            if part is None:
                return ""
            img_bytes = part.blob
            b64, mime = _compress_image_to_base64(img_bytes)
            result = f"![이미지](data:{mime};base64,{b64})"
            _img_cache[rId] = result
            return result
        except Exception:
            return ""

    def _is_on(rpr_elem, tag):
        """w:rPr 자식 요소의 서식이 활성화돼 있는지 확인.

        OOXML에서 <w:b/> 처럼 단독으로 존재하면 on,
        <w:b w:val="0"/> 이면 명시적 off다.
        """
        el = rpr_elem.find(tag)
        return el is not None and el.get(W_VAL, "1") not in ("0", "false", "off")

    def _para_to_md(p_elem) -> str:
        ppr = p_elem.find(W_PPR)
        style_id = ""
        if ppr is not None:
            ps = ppr.find(W_PSTYLE)
            if ps is not None:
                style_id = ps.get(W_VAL, "")
        style_name = style_map.get(style_id, "")

        # 이미지 포함 여부 확인
        img_parts: list[str] = []
        for drawing in p_elem.iter(W_DRAWING):
            for blip in drawing.iter(W_BLIP):
                rId = blip.get(R_EMBED, "")
                if rId:
                    md = _get_image_md(rId)
                    if md:
                        img_parts.append(md)

        raw = "".join(
            t.text or "" for r in p_elem.findall(W_R) for t in r.findall(W_T)
        ).strip()

        # 이미지만 있는 단락
        if img_parts and not raw:
            return "\n".join(img_parts)

        if not raw and not img_parts:
            return ""

        # 헤딩
        for key, prefix in HEADING_MAP.items():
            if style_name == key or style_name.startswith(key + " "):
                text_part = f"{prefix} {raw}" if raw else ""
                return "\n".join(filter(None, img_parts + [text_part]))

        # 목록
        if "List" in style_name or "목록" in style_name:
            text_part = f"- {raw}" if raw else ""
            return "\n".join(filter(None, img_parts + [text_part]))

        # 인라인 서식
        parts: list[str] = []
        for r in p_elem.findall(W_R):
            text = "".join(t.text or "" for t in r.findall(W_T))
            if not text:
                continue
            rpr = r.find(W_RPR)
            bold   = rpr is not None and _is_on(rpr, W_B)
            italic = rpr is not None and _is_on(rpr, W_I)
            if bold and italic:
                parts.append(f"***{text}***")
            elif bold:
                parts.append(f"**{text}**")
            elif italic:
                parts.append(f"*{text}*")
            else:
                parts.append(text)
        text_part = "".join(parts).strip()
        return "\n".join(filter(None, img_parts + [text_part]))

    def _tc_text(tc_elem) -> str:
        """표 셀(w:tc) 내 모든 단락 텍스트를 공백으로 이어붙여 반환.

        파이프 테이블 열 구분자 '|'가 셀 내용에 포함되면 GFM 테이블이
        깨지므로 전각 문자(｜)로 치환한다.
        """
        parts: list[str] = []
        for p in tc_elem.findall(f".//{W_P}"):
            t = "".join(
                txt.text or "" for r in p.findall(W_R) for txt in r.findall(W_T)
            ).strip()
            if t:
                parts.append(t)
        return " ".join(parts).replace("|", "｜")

    def _table_to_md(tbl_elem) -> list[str]:
        """w:tbl 요소를 GFM 파이프 테이블 행 목록으로 변환.

        첫 행 다음에 헤더/본문 구분선(| --- | --- |)을 자동 삽입한다.
        직계 w:tr 만 처리해 중첩 표(표 안의 표)의 행이 섞이는 것을 방지한다.
        """
        rows_md: list[str] = []
        # 직계 TR 자식만 처리 (중첩 표 혼입 방지)
        direct_trs = [c for c in tbl_elem if c.tag == W_TR]
        for row_idx, tr in enumerate(direct_trs):
            cells = tr.findall(W_TC)
            if not cells:
                continue
            cell_texts = [_tc_text(tc) for tc in cells]
            rows_md.append("| " + " | ".join(cell_texts) + " |")
            if row_idx == 0:
                rows_md.append("| " + " | ".join(["---"] * len(cell_texts)) + " |")
        return rows_md

    W_BR_MD   = qn("w:br")
    W_TYPE_MD = qn("w:type")
    W_LRPB_MD = qn("w:lastRenderedPageBreak")

    def _p_has_page_break(p_elem) -> bool:
        for elem in p_elem.iter():
            tag = elem.tag
            if tag == W_BR_MD and elem.get(W_TYPE_MD) == "page":
                return True
            if tag == W_LRPB_MD:
                return True
        return False

    lines: list[str] = []
    page_num = 1
    lines.append(f"**[페이지 {page_num}]**")

    for child in doc.element.body:
        tag = child.tag
        if tag == W_P:
            if _p_has_page_break(child):
                page_num += 1
                lines.append(f"\n**[페이지 {page_num}]**")
            lines.append(_para_to_md(child))
        elif tag == W_TBL:
            lines.append("")
            lines.extend(_table_to_md(child))
            lines.append("")

    # 연속 빈 줄 정리
    result: list[str] = []
    prev_blank = False
    for ln in lines:
        blank = not ln.strip()
        if blank and prev_blank:
            continue
        result.append(ln)
        prev_blank = blank
    return "\n".join(result).strip()


def _extract_markdown_from_hwpx(file_bytes: bytes) -> str:
    """HWPX(ZIP + OWPML XML) 파일을 마크다운으로 변환.

    HWPX 파일 구조:
      ├── Contents/section0.xml  ← 본문 섹션 (단락·표·이미지 포함)
      ├── BinData/BIN0001.jpg    ← 이미지 바이너리
      └── Header/header.xml      ← 스타일 정의 + binItem(이미지 ID↔파일 매핑)

    이미지 처리 흐름 (2단계):
      1) _raw_bin  : ZIP 내 모든 파일을 (경로/파일명/확장자 없는 이름) 3가지 키로 인덱싱
      2) _binitem_map : header.xml의 <hh:binItem id="BIN0001" href="BinData/...">
                        파싱으로 ID → (bytes, ext) 매핑 구성
      → _para_images() 에서 <hp:img binaryItemIDRef="BIN0001"> 속성값으로 조회

    * OWPML 공식 스펙에서 이미지 참조 속성명은 'binaryItemIDRef' 이므로
      'href', 'src' 가 아닌 이 이름으로만 탐색해야 한다.

    단락/표 처리:
      - 단락(hp:p): 스타일 ID → 헤딩 prefix 변환, 표 셀 내 단락 중복 방지
      - 표(hp:tbl): GFM 파이프 테이블로 변환, 중첩 표 스킵

    반환: 마크다운 문자열 (연속 빈 줄 정리 완료)
    """
    import zipfile
    import base64 as _b64
    from xml.etree import ElementTree as ET

    def local(tag: str) -> str:
        return tag.split("}")[-1] if "}" in tag else tag

    HEADING_PATTERNS: list[tuple[str, str]] = [
        ("제목 1", "#"), ("제목1", "#"), ("Heading 1", "#"), ("Heading1", "#"),
        ("제목 2", "##"), ("제목2", "##"), ("Heading 2", "##"), ("Heading2", "##"),
        ("제목 3", "###"), ("제목3", "###"), ("Heading 3", "###"), ("Heading3", "###"),
        ("제목 4", "####"), ("제목4", "####"),
        ("제목", "#"),
    ]

    def _style_id(p_elem) -> str:
        for child in p_elem:
            if local(child.tag) == "pPr":
                for sub in child:
                    if local(sub.tag) == "pStyle":
                        for attr, val in sub.attrib.items():
                            if local(attr) in ("val", "styleId", "id"):
                                return val
        return ""

    def _heading_prefix(sid: str, name_map: dict) -> str:
        name = name_map.get(sid, sid)
        for pattern, prefix in HEADING_PATTERNS:
            if name == pattern or name.startswith(pattern):
                return prefix
        return ""

    def _para_text(p_elem) -> str:
        """p 요소 텍스트 추출. tbl 내부(표 셀)는 건너뜀 — 표 안의 텍스트가 단락으로 중복 출력되는 것 방지."""
        texts: list[str] = []

        def _collect(elem) -> None:
            if local(elem.tag) == "t" and elem.text:
                texts.append(elem.text)
            for child in elem:
                if local(child.tag) != "tbl":  # 표 내부로는 내려가지 않음
                    _collect(child)

        for child in p_elem:
            if local(child.tag) != "tbl":
                _collect(child)
        return "".join(texts)

    def _cell_text(tc_elem) -> str:
        parts: list[str] = []
        for elem in tc_elem.iter():
            if local(elem.tag) == "p":
                t = _para_text(elem).strip()
                if t:
                    parts.append(t)
        return " ".join(parts).replace("|", "｜")

    def _tbl_to_md(tbl_elem) -> list[str]:
        rows_md: list[str] = []
        row_idx = 0
        # 직계 tr만 처리
        direct_trs = [c for c in tbl_elem if local(c.tag) == "tr"]
        for tr in direct_trs:
            cells = [c for c in tr if local(c.tag) == "tc"]
            if not cells:
                continue
            cell_texts = [_cell_text(c) for c in cells]
            rows_md.append("| " + " | ".join(cell_texts) + " |")
            if row_idx == 0:
                rows_md.append("| " + " | ".join(["---"] * len(cell_texts)) + " |")
            row_idx += 1
        return rows_md

    lines: list[str] = []

    with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
        # ── 레거시 BinData 인덱스 (하위 호환용, 현재는 _raw_bin으로 대체) ──
        # 파일명(확장자 제거) → (bytes, ext) 매핑
        # _raw_bin 이 더 넓은 경로/이름 변형을 커버하므로 실제 조회는 _raw_bin/_binitem_map 사용
        bin_data: dict[str, tuple] = {}
        for zname in z.namelist():
            if ("bindata" in zname.lower() or "BinData" in zname) and not zname.endswith("/"):
                bname = zname.split("/")[-1]
                if "." in bname:
                    bid, ext = bname.rsplit(".", 1)
                    ext = ext.lower()
                else:
                    bid, ext = bname, "png"
                try:
                    bin_data[bid] = (z.read(zname), ext)
                except Exception:
                    pass

        MIME_MAP = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                    "gif": "image/gif", "webp": "image/webp", "bmp": "image/bmp"}

        # ── 이미지 로드 단계 1: ZIP 내 모든 파일 3-key 인덱싱 ──────────────────
        # HWPX에서 이미지 파일 경로는 한컴 버전마다 다를 수 있으므로
        # 전체 경로 / 파일명 / 확장자 없는 이름 세 가지 키로 동시 등록해
        # 어떤 참조 방식이 와도 빠르게 조회할 수 있도록 한다.
        _raw_bin: dict[str, tuple] = {}  # {lowercase_key: (bytes, ext)}
        for zname in z.namelist():
            if zname.endswith("/"):
                continue
            bname = zname.split("/")[-1]
            if "." in bname:
                bid_raw, ext = bname.rsplit(".", 1)
            else:
                bid_raw, ext = bname, "png"
            try:
                data = z.read(zname)
                _raw_bin[zname.lower()] = (data, ext.lower())      # 전체 경로 키
                _raw_bin[bname.lower()] = (data, ext.lower())      # 파일명만 키
                _raw_bin[bid_raw.lower()] = (data, ext.lower())    # 확장자 없는 이름 키
            except Exception:
                pass

        # ── 이미지 로드 단계 2: header.xml binItem → ID 매핑 구성 ───────────────
        # OWPML 스펙: <hh:binItem id="BIN0001" href="BinData/BIN0001.jpg" type="Embed"/>
        # 본문 XML의 <hp:img binaryItemIDRef="BIN0001"> 가 이 ID를 참조한다.
        # ID(대소문자 양쪽)를 키로 등록해 대소문자 불일치를 방지한다.
        _binitem_map: dict[str, tuple] = {}  # {id_string: (bytes, ext)}
        for zname in z.namelist():
            if ("header" in zname.lower() or "head" in zname.lower()) and zname.endswith(".xml"):
                try:
                    with z.open(zname) as f:
                        ht = ET.fromstring(f.read())
                    for elem in ht.iter():
                        if local(elem.tag) == "binItem":
                            bitem_id, href = None, None
                            for attr, val in elem.attrib.items():
                                la = local(attr)
                                if la == "id":
                                    bitem_id = val
                                elif la == "href":
                                    href = val
                            if bitem_id and href:
                                # 전체 경로 조회 → 실패 시 파일명만으로 재시도
                                entry = _raw_bin.get(href.lower()) or _raw_bin.get(href.split("/")[-1].lower())
                                if entry:
                                    _binitem_map[bitem_id] = entry
                                    _binitem_map[bitem_id.lower()] = entry  # 소문자 키 중복 등록
                except Exception:
                    pass

        def _get_img_md_by_ref(ref: str) -> str:
            """binaryItemIDRef 값으로 Base64 마크다운 이미지 태그를 반환.

            조회 우선순위:
              1. _binitem_map[ref]        — header.xml binItem ID 정확 매칭
              2. _binitem_map[ref.lower()]— 대소문자 무시 매칭
              3. _raw_bin[파일명]         — ID가 파일명 형태일 때 직접 조회
              4. _raw_bin[확장자 제거명]  — 마지막 fallback

            실패 시 빈 문자열 반환 (이미지 누락 처리는 호출 측에서 담당).
            """
            # 1~2단계: binItem ID 기반 조회
            entry = _binitem_map.get(ref) or _binitem_map.get(ref.lower())
            if not entry:
                # 3~4단계: 파일명/확장자 제거 이름으로 직접 조회
                fname = ref.replace("\\", "/").split("/")[-1].lower()
                bid_only = fname.rsplit(".", 1)[0] if "." in fname else fname
                entry = _raw_bin.get(fname) or _raw_bin.get(bid_only)
            if not entry:
                return ""
            img_bytes, _ext = entry
            b64, mime = _compress_image_to_base64(img_bytes)
            return f"![이미지](data:{mime};base64,{b64})"

        def _para_images(p_elem) -> list[str]:
            """p 요소 내 hp:img 태그(binaryItemIDRef 속성)를 찾아 Base64 마크다운 반환.
            HWPX 구조: <hp:picture><hp:img binaryItemIDRef="BIN0001" .../></hp:picture>"""
            imgs: list[str] = []
            seen: set[str] = set()
            for child in p_elem.iter():
                ltag = local(child.tag)
                if ltag == "img":
                    ref = None
                    for attr, val in child.attrib.items():
                        if local(attr) == "binaryItemIDRef" and val:
                            ref = val
                            break
                    if ref and ref not in seen:
                        seen.add(ref)
                        md = _get_img_md_by_ref(ref)
                        if md:
                            imgs.append(md)
            return imgs

        # 스타일 ID → 이름 맵 로드 (header.xml의 <style> 요소에서 수집)
        # 헤딩 판별에 사용: style_name_map["s1"] == "제목 1" 이면 # 프리픽스 적용
        style_name_map: dict[str, str] = {}
        for zname in z.namelist():
            if "header" in zname.lower() and zname.endswith(".xml"):
                try:
                    with z.open(zname) as f:
                        sr = ET.fromstring(f.read())
                    for elem in sr.iter():
                        if local(elem.tag) == "style":
                            sid, sname = None, None
                            for attr, val in elem.attrib.items():
                                if local(attr) in ("id", "styleId"):
                                    sid = val
                                elif local(attr) == "name":
                                    sname = val
                            for child in elem:
                                ltag = local(child.tag)
                                if ltag in ("id", "styleId") and child.text:
                                    sid = child.text.strip()
                                elif ltag in ("name", "localName") and child.text and not sname:
                                    sname = child.text.strip()
                            if sid and sname:
                                style_name_map[sid] = sname
                except Exception:
                    pass

        # HWPX 본문 섹션 파일 목록 (Contents/section0.xml, section1.xml ...)
        # 섹션이 없으면 Contents 폴더의 모든 XML로 폴백
        section_files = sorted(
            n for n in z.namelist()
            if n.startswith("Contents/") and n.endswith(".xml") and "section" in n.lower()
        )
        if not section_files:
            section_files = sorted(
                n for n in z.namelist() if "Contents" in n and n.endswith(".xml")
            )

        for sec_file in section_files:
            with z.open(sec_file) as f:
                try:
                    root = ET.fromstring(f.read())
                except Exception:
                    continue

            # parent_map: 자식 → 부모 역방향 매핑
            # _is_inside_tc() 에서 특정 요소가 표 셀(tc) 안에 있는지 빠르게 판단하기 위해 사용
            parent_map: dict = {}
            for parent in root.iter():
                for child in parent:
                    parent_map[child] = parent

            def _is_inside_tc(elem) -> bool:
                """elem 이 hp:tc(표 셀) 의 자손이면 True 반환.

                표 셀 내부 단락은 _tbl_to_md() 에서 별도로 처리되므로
                최상위 순회 시 중복 출력을 막기 위해 이 함수로 필터링한다.
                """
                current = elem
                while current in parent_map:
                    current = parent_map[current]
                    if local(current.tag) == "tc":
                        return True
                return False

            # 섹션 XML을 문서 순서 그대로 순회
            # 표 셀 바깥의 단락(p)과 표(tbl)만 처리한다
            # (표 셀 안의 p는 _tbl_to_md → _cell_text 에서 처리)
            for elem in root.iter():
                etag = local(elem.tag)
                if etag == "p":
                    if _is_inside_tc(elem):
                        continue
                    img_mds = _para_images(elem)
                    text = _para_text(elem).strip()
                    if img_mds:
                        lines.extend(img_mds)
                    if not text:
                        if not img_mds:
                            lines.append("")
                        continue
                    sid = _style_id(elem)
                    prefix = _heading_prefix(sid, style_name_map)
                    lines.append(f"{prefix} {text}" if prefix else text)
                elif etag == "tbl":
                    if _is_inside_tc(elem):
                        continue  # 중첩 표 스킵
                    tbl_lines = _tbl_to_md(elem)
                    if tbl_lines:
                        lines.append("")
                        lines.extend(tbl_lines)
                        lines.append("")

    # 연속 빈 줄 정리
    result: list[str] = []
    prev_blank = False
    for ln in lines:
        blank = not ln.strip()
        if blank and prev_blank:
            continue
        result.append(ln)
        prev_blank = blank
    return "\n".join(result).strip()


def _extract_text_from_image(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """이미지에서 GPT-4o Vision으로 텍스트/내용 + 구조화 메타데이터 추출.
    메타데이터는 [IMAGE_META]...[/IMAGE_META] 태그로 텍스트 앞에 삽입된다.
    Returns: (text_with_meta, 0)
    """
    import base64 as _base64
    ext = filename.lower().rsplit(".", 1)[-1]
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                "gif": "image/gif", "webp": "image/webp"}
    mime_type = mime_map.get(ext, "image/jpeg")
    b64 = _base64.b64encode(file_bytes).decode("utf-8")

    prompt = """이 이미지를 분석하여 아래 JSON 형식으로 응답하세요.

{
  "subjects": ["이미지에 등장하는 주요 대상들 (사람, 동물, 사물, 장소, 브랜드, 개념 등)"],
  "image_type": "photo | chart | diagram | screenshot | document | illustration",
  "visual_context": "이미지 전체 상황을 한 문장으로 (예: 실내에서 후드를 입고 있는 소형 강아지)",
  "extracted_text": "이미지에 포함된 모든 텍스트를 그대로 추출 (없으면 빈 문자열)",
  "description": "이미지 내용을 시각적 요소 중심으로 상세히 설명 (수치, 색상, 구조, 레이아웃 포함)"
}

image_type 선택 기준:
- photo: 실제 사진 (인물, 동물, 풍경, 사물)
- chart: 데이터 차트, 그래프, 표
- diagram: 흐름도, 아키텍처, 마인드맵, 구조도
- screenshot: 컴퓨터 화면 캡처 (UI, 코드, 웹페이지)
- document: 스캔된 문서, 손글씨, 텍스트 위주 이미지
- illustration: 그림, 일러스트, 아이콘

반드시 JSON만 응답하세요."""

    raw = _call_llm_vision(
        image_b64=b64,
        mime_type=mime_type,
        prompt=prompt,
        model="claude-sonnet-4-6",
        json_mode=True,
        max_tokens=2000,
    )
    try:
        parsed = json.loads(raw)
    except Exception:
        parsed = {}

    subjects = parsed.get("subjects", [])
    image_type = parsed.get("image_type", "photo")
    visual_context = parsed.get("visual_context", "")
    extracted_text = parsed.get("extracted_text", "")
    description = parsed.get("description", "")

    # 메타데이터를 특수 태그로 인코딩하여 텍스트 앞에 삽입
    meta_json = json.dumps({
        "subjects": subjects,
        "image_type": image_type,
        "visual_context": visual_context,
    }, ensure_ascii=False)
    meta_block = f"[IMAGE_META]{meta_json}[/IMAGE_META]"

    body = "\n\n".join(filter(None, [extracted_text, description]))
    full_text = f"{meta_block}\n\n{body}" if body else meta_block

    return full_text, 0


def _extract_text_from_video(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """비디오/오디오에서 OpenAI Whisper로 STT. (text, 0) 반환."""
    segments, _ = _extract_media_transcript(file_bytes, filename)
    text = "\n".join(segment["text"] for segment in segments if segment.get("text"))
    return text, 0


# ─────────────────────────────────────────────
# 인덱싱 (공통)
# ─────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "pptx", "ppt", "hwp", "hwpx",
    "jpg", "jpeg", "png", "gif", "webp",
    "mp4", "mov", "avi", "mkv", "webm", "mp3", "m4a",
}

VIDEO_AUDIO_EXTENSIONS = {"mp4", "mov", "avi", "mkv", "webm", "mp3", "m4a"}
IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "gif", "webp"}


def ingest_document(file_bytes: bytes, doc_id: str, filename: str = "", pptx_vision: dict | None = None, pdf_vision: dict | None = None, hwpx_page_count: int = 0, hwpx_pdf_bytes: bytes | None = None) -> tuple[int, int]:
    """파일을 청크로 분할하고 임베딩을 포함하여 Supabase에 저장.

    Returns:
        (chunk_count, page_count)
    """
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else "pdf"

    media_chunks: list[str] | None = None

    if ext == "pdf":
        text, page_count = _extract_text_from_pdf(file_bytes, vision_descriptions=pdf_vision or {})
    elif ext == "docx":
        text, page_count = _extract_text_from_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        text, page_count = _extract_text_from_pptx(file_bytes, vision_descriptions=pptx_vision or {})
    elif ext == "hwpx":
        if hwpx_pdf_bytes:
            # pyhwpx로 생성된 PDF에서 정확한 페이지별 텍스트 추출
            text, page_count = _extract_text_from_pdf(hwpx_pdf_bytes)
        else:
            text, page_count = _extract_text_from_hwpx(file_bytes)
            # PDF 없을 때: 단락을 페이지 수로 균등 분배해 [페이지 N] 마커 삽입
            if hwpx_page_count > 0 and text:
                raw_lines = [l for l in text.split('\n') if l.strip()]
                total = len(raw_lines)
                paras_per_page = max(1, total // hwpx_page_count)
                paged: list[str] = []
                for p in range(1, hwpx_page_count + 1):
                    start = (p - 1) * paras_per_page
                    end = start + paras_per_page if p < hwpx_page_count else total
                    block = raw_lines[start:end]
                    if block:
                        paged.append(f"[페이지 {p}]\n" + "\n".join(block))
                text = "\n\n".join(paged)
                page_count = hwpx_page_count
    elif ext == "hwp":
        text, page_count = _extract_text_from_hwp(file_bytes)
    elif ext in IMAGE_EXTENSIONS:
        text, page_count = _extract_text_from_image(file_bytes, filename)
    elif ext in VIDEO_AUDIO_EXTENSIONS:
        segments, page_count = _extract_media_transcript(file_bytes, filename)
        media_chunks = _build_media_chunks(segments, filename)
        text = "\n".join(segment["text"] for segment in segments if segment.get("text"))
    else:
        raise ValueError(f"지원하지 않는 파일 형식: .{ext}")

    text = _hard_sanitize(_sanitize(text)).strip()
    if not text:
        return 0, page_count

    chunks: list[str] = []
    if ext in VIDEO_AUDIO_EXTENSIONS:
        chunks = media_chunks or []
    elif ext == "pdf":
        # PDF: [페이지 N] 마커 기준으로 페이지별 청킹
        page_parts = re.split(r'(?=\[페이지\s*\d+\])', text)
        for part in page_parts:
            part = _hard_sanitize(_sanitize(part)).strip()
            if not part:
                continue
            if len(part) <= CHUNK_SIZE * 2:
                chunks.append(part)
            else:
                header_match = re.match(r'(\[페이지\s*\d+\])', part)
                header = header_match.group(1) + "\n" if header_match else ""
                body = part[len(header):]
                step = CHUNK_SIZE - CHUNK_OVERLAP
                for i in range(0, len(body), step):
                    sub = _hard_sanitize(_sanitize(header + body[i: i + CHUNK_SIZE]))
                    if sub.strip():
                        chunks.append(sub)
    elif ext in ("pptx", "ppt"):
        # PPT: [슬라이드 N] 마커 기준으로 슬라이드별 청킹
        # → 슬라이드 번호가 청크마다 명확히 유지됨
        slide_parts = re.split(r'(?=\[슬라이드\s*\d+\])', text)
        for part in slide_parts:
            part = _hard_sanitize(_sanitize(part)).strip()
            if not part:
                continue
            # 슬라이드가 너무 길면 CHUNK_SIZE 단위로 추가 분할 (슬라이드 마커 보존)
            if len(part) <= CHUNK_SIZE * 2:
                chunks.append(part)
            else:
                header_match = re.match(r'(\[슬라이드\s*\d+\])', part)
                header = header_match.group(1) + "\n" if header_match else ""
                body = part[len(header):]
                step = CHUNK_SIZE - CHUNK_OVERLAP
                for i in range(0, len(body), step):
                    sub = _hard_sanitize(_sanitize(header + body[i: i + CHUNK_SIZE]))
                    if sub.strip():
                        chunks.append(sub)
    elif ext in ("docx", "hwpx"):
        # DOCX/HWPX: [페이지 N] 마커 기준으로 페이지별 청킹 (PDF와 동일 방식)
        page_parts = re.split(r'(?=\[페이지\s*\d+\])', text)
        for part in page_parts:
            part = _hard_sanitize(_sanitize(part)).strip()
            if not part:
                continue
            if len(part) <= CHUNK_SIZE * 2:
                chunks.append(part)
            else:
                header_match = re.match(r'(\[페이지\s*\d+\])', part)
                header = header_match.group(1) + "\n" if header_match else ""
                body = part[len(header):]
                step = CHUNK_SIZE - CHUNK_OVERLAP
                for i in range(0, len(body), step):
                    sub = _hard_sanitize(_sanitize(header + body[i: i + CHUNK_SIZE]))
                    if sub.strip():
                        chunks.append(sub)
    else:
        # DOCX / HWPX / HWP 등 — 문단(빈 줄) 경계 기준 청킹
        # 글자 수로 단순 자르면 문단 중간에서 끊겨 스크롤 위치 탐색 실패
        # → 빈 줄(\n\n) 또는 개행(\n) 기준으로 문단 목록 추출 후 합산
        paragraphs = [p.strip() for p in re.split(r'\n\n+', text) if p.strip()]
        if not paragraphs:
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]

        current_parts: list[str] = []
        current_len = 0
        overlap_parts: list[str] = []  # 다음 청크 시작에 붙일 오버랩 문단

        for para in paragraphs:
            para_len = len(para)

            # 단일 문단이 CHUNK_SIZE를 초과하면 그 자체로 하나의 청크
            if para_len > CHUNK_SIZE:
                # 지금까지 모은 내용을 먼저 청크로 확정
                if current_parts:
                    chunk = _hard_sanitize(_sanitize("\n\n".join(current_parts)))
                    if chunk.strip():
                        chunks.append(chunk)
                    overlap_parts = current_parts[-1:]  # 마지막 문단만 오버랩
                    current_parts = []
                    current_len = 0

                # 긴 문단 자체를 청크로 추가
                chunk = _hard_sanitize(_sanitize(para))
                if chunk.strip():
                    chunks.append(chunk)
                overlap_parts = [para]
                continue

            # 현재 문단을 추가했을 때 CHUNK_SIZE 초과 여부 확인
            sep_len = 2 if current_parts else 0  # "\n\n" 구분자
            if current_len + sep_len + para_len > CHUNK_SIZE and current_parts:
                # 현재까지 모은 내용으로 청크 확정
                chunk = _hard_sanitize(_sanitize("\n\n".join(current_parts)))
                if chunk.strip():
                    chunks.append(chunk)
                # 오버랩: 직전 청크의 마지막 문단(들)을 새 청크 시작에 포함
                current_parts = list(overlap_parts)
                current_len = sum(len(p) for p in current_parts) + 2 * max(0, len(current_parts) - 1)
                overlap_parts = []

            current_parts.append(para)
            current_len += sep_len + para_len
            # 다음 청크 오버랩을 위해 마지막 1개 문단 기억
            overlap_parts = current_parts[-1:]

        # 남은 문단 처리
        if current_parts:
            chunk = _hard_sanitize(_sanitize("\n\n".join(current_parts)))
            if chunk.strip():
                chunks.append(chunk)

    if not chunks:
        return 0, page_count

    # 임베딩 생성
    try:
        embeddings = _embed_batch(chunks)
    except Exception as e:
        print(f"[WARNING] 임베딩 생성 실패, 임베딩 없이 저장: {e}")
        embeddings = [None] * len(chunks)  # type: ignore

    records = [
        {
            "doc_id": doc_id,
            "chunk_index": idx,
            "content": _hard_sanitize(content),
            "embedding": emb,
        }
        for idx, (content, emb) in enumerate(zip(chunks, embeddings))
    ]

    _db_safe_insert("document_chunks", records)
    print(f"[INFO] ingest_document: {len(records)} chunks 저장 완료 ({ext})")
    return len(chunks), page_count


# ─────────────────────────────────────────────
# 인덱싱 (URL)
# ─────────────────────────────────────────────

def _extract_youtube_video_id(url: str) -> str | None:
    """YouTube URL에서 video ID 추출. 아니면 None."""
    import re
    patterns = [
        r"(?:youtube\.com/watch\?.*v=|youtu\.be/)([A-Za-z0-9_-]{11})",
        r"youtube\.com/embed/([A-Za-z0-9_-]{11})",
        r"youtube\.com/shorts/([A-Za-z0-9_-]{11})",
    ]
    for pattern in patterns:
        m = re.search(pattern, url)
        if m:
            return m.group(1)
    return None


def _is_youtube_url(url: str) -> bool:
    return _extract_youtube_video_id(url) is not None


def _fetch_youtube_transcript(video_id: str) -> str:
    """youtube_transcript_api로 자막 텍스트 추출 (한국어 → 영어 순)."""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled
    except ImportError as e:
        raise ValueError("유튜브 자막 추출 기능이 서버에 설치되어 있지 않습니다. youtube-transcript-api를 설치해주세요.") from e
    api = YouTubeTranscriptApi()
    try:
        transcript = api.fetch(video_id, languages=["ko", "ko-KR", "en", "en-US"])
    except NoTranscriptFound:
        try:
            transcript_list = api.list(video_id)
            # 자동 생성 포함 아무 언어나
            for t in transcript_list:
                transcript = api.fetch(video_id, languages=[t.language_code])
                break
            else:
                raise ValueError("유튜브 영상에 자막이 없습니다.")
        except TranscriptsDisabled:
            raise ValueError("이 유튜브 영상은 자막이 비활성화되어 있습니다.")
    except TranscriptsDisabled:
        raise ValueError("이 유튜브 영상은 자막이 비활성화되어 있습니다.")

    text = " ".join(s.text.replace("\n", " ") for s in transcript.snippets)
    return text.strip()


def _fetch_youtube_transcript_segments(video_id: str) -> list[dict[str, Any]]:
    """youtube_transcript_api로 자막 세그먼트 추출."""
    try:
        from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled
    except ImportError as e:
        raise ValueError("유튜브 자막 추출 기능이 서버에 설치되어 있지 않습니다. youtube-transcript-api를 설치해주세요.") from e

    api = YouTubeTranscriptApi()
    try:
        transcript = api.fetch(video_id, languages=["ko", "ko-KR", "en", "en-US"])
    except NoTranscriptFound:
        try:
            transcript_list = api.list(video_id)
            transcript = None
            for t in transcript_list:
                transcript = api.fetch(video_id, languages=[t.language_code])
                break
            if transcript is None:
                raise ValueError("유튜브 영상에 자막이 없습니다.")
        except TranscriptsDisabled:
            raise ValueError("이 유튜브 영상은 자막이 비활성화되어 있습니다.")
    except TranscriptsDisabled:
        raise ValueError("이 유튜브 영상은 자막이 비활성화되어 있습니다.")

    segments: list[dict[str, Any]] = []
    snippets = getattr(transcript, "snippets", None) or []
    for snippet in snippets:
        text = _hard_sanitize(_sanitize(getattr(snippet, "text", "") or "")).replace("\n", " ").strip()
        if not text:
            continue
        start = getattr(snippet, "start", None)
        duration = getattr(snippet, "duration", None)
        if not isinstance(start, (int, float)):
            start = getattr(snippet, "start_sec", 0)
        if not isinstance(duration, (int, float)):
            duration = getattr(snippet, "duration_sec", 0)
        start_sec = float(start or 0)
        end_sec = start_sec + float(duration or 0)
        segments.append({
            "start_sec": max(0.0, start_sec),
            "end_sec": max(start_sec, end_sec),
            "text": text,
        })
    return segments


# ── URL 스크래핑 헬퍼 ────────────────────────────────────────

def _decode_html(raw_bytes: bytes, content_type: str = "") -> str:
    """HTTP 응답 바이트를 올바른 인코딩으로 디코딩."""
    import re as _re
    charset: str | None = None
    for part in content_type.split(";"):
        part = part.strip()
        if part.lower().startswith("charset="):
            charset = part[8:].strip().strip('"\'')
            break
    if not charset:
        head = raw_bytes[:4096].decode("ascii", errors="ignore").lower()
        for pat in [r'charset=["\']?([a-z0-9\-_]+)', r'content=["\'][^"\']*charset=([a-z0-9\-_]+)']:
            m = _re.search(pat, head)
            if m:
                charset = m.group(1)
                break
    if charset:
        try:
            return raw_bytes.decode(charset, errors="replace")
        except (LookupError, UnicodeDecodeError):
            pass
    try:
        return raw_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return raw_bytes.decode("latin-1", errors="replace")


def _text_quality_score(t: str) -> float:
    """추출된 텍스트의 품질을 0~1로 평가."""
    if not t or len(t) < 20:
        return 0.0
    replacement_chars = t.count("\ufffd")
    if replacement_chars / len(t) > 0.05:
        return 0.0
    non_printable = sum(1 for c in t if ord(c) < 32 and c not in "\n\r\t")
    if non_printable / len(t) > 0.03:
        return 0.0
    readable = sum(1 for c in t if c.isalpha() or c.isdigit() or c in " \n\t.,;:!?")
    if readable / len(t) < 0.4:
        return 0.0
    if t.count(" ") / len(t) < 0.04:
        return 0.0
    import re as _re
    js_patterns = [r"\bfunction\s*\(", r"\bvar\s+\w+\s*=", r"document\.getElementById",
                   r"window\.onload", r"jQuery\(", r"\.addEventListener\("]
    js_hits = sum(1 for p in js_patterns if _re.search(p, t))
    if js_hits >= 3:
        return 0.2
    special = sum(1 for c in t if not c.isalnum() and c not in " \n\r\t.,;:!?-_'\"()[]{}/@#$%^&*+=<>/\\|`~")
    score = 1.0 - min(special / max(len(t), 1), 0.5)
    return round(score, 3)


def _clean_text(t: str) -> str:
    """JS 코드 라인, base64, 독립 URL 라인, HTML 엔티티 제거."""
    import re as _re
    import html as _html
    t = _html.unescape(t)
    lines = t.splitlines()
    cleaned = []
    js_line_re = _re.compile(
        r"^\s*(var |let |const |function |if\s*\(|else\s*{|\}|\);|//|/\*|\*/|import |export |return |document\.|window\.|jQuery|\$\()"
    )
    base64_re = _re.compile(r"^[A-Za-z0-9+/]{60,}={0,2}$")
    url_only_re = _re.compile(r"^\s*https?://\S+\s*$")
    for line in lines:
        stripped = line.strip()
        if not stripped:
            cleaned.append("")
            continue
        if js_line_re.match(stripped):
            continue
        if base64_re.match(stripped):
            continue
        if url_only_re.match(stripped):
            continue
        cleaned.append(line)
    return "\n".join(cleaned).strip()


def _is_bot_protected(html: str) -> bool:
    """Cloudflare/DDoS-Guard 차단 페이지 여부 감지."""
    markers = ["_cf_chl_", "jschl_vc", "__ddg", "DDoS-Guard",
               "cf-browser-verification", "challenge-form", "cf_chl_prog"]
    return any(m in html for m in markers)


def _is_spa(html: str) -> bool:
    """Next.js/Nuxt/React 등 SPA 감지."""
    markers = ["__NEXT_DATA__", "__NUXT__", "ng-version=",
               "data-reactroot", "window.__INITIAL_STATE__", "__REACT_QUERY_STATE__"]
    return any(m in html for m in markers)


def _fetch_via_jina(target_url: str) -> str:
    """r.jina.ai를 통해 순수 텍스트 추출 (JS 렌더링, 봇 차단 우회)."""
    import httpx as _httpx
    try:
        jina_url = f"https://r.jina.ai/{target_url}"
        resp = _httpx.get(jina_url, headers={
            "Accept": "text/plain",
            "X-Return-Format": "text",
            "X-With-Generated-Alt": "true",
        }, timeout=45, follow_redirects=True)
        if resp.status_code == 200 and len(resp.text.strip()) > 50:
            return _clean_text(resp.text.strip())
    except Exception as e:
        print(f"[INFO] Jina Reader 실패: {e}")
    return ""


def _fetch_via_google_cache(target_url: str) -> str:
    """Google 캐시에서 텍스트 추출 폴백."""
    import httpx as _httpx
    from bs4 import BeautifulSoup
    try:
        cache_url = f"https://webcache.googleusercontent.com/search?q=cache:{target_url}"
        resp = _httpx.get(cache_url, headers={
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                          "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
        }, timeout=20, follow_redirects=True)
        if resp.status_code == 200:
            soup = BeautifulSoup(resp.text, "html.parser")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
            return _clean_text(text) if len(text) > 100 else ""
    except Exception as e:
        print(f"[INFO] Google Cache 실패: {e}")
    return ""


def ingest_url(url: str, doc_id: str) -> tuple[int, int]:
    """URL에서 텍스트를 추출하여 임베딩과 함께 Supabase document_chunks에 저장.

    Returns:
        (chunk_count, 0)
    Raises:
        ValueError: URL 접근 실패 또는 텍스트 부족
    """
    import httpx
    from urllib.parse import unquote, urlparse

    # ── YouTube 처리 ──
    video_id = _extract_youtube_video_id(url)
    if video_id:
        try:
            segments = _fetch_youtube_transcript_segments(video_id)
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"유튜브 자막 추출 실패: {str(e)}")
        if not segments:
            raise ValueError("유튜브 영상에서 충분한 자막을 찾지 못했습니다.")

        chunks = _build_media_chunks(segments, f"{video_id}.youtube")

        try:
            embeddings = _embed_batch(chunks)
        except Exception as e:
            print(f"[WARNING] 임베딩 생성 실패, 임베딩 없이 저장: {e}")
            embeddings = [None] * len(chunks)  # type: ignore

        records = [
            {
                "doc_id": doc_id,
                "chunk_index": idx,
                "content": content,
                "embedding": emb,
            }
            for idx, (content, emb) in enumerate(zip(chunks, embeddings))
        ]
        _db_safe_insert("document_chunks", records)
        print(f"[INFO] ingest_url: {len(records)} YouTube transcript chunks 저장 완료")
        return len(records), 0
    parsed_url = urlparse(url)
    media_filename = unquote(parsed_url.path.split("/")[-1]) if parsed_url.path else ""
    media_ext = media_filename.rsplit(".", 1)[-1].lower() if "." in media_filename else ""
    if media_ext in VIDEO_AUDIO_EXTENSIONS:
        try:
            response = httpx.get(url, timeout=60, follow_redirects=True)
            response.raise_for_status()
        except httpx.HTTPStatusError as e:
            raise ValueError(f"미디어 URL 접근 실패: HTTP {e.response.status_code}")
        except Exception as e:
            raise ValueError(f"미디어 URL 접근 실패: {str(e)}")

        file_bytes = response.content or b""
        if not file_bytes:
            raise ValueError("미디어 URL에서 파일을 불러오지 못했습니다.")
        if len(file_bytes) > 25 * 1024 * 1024:
            raise ValueError("비디오/오디오 URL 파일은 25MB 이하일 때만 전사할 수 있습니다.")

        filename = media_filename or f"url-media.{media_ext}"
        segments, _ = _extract_media_transcript(file_bytes, filename)
        if not segments:
            raise ValueError("비디오/오디오 URL에서 전사할 수 있는 음성 내용을 찾지 못했습니다.")

        chunks = _build_media_chunks(segments, filename)

        try:
            embeddings = _embed_batch(chunks)
        except Exception as e:
            print(f"[WARNING] 임베딩 생성 실패, 임베딩 없이 저장: {e}")
            embeddings = [None] * len(chunks)  # type: ignore

        records = [
            {
                "doc_id": doc_id,
                "chunk_index": idx,
                "content": content,
                "embedding": emb,
            }
            for idx, (content, emb) in enumerate(zip(chunks, embeddings))
        ]
        _db_safe_insert("document_chunks", records)
        print(f"[INFO] ingest_url: {len(records)} media transcript chunks 저장 완료")
        return len(records), 0
    else:
        # ── 일반 웹페이지 처리 ──
        from bs4 import BeautifulSoup
        from urllib.parse import urljoin, urlparse

        req_headers = {
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/122.0.0.0 Safari/537.36"
            ),
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8",
            "Accept-Language": "ko-KR,ko;q=0.9,en-US;q=0.8,en;q=0.7",
            "Accept-Encoding": "gzip, deflate, br",
            "Connection": "keep-alive",
            "Upgrade-Insecure-Requests": "1",
            "Sec-Fetch-Dest": "document",
            "Sec-Fetch-Mode": "navigate",
            "Sec-Fetch-Site": "none",
        }

        # ── 1단계: 직접 HTTP 요청 ──
        html_content = ""
        text = ""
        skip_to_jina = False
        fetch_failed = False

        try:
            response = httpx.get(url, headers=req_headers, timeout=30, follow_redirects=True)
            response.raise_for_status()
            raw = response.content
            # 바이너리 감지: null byte > 0.5% → Jina로 건너뜀
            null_ratio = raw.count(b"\x00") / max(len(raw), 1)
            if null_ratio > 0.005:
                print(f"[INFO] ingest_url: 바이너리 응답 감지 ({null_ratio:.3f}), Jina로 전환")
                skip_to_jina = True
            else:
                html_content = _decode_html(raw, response.headers.get("content-type", ""))
                # 봇 차단 / SPA 감지
                if _is_bot_protected(html_content):
                    print(f"[INFO] ingest_url: 봇 차단 페이지 감지, Jina로 전환")
                    skip_to_jina = True
                elif _is_spa(html_content):
                    print(f"[INFO] ingest_url: SPA 감지, Jina로 전환")
                    skip_to_jina = True
        except httpx.HTTPStatusError as e:
            print(f"[INFO] ingest_url: HTTP {e.response.status_code}, Jina로 전환")
            fetch_failed = True
        except Exception as e:
            print(f"[INFO] ingest_url: 직접 요청 실패 ({e}), Jina로 전환")
            fetch_failed = True

        def _extract_meta_local(html: str) -> str:
            try:
                soup = BeautifulSoup(html, "html.parser")
                parts = []
                title = soup.find("title")
                if title and title.get_text(strip=True):
                    parts.append(title.get_text(strip=True))
                for attr in [("name", "description"), ("property", "og:description"),
                              ("property", "og:title"), ("name", "keywords")]:
                    tag = soup.find("meta", attrs={attr[0]: attr[1]})
                    if tag and tag.get("content", "").strip():
                        parts.append(tag["content"].strip())
                return "\n".join(parts)
            except Exception:
                return ""

        def _extract_body_local(html: str) -> str:
            """HTML에서 본문 텍스트 추출. trafilatura → 콘텐츠 셀렉터 → p태그 → 전체 순."""
            t = ""
            try:
                import trafilatura
                t = trafilatura.extract(html, include_tables=True, favor_recall=True) or ""
            except Exception:
                pass
            if len(t.strip()) < 100:
                try:
                    soup = BeautifulSoup(html, "html.parser")
                    for tag in soup(["script", "style", "nav", "footer", "header", "aside", "menu"]):
                        tag.decompose()
                    content_text = ""
                    for selector in [
                        ".se-main-container", "#postViewArea", ".se_component_wrap",
                        "article", "main", "[role='main']",
                        ".content", "#content", ".post", ".article",
                        ".entry-content", "#article-body", ".post-content",
                        ".tt_article_useless_p_margin", "#article-view",
                    ]:
                        nodes = soup.select(selector)
                        if nodes:
                            content_text = "\n".join(n.get_text(separator="\n", strip=True) for n in nodes)
                            if len(content_text) > 200:
                                break
                    if len(content_text) < 200:
                        paragraphs = [p.get_text(strip=True) for p in soup.find_all("p") if len(p.get_text(strip=True)) > 20]
                        content_text = "\n".join(paragraphs)
                    if len(content_text) < 200:
                        content_text = soup.get_text(separator="\n", strip=True)
                    if len(content_text) > len(t):
                        t = content_text
                except Exception:
                    pass
            return t.strip()

        def _follow_iframe_local(html: str, base_url: str) -> str:
            try:
                soup = BeautifulSoup(html, "html.parser")
                for iframe in soup.find_all("iframe"):
                    src = iframe.get("src", "")
                    if not src or src.startswith("javascript"):
                        continue
                    iframe_url = urljoin(base_url, src)
                    base_host = urlparse(base_url).netloc
                    iframe_host = urlparse(iframe_url).netloc
                    if base_host not in iframe_host and iframe_host not in base_host:
                        continue
                    try:
                        iframe_resp = httpx.get(iframe_url, headers={**req_headers, "Referer": base_url}, timeout=20, follow_redirects=True)
                        if iframe_resp.status_code == 200:
                            t = _extract_body_local(iframe_resp.text)
                            if len(t) > 100:
                                return t
                    except Exception:
                        continue
            except Exception:
                pass
            return ""

        # 직접 추출 시도 (Jina로 건너뛰지 않는 경우)
        if not skip_to_jina and not fetch_failed and html_content:
            meta_text = _extract_meta_local(html_content)
            text = _extract_body_local(html_content)
            if len(text) < 100:
                iframe_text = _follow_iframe_local(html_content, url)
                if len(iframe_text) > len(text):
                    text = iframe_text
            if len(text) < 100:
                try:
                    import trafilatura
                    downloaded = trafilatura.fetch_url(url)
                    if downloaded:
                        t2 = trafilatura.extract(downloaded, favor_recall=True) or ""
                        if len(t2) > len(text):
                            text = t2
                except Exception:
                    pass
            if meta_text:
                text = (meta_text + "\n\n" + text).strip()

        # ── 2단계: 품질 부족하면 Jina Reader ──
        quality = _text_quality_score(text)
        if skip_to_jina or fetch_failed or len(text) < 200 or quality < 0.4:
            print(f"[INFO] ingest_url: Jina Reader 시도 (text={len(text)}자, quality={quality})")
            jina_text = _fetch_via_jina(url)
            jina_quality = _text_quality_score(jina_text)
            if jina_text and (jina_quality > quality or len(jina_text) > len(text) * 1.2):
                print(f"[INFO] ingest_url: Jina 결과 사용 ({len(jina_text)}자, quality={jina_quality})")
                text = jina_text

        # ── 3단계: 여전히 부족하면 Google Cache ──
        if len(text) < 100:
            print(f"[INFO] ingest_url: Google Cache 폴백 시도")
            cache_text = _fetch_via_google_cache(url)
            if len(cache_text) > len(text):
                text = cache_text

        text = _clean_text(text)

    text = _hard_sanitize(_sanitize(text)).strip()
    if len(text) < 50:
        raise ValueError(
            "URL에서 충분한 텍스트를 추출할 수 없습니다. "
            "로그인이 필요하거나, JavaScript로 렌더링되거나, 스크래핑을 차단하는 페이지일 수 있습니다. "
            f"(추출된 텍스트: {len(text)}자)"
        )

    # 청킹
    chunks: list[str] = []
    step = CHUNK_SIZE - CHUNK_OVERLAP
    for i in range(0, len(text), step):
        chunk = _hard_sanitize(_sanitize(text[i: i + CHUNK_SIZE]))
        if chunk.strip():
            chunks.append(chunk)

    if not chunks:
        raise ValueError("텍스트 청킹 실패")

    # 임베딩 생성
    try:
        embeddings = _embed_batch(chunks)
    except Exception as e:
        print(f"[WARNING] 임베딩 생성 실패, 임베딩 없이 저장: {e}")
        embeddings = [None] * len(chunks)  # type: ignore

    records = [
        {
            "doc_id": doc_id,
            "chunk_index": idx,
            "content": content,
            "embedding": emb,
        }
        for idx, (content, emb) in enumerate(zip(chunks, embeddings))
    ]
    _db_safe_insert("document_chunks", records)
    print(f"[INFO] ingest_url: {len(records)} chunks 저장 완료")
    return len(chunks), 0


# ─────────────────────────────────────────────
# 컨텍스트 조회
# ─────────────────────────────────────────────

def _get_context_semantic(
    doc_ids: list[str],
    question: str,
    top_k: int = TOP_K,
    labeled: bool = False,
) -> tuple[str, list[dict]]:
    """임베딩 기반 코사인 유사도로 각 문서에서 top-k 청크를 선택하여 컨텍스트 구성.
    (context_str, source_chunks) 튜플 반환."""
    if not doc_ids or not question:
        return _get_context(doc_ids), []

    # 질문 임베딩
    try:
        query_embedding = _embed_batch([question])[0]
    except Exception as e:
        print(f"[WARNING] 질문 임베딩 실패, 순차 검색으로 fallback: {e}")
        return _get_context(doc_ids, labeled=labeled), []

    # 문서명 조회 (항상 수행 — 출처 추적용)
    doc_name_map: dict[str, str] = {}
    try:
        name_result = supabase_admin.table("documents").select("id, filename").in_("id", doc_ids).execute()
        doc_name_map = {row["id"]: row["filename"] for row in name_result.data}
    except Exception:
        pass

    parts: list[str] = []
    source_chunks: list[dict] = []

    for doc_id in doc_ids:
        # 해당 문서의 청크 + 임베딩 조회
        try:
            result = (
                supabase_admin.table("document_chunks")
                .select("content, embedding, chunk_index")
                .eq("doc_id", doc_id)
                .order("chunk_index")
                .execute()
            )
        except Exception:
            continue

        rows = result.data
        if not rows:
            continue

        # 임베딩이 있는 행만 유사도 검색, 없으면 순서대로
        rows_with_emb = [r for r in rows if r.get("embedding")]
        if rows_with_emb:
            embeddings = [r["embedding"] for r in rows_with_emb]
            scores = _cosine_similarity(query_embedding, embeddings)
            top_indices = np.argsort(scores)[::-1][:top_k]
            # 원래 chunk_index 순서로 정렬
            selected_with_idx = sorted(
                [(rows_with_emb[i]["chunk_index"], rows_with_emb[i]["content"]) for i in top_indices],
                key=lambda x: x[0],
            )
            selected_chunks = [_strip_media_metadata(c) for _, c in selected_with_idx]
        else:
            # 임베딩 없으면 앞에서 top_k개
            selected_with_idx = [(r["chunk_index"], _strip_media_metadata(r["content"])) for r in rows[:top_k]]
            selected_chunks = [text for _, text in selected_with_idx]

        if not selected_chunks:
            continue

        # 출처 청크 메타데이터 수집 (전역 1-based 번호 부여)
        doc_filename = doc_name_map.get(doc_id, doc_id)

        # /chunks 엔드포인트는 overlap(100자) 제거 후 "\n\n".join 으로 텍스트를 복원.
        # 각 청크의 char_offset 을 동일한 방식으로 계산한다.
        running = 0
        offset_map: dict[int, int] = {}
        length_map: dict[int, int] = {}  # overlap 제거 후 실제 표시 길이
        for j, row in enumerate(rows):
            ci = row["chunk_index"]
            content: str = row["content"]
            display = content if j == 0 else (content[CHUNK_OVERLAP:] if len(content) > CHUNK_OVERLAP else content)
            offset_map[ci] = running
            length_map[ci] = len(display)
            running += len(display)
            # separator 없음 — documents.py 와 동일하게 ""로 join

        numbered_parts: list[str] = []
        for ci, text in selected_with_idx:
            num = len(source_chunks) + 1
            source_chunks.append({
                "num": num,
                "doc_id": doc_id,
                "filename": doc_filename,
                "chunk_index": ci,
                "text": text[:200],  # 툴팁 미리보기용
                "full_text": text,   # evidence extraction용 (응답 전 제거됨)
                "char_offset": offset_map.get(ci, -1),
                "char_length": length_map.get(ci, len(text)),  # AI가 실제 읽은 청크 전체 범위
            })
            numbered_parts.append(f"[출처 {num}]\n{text}")

        doc_text = "\n\n".join(numbered_parts)
        if labeled:
            doc_name = doc_name_map.get(doc_id, doc_id)
            parts.append(f"[문서명: {doc_name}]\n{doc_text}")
        else:
            parts.append(doc_text)

    separator = "\n\n=====\n\n" if labeled else "\n\n"
    return separator.join(parts), source_chunks


def _get_semantic_rows(doc_ids: list[str], question: str, top_k: int = TOP_K) -> list[dict[str, Any]]:
    if not doc_ids or not question:
        return []

    try:
        query_embedding = _embed_batch([question])[0]
    except Exception:
        query_embedding = None

    name_result = supabase_admin.table("documents").select("id, filename").in_("id", doc_ids).execute()
    doc_name_map = {row["id"]: row["filename"] for row in (name_result.data or [])}

    selected: list[dict[str, Any]] = []
    for doc_id in doc_ids:
        result = (
            supabase_admin.table("document_chunks")
            .select("content, embedding, chunk_index")
            .eq("doc_id", doc_id)
            .order("chunk_index")
            .execute()
        )
        rows = result.data or []
        if not rows:
            continue

        chosen_rows: list[dict[str, Any]]
        if query_embedding is not None:
            rows_with_emb = [r for r in rows if r.get("embedding")]
            if rows_with_emb:
                embeddings = [r["embedding"] for r in rows_with_emb]
                scores = _cosine_similarity(query_embedding, embeddings)
                top_indices = np.argsort(scores)[::-1][:top_k]
                chosen_rows = [rows_with_emb[i] for i in top_indices]
                chosen_rows.sort(key=lambda item: item.get("chunk_index", 0))
            else:
                chosen_rows = rows[:top_k]
        else:
            chosen_rows = rows[:top_k]

        for row in chosen_rows:
            selected.append({
                "doc_id": doc_id,
                "filename": doc_name_map.get(doc_id, doc_id),
                "chunk_index": row.get("chunk_index", 0),
                "total_chunks": len(rows),
                "content": row.get("content", ""),
            })
    return selected


def _build_media_references(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    references: list[dict[str, Any]] = []
    for row in rows:
        meta = _parse_media_metadata(row.get("content", ""))
        references.append({
            "doc_id": row.get("doc_id", ""),
            "filename": row.get("filename", ""),
            "chunk_index": row.get("chunk_index", 0),
            "total_chunks": row.get("total_chunks", 0),
            "start_sec": meta.get("start_sec"),
            "end_sec": meta.get("end_sec"),
            "excerpt": _strip_media_metadata(row.get("content", ""))[:320],
        })
    return references


def _format_reference_lines(references: list[dict[str, Any]], limit: int = 6) -> str:
    lines: list[str] = []
    for ref in references[:limit]:
        start_sec = ref.get("start_sec")
        end_sec = ref.get("end_sec")
        if not isinstance(start_sec, (int, float)):
            continue
        if isinstance(end_sec, (int, float)) and end_sec > start_sec:
            time_label = f"{_format_media_timestamp(start_sec)} - {_format_media_timestamp(end_sec)}"
        else:
            time_label = _format_media_timestamp(start_sec)
        lines.append(f"- {time_label}: {ref.get('excerpt', '')}")
    return "\n".join(lines)


def _get_ppt_full_context(doc_ids: list[str], max_chars: int = 48000) -> str:
    """PPT 전체 질문용: 슬라이드별 청크를 균등하게 샘플링해서 반환.
    앞 슬라이드에 치우치지 않도록 전체 슬라이드를 고르게 포함."""
    if not doc_ids:
        return ""
    result = (
        supabase_admin.table("document_chunks")
        .select("content, chunk_index")
        .in_("doc_id", doc_ids)
        .order("chunk_index")
        .execute()
    )
    all_chunks = [row["content"] for row in result.data if row.get("content", "").strip()]
    if not all_chunks:
        return ""

    total = len(all_chunks)
    # 전체가 max_chars 이하면 그냥 전부 반환
    full_text = "\n\n".join(all_chunks)
    if len(full_text) <= max_chars:
        return full_text

    # 균등 샘플링: 전체에서 고르게 N개 선택
    avg_len = sum(len(c) for c in all_chunks) / total
    n_samples = max(1, int(max_chars / max(avg_len, 1)))
    n_samples = min(n_samples, total)

    if n_samples >= total:
        return full_text[:max_chars]

    step = total / n_samples
    sampled = [all_chunks[int(i * step)] for i in range(n_samples)]
    return "\n\n".join(sampled)[:max_chars]


def _get_pdf_full_context(doc_ids: list[str], max_chars: int = 48000) -> str:
    """PDF 전체 질문용: 페이지별 청크를 균등하게 샘플링해서 반환."""
    if not doc_ids:
        return ""
    result = (
        supabase_admin.table("document_chunks")
        .select("content, chunk_index")
        .in_("doc_id", doc_ids)
        .order("chunk_index")
        .execute()
    )
    all_chunks = [row["content"] for row in result.data if row.get("content", "").strip()]
    if not all_chunks:
        return ""

    full_text = "\n\n".join(all_chunks)
    if len(full_text) <= max_chars:
        return full_text

    avg_len = sum(len(c) for c in all_chunks) / len(all_chunks)
    n_samples = max(1, int(max_chars / max(avg_len, 1)))
    n_samples = min(n_samples, len(all_chunks))

    step = len(all_chunks) / n_samples
    sampled = [all_chunks[int(i * step)] for i in range(n_samples)]
    return "\n\n".join(sampled)[:max_chars]


def _get_context(doc_ids: list[str], max_chars: int = 10000, labeled: bool = False) -> str:
    """순차적으로 청크를 가져와 컨텍스트 문자열 반환 (임베딩 없는 fallback)."""
    if not doc_ids:
        return ""

    if labeled:
        # 문서별 레이블 붙이기
        try:
            name_result = supabase_admin.table("documents").select("id, filename").in_("id", doc_ids).execute()
            doc_name_map = {row["id"]: row["filename"] for row in name_result.data}
        except Exception:
            doc_name_map = {}

        parts = []
        per_doc_chars = max_chars // max(len(doc_ids), 1)
        for doc_id in doc_ids:
            result = (
                supabase_admin.table("document_chunks")
                .select("content")
                .eq("doc_id", doc_id)
                .order("chunk_index")
                .execute()
            )
            chunks = [_strip_media_metadata(row["content"]) for row in result.data]
            if not chunks:
                continue
            doc_text = "\n\n".join(chunks)[:per_doc_chars]
            doc_name = doc_name_map.get(doc_id, doc_id)
            parts.append(f"[문서명: {doc_name}]\n{doc_text}")
        return "\n\n=====\n\n".join(parts)

    result = (
        supabase_admin.table("document_chunks")
        .select("content")
        .in_("doc_id", doc_ids)
        .order("doc_id")
        .order("chunk_index")
        .execute()
    )
    chunks = [_strip_media_metadata(row["content"]) for row in result.data]
    return "\n\n".join(chunks)[:max_chars]


def _get_filenames(doc_ids: list[str]) -> list[str]:
    if not doc_ids:
        return []
    result = supabase_admin.table("documents").select("filename").in_("id", doc_ids).execute()
    return [row["filename"] for row in result.data]


def _get_doc_filename_map(doc_ids: list[str]) -> dict[str, str]:
    if not doc_ids:
        return {}
    result = supabase_admin.table("documents").select("id, filename").in_("id", doc_ids).execute()
    return {row["id"]: row["filename"] for row in (result.data or [])}


def _get_media_doc_ids(doc_ids: list[str]) -> list[str]:
    if not doc_ids:
        return []
    result = supabase_admin.table("documents").select("id, filename, storage_path").in_("id", doc_ids).execute()
    media_doc_ids: list[str] = []
    for row in (result.data or []):
        filename = row.get("filename", "") or ""
        storage_path = row.get("storage_path", "") or ""
        ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
        if ext in VIDEO_AUDIO_EXTENSIONS or _is_youtube_url(storage_path):
            media_doc_ids.append(row["id"])
    return media_doc_ids


def _get_youtube_doc_ids(doc_ids: list[str]) -> list[str]:
    if not doc_ids:
        return []
    result = supabase_admin.table("documents").select("id, storage_path").in_("id", doc_ids).execute()
    youtube_doc_ids: list[str] = []
    for row in (result.data or []):
        storage_path = row.get("storage_path", "") or ""
        if _is_youtube_url(storage_path):
            youtube_doc_ids.append(row["id"])
    return youtube_doc_ids


# ─────────────────────────────────────────────
# 챗
# ─────────────────────────────────────────────

LEVEL_PROMPTS = {
    "beginner": "쉽고 친절하게, 예시를 들어 입문자 수준으로 설명해주세요.",
    "intermediate": "핵심 개념 위주로 중급 학습자에게 적합하게 설명해주세요.",
    "advanced": "심화 분석과 비판적 관점을 포함하여 전문가 수준으로 설명해주세요.",
}


_IMAGE_MIME_MAP = {
    "jpg": "image/jpeg", "jpeg": "image/jpeg",
    "png": "image/png", "gif": "image/gif", "webp": "image/webp",
}

def _get_image_docs(doc_ids: list[str]) -> list[dict]:
    """doc_ids 중 이미지 파일인 것의 메타데이터 반환.
    반환: [{"doc_id", "filename", "storage_path", "mime_type"}]
    """
    if not doc_ids:
        return []
    result = supabase_admin.table("documents") \
        .select("id, filename, storage_path") \
        .in_("id", doc_ids).execute()
    image_docs = []
    for row in result.data:
        fn = row.get("filename", "")
        ext = fn.lower().rsplit(".", 1)[-1] if "." in fn else ""
        if ext in IMAGE_EXTENSIONS and row.get("storage_path"):
            image_docs.append({
                "doc_id": row["id"],
                "filename": fn,
                "storage_path": row["storage_path"],
                "mime_type": _IMAGE_MIME_MAP.get(ext, "image/jpeg"),
            })
    return image_docs


def _check_image_relevance(question: str, subjects: list[str], visual_context: str) -> bool:
    """질문이 이미지 대상/내용과 관련 있는지 YES/NO로 판단. 관련 있으면 True."""
    # ── Python 사전 판단 (API 호출 없음) ──────────────
    # 지시어("이 개", "이 사진" 등) 또는 시각적 속성 키워드가 있으면 무조건 관련 있음
    if _IMAGE_DIRECT_REF_RE.search(question):
        return True
    # 메타데이터 없으면 관련 있다고 가정 (안전 기본값)
    if not subjects and not visual_context:
        return True

    subjects_str = ", ".join(subjects) if subjects else "불명확"
    prompt = f"""이미지 정보:
- 등장 대상: {subjects_str}
- 상황: {visual_context}

사용자 질문: "{question}"

판단 기준:
관련 있음(YES): 아래 중 하나라도 해당
  - 이미지 대상의 외형·색상·자세·표정·크기 등 시각적 속성
  - 이미지를 직접 가리키는 표현 ("이 개", "이 사진", "여기서" 등)
  - 이미지 대상의 특성·습성·행동·수명·먹이 등 관련 지식
  - 이미지 대상이 입고 있는 것, 들고 있는 것, 주변 환경
  - 이미지 주제와 맥락적으로 연결되는 배경지식
  - 판단이 애매한 경우

관련 없음(NO): 이미지 속 어떤 대상과도 완전히 무관한 다른 주제
  (예: 개 사진 → "파이썬 코드 짜줘" / 음식 사진 → "주식 투자 알려줘")

"YES" 또는 "NO"만 대답하세요."""

    try:
        answer = _call_llm(
            messages=[{"role": "user", "content": prompt}],
            model=_DEFAULT_MODEL,
            max_tokens=5,
            temperature=0,
        ).strip().upper()
        return answer != "NO"
    except Exception:
        return True  # 오류 시 관련 있다고 가정


def _get_image_metadata(doc_ids: list[str]) -> dict[str, dict]:
    """doc_ids에 해당하는 이미지 메타데이터를 청크에서 파싱해 반환.
    반환: {doc_id: {"subjects": [...], "image_type": "...", "visual_context": "...", "extracted_text": "..."}}
    """
    if not doc_ids:
        return {}
    import re
    result = supabase_admin.table("document_chunks") \
        .select("doc_id, content") \
        .in_("doc_id", doc_ids) \
        .execute()
    meta: dict[str, dict] = {}
    for row in result.data:
        doc_id = row["doc_id"]
        content = row.get("content", "")
        m = re.search(r"\[IMAGE_META\](.*?)\[/IMAGE_META\](.*)", content, re.DOTALL)
        if m and doc_id not in meta:
            try:
                parsed = json.loads(m.group(1))
                # [IMAGE_META] 태그 이후의 본문에서 첫 단락(extracted_text)을 추출
                # 저장 형식: [IMAGE_META]...[/IMAGE_META]\n\n{extracted_text}\n\n{description}
                body = m.group(2).strip()
                extracted_text = body.split("\n\n")[0].strip() if body else ""
                parsed["extracted_text"] = extracted_text
                meta[doc_id] = parsed
            except Exception:
                pass
    return meta


def _download_image_b64(storage_path: str) -> bytes | None:
    """Supabase Storage에서 이미지를 다운로드해 bytes 반환. 실패 시 None."""
    try:
        return supabase_admin.storage.from_("documents").download(storage_path)
    except Exception:
        return None


def _extract_evidence_passages(
    answer: str,
    source_chunks: list[dict],
    client: "OpenAI",
    model: str,
) -> None:
    """
    LLM 답변에서 인용된 각 청크의 정확한 근거 구절을 추출하여
    char_offset / char_length 를 in-place 로 정밀화한다.
    각 청크별 LLM 호출은 ThreadPoolExecutor 로 병렬 실행.
    """
    cited_nums = set(int(m) for m in re.findall(r'\[(\d+)\]', answer))

    def extract_one(chunk: dict) -> None:
        num = chunk["num"]
        if num not in cited_nums:
            return

        full_text: str = chunk.get("full_text", "")
        if not full_text:
            return

        # [N] 바로 앞 문장을 컨텍스트로 수집 (최대 200자)
        pattern = re.compile(r'([^\n]{5,200})\s*\[' + str(num) + r'\]')
        m = pattern.search(answer)
        context_sentence = m.group(1).strip() if m else ""

        try:
            print(f"[evidence] num={num} context_sentence={context_sentence[:80]!r}")
            resp = client.chat.completions.create(
                model=model,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "아래 【문서 청크】에서 【답변 문장】의 근거가 된 구절을 찾아, "
                            "문서 청크 원문에서 그대로 복사하여 출력하세요. "
                            "변형하거나 요약하지 말고, 해당 구절만 출력하세요."
                        ),
                    },
                    {
                        "role": "user",
                        "content": (
                            f"【답변 문장】\n{context_sentence}\n\n"
                            f"【문서 청크】\n{full_text}"
                        ),
                    },
                ],
                temperature=0.0,
                max_tokens=300,
            )
            extracted = (resp.choices[0].message.content or "").strip().strip('"').strip("'").strip()
            print(f"[evidence] num={num} extracted={extracted[:80]!r} len={len(extracted)}")

            if not extracted or len(extracted) < 5:
                print(f"[evidence] num={num} extracted too short, skipping")
                return

            # 청크 내에서 해당 구절의 시작 위치 탐색
            chunk_start: int = chunk["char_offset"]
            ci: int = chunk.get("chunk_index", 0)
            overlap_adj = 0 if ci == 0 else CHUNK_OVERLAP

            search_key = extracted[:60]
            idx = full_text.lower().find(search_key.lower())
            print(f"[evidence] num={num} idx={idx} chunk_start={chunk_start} ci={ci} overlap_adj={overlap_adj}")
            if idx >= 0:
                doc_pos = chunk_start + (idx - overlap_adj)
                if doc_pos < 0:
                    doc_pos = chunk_start
                chunk["char_offset"] = doc_pos
                chunk["char_length"] = min(len(extracted), len(full_text) - idx)
                print(f"[evidence] num={num} → offset={doc_pos} length={chunk['char_length']}")
        except Exception as e:
            print(f"[evidence] num={num} EXCEPTION: {e}")

    with ThreadPoolExecutor(max_workers=5) as executor:
        futures = [executor.submit(extract_one, chunk) for chunk in source_chunks]
        for f in as_completed(futures):
            try:
                f.result()
            except Exception:
                pass


def _get_slide_image_b64(doc_id: str, slide_num: int) -> str | None:
    """Supabase Storage에서 슬라이드 이미지를 다운로드해 base64로 반환."""
    import base64 as _b64
    path = f"slide-assets/{doc_id}/{slide_num}.webp"
    try:
        raw = supabase_admin.storage.from_("documents").download(path)
        if raw:
            return _b64.b64encode(raw).decode("utf-8")
    except Exception:
        pass
    return None


def chat_with_docs(
    doc_ids: list[str],
    question: str,
    model: str = _DEFAULT_MODEL,
    level: str = "intermediate",
    chat_history: list | None = None,
    current_slide: int | None = None,
    stream: bool = False,
    asked_questions: list[str] | None = None,
) -> tuple[str, list[str], list[dict[str, Any]]]:
    """문서 기반 RAG 질의응답. (answer, sources, references) 반환.
    stream=True 이면 (generator, [], []) 를 반환하며 generator가 SSE 문자열을 yield."""
    import base64

    is_multi = len(doc_ids) > 1
    level_hint = LEVEL_PROMPTS.get(level, LEVEL_PROMPTS["intermediate"])
    filename_map = _get_doc_filename_map(doc_ids)
    doc_filenames = [filename_map.get(doc_id, doc_id) for doc_id in doc_ids]

    # URL 소스 / PPT 여부 확인
    try:
        doc_meta_result = supabase_admin.table("documents").select("id, filename, storage_path").in_("id", doc_ids).execute()
        doc_meta_map = {row["id"]: row for row in doc_meta_result.data}
    except Exception:
        doc_meta_map = {}
    is_url_doc = any(
        str(doc_meta_map.get(d, {}).get("storage_path", "")).startswith("http")
        for d in doc_ids
    )
    is_ppt_doc = (not is_url_doc) and any(
        str(doc_meta_map.get(d, {}).get("filename", "")).lower().rsplit(".", 1)[-1] in ("pptx", "ppt")
        for d in doc_ids
    )
    is_pdf_doc = (not is_url_doc) and (not is_ppt_doc) and any(
        str(doc_meta_map.get(d, {}).get("filename", "")).lower().rsplit(".", 1)[-1] == "pdf"
        for d in doc_ids
    )
    is_docx_doc = (not is_url_doc) and (not is_ppt_doc) and (not is_pdf_doc) and any(
        str(doc_meta_map.get(d, {}).get("filename", "")).lower().rsplit(".", 1)[-1] == "docx"
        for d in doc_ids
    )
    is_hwpx_doc = (not is_url_doc) and (not is_ppt_doc) and (not is_pdf_doc) and (not is_docx_doc) and any(
        str(doc_meta_map.get(d, {}).get("filename", "")).lower().rsplit(".", 1)[-1] in ("hwpx", "hwp")
        for d in doc_ids
    )
    # 전체 문서 조회 의도 감지 (PPT / PDF / DOCX / HWPX 공통)
    _FULL_DOC_RE = re.compile(
        r"(전체|모든\s*슬라이드|슬라이드\s*별|슬라이드별|처음부터|이어서|다음\s*슬라이드|전부|순서대로|차례로|전체적으로\s*설명|모두\s*설명"
        r"|핵심\s*개념|주요\s*개념|핵심\s*내용|주요\s*내용|요약|정리|개요|전반적|전반|내용\s*설명|설명해\s*줘|설명해줘"
        r"|무엇|무슨\s*내용|어떤\s*내용|다루|소개|강의\s*내용|자료\s*내용|ppt\s*내용|pdf\s*내용|전반적\s*내용"
        r"|페이지\s*별|페이지별|모든\s*페이지|각\s*페이지|docx\s*내용|hwpx\s*내용|hwp\s*내용)",
        re.IGNORECASE,
    )
    is_full_doc_query = (is_ppt_doc or is_pdf_doc or is_docx_doc or is_hwpx_doc) and bool(_FULL_DOC_RE.search(question))

    media_doc_ids = _get_media_doc_ids(doc_ids)
    youtube_doc_ids = _get_youtube_doc_ids(doc_ids)
    youtube_only = False
    semantic_rows = _get_semantic_rows(doc_ids, question, top_k=TOP_K)
    references = _build_media_references(semantic_rows)

    # ── 이미지 문서 처리 ──────────────────────────
    image_docs = _get_image_docs(doc_ids)
    img_meta_map = _get_image_metadata([img["doc_id"] for img in image_docs])
    non_image_ids = [
        d for d in doc_ids
        if d not in {img["doc_id"] for img in image_docs}
    ]
    has_images = len(image_docs) > 0
    has_non_image = len(non_image_ids) > 0

    # 이미지 base64 수집 (최대 1568px 리사이즈 + JPEG 압축 후 전송)
    # OpenAI detail:high는 512px 타일 단위 과금이므로 원본 고해상도는 토큰 낭비
    image_parts: list[dict] = []
    for img in image_docs:
        raw = _download_image_b64(img["storage_path"])
        if raw:
            b64, mime = _compress_image_to_base64(raw, max_width=1568, quality=85)
            image_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"},
            })

    # ── 컨텍스트 구성 ─────────────────────────────
    context_ids = non_image_ids if has_non_image else doc_ids
    media_context_ids = [doc_id for doc_id in context_ids if doc_id in media_doc_ids]
    semantic_context_ids = [doc_id for doc_id in context_ids if doc_id not in media_context_ids]

    context_parts: list[str] = []
    source_chunks: list[dict[str, Any]] = []
    if media_context_ids:
        media_context = _get_context(media_context_ids, max_chars=24000, labeled=is_multi)
        if media_context:
            context_parts.append(media_context)
    if semantic_context_ids:
        if is_full_doc_query and is_ppt_doc:
            # PPT 전체 슬라이드 요청: 슬라이드별 균등 샘플링 (치우침 방지)
            semantic_context = _get_ppt_full_context(semantic_context_ids)
            source_chunks = []
        elif is_full_doc_query and is_pdf_doc:
            # PDF 전체 페이지 요청: 페이지별 균등 샘플링
            semantic_context = _get_pdf_full_context(semantic_context_ids)
            source_chunks = []
        elif is_full_doc_query and is_docx_doc:
            # DOCX 전체 페이지 요청: 페이지별 균등 샘플링 (PDF와 동일 방식)
            semantic_context = _get_pdf_full_context(semantic_context_ids)
            source_chunks = []
        elif is_full_doc_query and is_hwpx_doc:
            # HWPX 전체 페이지 요청: 페이지별 균등 샘플링 (PDF와 동일 방식)
            semantic_context = _get_pdf_full_context(semantic_context_ids)
            source_chunks = []
        else:
            semantic_context, source_chunks = _get_context_semantic(semantic_context_ids, question, top_k=TOP_K, labeled=is_multi)
            # DOCX/HWPX는 [페이지 N] 마커 기반 인용 → chunk [1][2] 인용 제거
            if is_docx_doc or is_hwpx_doc:
                source_chunks = []
            # 현재 슬라이드 청크를 컨텍스트 앞에 보강
            if is_ppt_doc and current_slide:
                try:
                    slide_marker = f"[슬라이드 {current_slide}]"
                    slide_rows = (
                        supabase_admin.table("document_chunks")
                        .select("chunk_index, content")
                        .in_("doc_id", semantic_context_ids)
                        .ilike("content", f"%{slide_marker}%")
                        .limit(3)
                        .execute()
                    )
                    slide_texts = [r["content"] for r in (slide_rows.data or []) if r.get("content")]
                    if slide_texts:
                        slide_context = f"【현재 슬라이드 {current_slide} 내용】\n" + "\n---\n".join(slide_texts)
                        semantic_context = slide_context + "\n\n" + semantic_context if semantic_context else slide_context
                except Exception:
                    pass
            if is_pdf_doc and current_slide:
                try:
                    page_marker = f"[페이지 {current_slide}]"
                    page_rows = (
                        supabase_admin.table("document_chunks")
                        .select("chunk_index, content")
                        .in_("doc_id", semantic_context_ids)
                        .ilike("content", f"%{page_marker}%")
                        .limit(3)
                        .execute()
                    )
                    page_texts = [r["content"] for r in (page_rows.data or []) if r.get("content")]
                    if page_texts:
                        page_context = f"【현재 페이지 {current_slide} 내용】\n" + "\n---\n".join(page_texts)
                        semantic_context = page_context + "\n\n" + semantic_context if semantic_context else page_context
                except Exception:
                    pass
        if semantic_context:
            context_parts.append(semantic_context)

    context = "\n\n=====\n\n".join(part for part in context_parts if part)

    # 비디오/오디오 여부
    has_video = bool(media_doc_ids)
    video_note = ""
    if has_video:
        joined_ref_lines = _format_reference_lines(references)
        if youtube_only:
            video_note = (
                "\n\n중요: 이 문서는 유튜브 자막 기반 문서입니다. "
                "1차 답변에서는 먼저 자막 텍스트만 읽고 질문에 정확히 답하세요. "
                "이 단계에서는 타임스탬프를 억지로 넣지 말고, 질문에 대한 설명만 자연스럽게 정리하세요. "
                "추측하지 말고 자막에 근거가 있는 내용만 답하세요. "
                "답변은 Markdown 형식으로 작성하고, 핵심 개념·중요 결론·주의할 부분은 **굵게** 표시하세요. "
                "문단 사이에는 빈 줄을 넣으세요. "
                "설명이 여러 포인트로 나뉘면 불릿 목록(*)이나 번호 목록을 사용하되, 각 항목은 반드시 새로운 줄에서 시작하세요. "
                "불릿 목록을 시작하기 전에도 빈 줄을 넣으세요. "
                "핵심 주제가 바뀌면 빈 줄을 넣어 문단을 나누고, 한 문단을 너무 길게 이어 쓰지 마세요."
            )
        else:
            video_note = (
                "\n\n중요: 비디오 또는 오디오 파일이 포함되어 있습니다. "
                "답변은 반드시 전사 텍스트를 바탕으로 작성하고, 관련 내용이 있으면 답변 본문 안에 바로 타임스탬프를 넣어주세요. "
                "예: 드래그 장면(0:17): ..., 사라지는 장면(0:23-0:31): ... "
                "답변 아래에 별도의 '관련 시점' 목록은 만들지 마세요. "
                "타임스탬프를 고를 때는 어떤 용어를 소개하거나 이름을 붙이는 시점이 아니라, "
                "질문과 직접 관련된 내용이 실제로 설명되거나 시연되기 시작하는 지점을 기준으로 삼으세요. "
                "즉 제목을 붙이는 문장보다, 의미를 설명하는 문장, 행동이 드러나는 문장, 핵심 내용이 시작되는 문장을 우선하세요. "
                "질문이 어떤 개념, 장면, 행동, 차이, 이유, 방법, 흐름, 인물, 사건, 주제의 위치를 묻는 경우에도 같은 원칙을 적용하세요. "
                "질문과 직접 대응되는 설명 구간을 먼저 찾고 그 구간의 시간을 붙이세요. "
                "근거가 여러 개면 시간들을 나열만 하지 말고, 정말 이어지는 설명이면 범위(예: 0:35-0:59)로 묶어서 쓰세요. "
                "답변이 여러 포인트로 나뉘는 주제라면 짧은 도입 설명 뒤에 '-' 불릿 목록으로 나눠 설명해도 됩니다. "
                "예: 한 문장 요약 후, '* 항목명(시간): 설명' 형식으로 2~4개 핵심 포인트를 정리하세요. "
                "비교, 원인, 특징, 단계, 사례처럼 여러 요소를 구분해서 설명하는 질문에서는 불릿 형식을 우선적으로 사용하세요. "
                "추측하지 말고 전사에 근거가 있는 내용만 답하세요. "
                "답변은 Markdown 형식으로 작성하고, 핵심 개념·중요 결론·주의할 부분은 **굵게** 표시하세요. "
                "문단 사이에는 빈 줄을 넣으세요. "
                "설명이 여러 포인트로 나뉘면 불릿 목록(*)이나 번호 목록을 사용하되, 각 항목은 반드시 새로운 줄에서 시작하세요. "
                "불릿 목록을 시작하기 전에도 빈 줄을 넣으세요. "
                "핵심 주제가 바뀌면 빈 줄을 넣어 문단을 나누고, 한 문단을 너무 길게 이어 쓰지 마세요."
            )
        if joined_ref_lines:
            video_note += f"\n\n관련 전사 구간:\n{joined_ref_lines}"

    # ── 이미지 타입별 분석 가이드 ─────────────────
    IMAGE_TYPE_GUIDE = {
        "photo": "피사체의 외형·행동·감정을 묘사하고, 배경과 전체적인 분위기도 언급하세요.",
        "chart": "축 레이블, 데이터 값, 단위, 추세, 최댓값·최솟값을 정확히 언급하세요. 수치는 반드시 이미지에서 보이는 그대로 인용하세요.",
        "diagram": "각 구성 요소의 이름·역할과 요소 간 관계·흐름 방향을 순서대로 설명하세요.",
        "screenshot": "화면에 보이는 UI 요소, 코드, 오류 메시지, 설정 값을 그대로 인용하고 기술적으로 분석하세요.",
        "document": "문서에 적힌 텍스트를 정확히 인용하고, 문서의 구조(제목·항목·서명 등)를 설명하세요.",
        "illustration": "그림의 스타일, 색상, 표현하는 개념이나 메시지를 설명하세요.",
    }

    # ── 시스템 프롬프트 ───────────────────────────
    if has_images and not has_non_image:
        # 이미지 전용 모드 — 메타데이터로 subjects·type 주입
        image_names = ", ".join(f"'{img['filename']}'" for img in image_docs)

        # 모든 이미지의 subjects 통합
        all_subjects: list[str] = []
        all_types: list[str] = []
        for img in image_docs:
            meta = img_meta_map.get(img["doc_id"], {})
            all_subjects.extend(meta.get("subjects", []))
            if meta.get("image_type"):
                all_types.append(meta["image_type"])
            # visual_context를 context에 보조로 추가
        subjects_str = ", ".join(dict.fromkeys(all_subjects)) or "파악 중"
        primary_type = all_types[0] if all_types else "photo"
        type_guide = IMAGE_TYPE_GUIDE.get(primary_type, IMAGE_TYPE_GUIDE["photo"])

        # 관련성 사전 체크 (답변 생성과 완전히 분리)
        visual_context = ""
        for img in image_docs:
            meta = img_meta_map.get(img["doc_id"], {})
            if meta.get("visual_context"):
                visual_context = meta["visual_context"]
                break
        is_related = _check_image_relevance(question, all_subjects, visual_context)

        # 이미지 속 OCR 텍스트만 보조로 제공 (description은 모델이 이미지 직접 보므로 제외)
        ocr_parts: list[str] = []
        for img in image_docs:
            ocr = img_meta_map.get(img["doc_id"], {}).get("extracted_text", "").strip()
            if ocr:
                ocr_parts.append(ocr)
        ocr_text = "\n\n".join(ocr_parts)

        system_msg = f"""당신은 이미지를 깊이 분석하는 전문 비주얼 학습 튜터입니다.
첨부된 이미지({image_names})를 직접 보고 질문에 성실히 답변하세요.

【이미지 정보】
- 유형: {primary_type}
- 등장 대상: {subjects_str}

【{primary_type} 유형 분석 원칙】
{type_guide}
불명확한 부분은 "이미지에서 명확히 확인되지 않습니다"라고 솔직히 말하세요.

답변 시 {level_hint}"""
        if ocr_text:
            system_msg += f"\n\n【이미지 내 텍스트(OCR)】\n{ocr_text}"
        if not is_related:
            system_msg += '\n\n【출력 형식】답변 맨 앞에 반드시 이 문장을 먼저 쓰세요: "📌 이 질문은 업로드된 이미지와 직접적인 관련은 없지만, 알고 계시면 도움이 될 것 같아 답변드립니다."'

    elif has_images and has_non_image:
        # 혼합 모드 (이미지 + 다른 문서)
        names_str = ", ".join(f"'{n}'" for n in doc_filenames)
        system_msg = f"""당신은 학습 자료를 분석하는 AI 학습 코치입니다.
총 {len(doc_ids)}개의 자료({names_str})가 제공됩니다. 이미지는 직접 확인하고, 문서는 아래 컨텍스트를 참조하세요.
이미지와 문서 내용을 통합하여 종합적으로 답변하세요.
답변 시 {level_hint}{video_note}

<context>
{context}
</context>"""

    # 이미지 전용 케이스는 위에서 이미 system_msg 설정 완료 — 아래 블록 건너뜀
    if has_images and not has_non_image:
        pass

    elif is_multi:
        doc_names = doc_filenames
        names_str = ", ".join(f"'{n}'" for n in doc_names)
        system_msg = f"""당신은 학습 자료를 분석하는 AI 학습 코치입니다.
총 {len(doc_ids)}개의 문서({names_str})가 제공됩니다.
반드시 각 문서를 모두 참조하여 답변하세요.
'**[문서명]** 에서는 ~', '**[문서명]** 에 따르면 ~' 형식으로 각 문서의 내용을 명확히 구분하여 서술하세요.
답변 시 {level_hint}
제공된 문서 내용에서 최대한 찾아서 답변하세요.{video_note}

<context>
{context}
</context>"""

    elif is_ppt_doc:
        # PPT 전용 프롬프트
        ppt_name = doc_filenames[0] if doc_filenames else "슬라이드 자료"
        current_slide_hint = (
            f"\n【현재 슬라이드】\n학생이 현재 [슬라이드 {current_slide}]를 보고 있습니다. "
            f"질문이 특정 슬라이드를 명시하지 않는 경우, [슬라이드 {current_slide}]의 내용을 우선적으로 참조하여 답변하세요.\n"
        ) if current_slide else ""
        system_msg = f"""당신은 '{ppt_name}' PPT 자료를 바탕으로 학생의 질문에 깊이 있게 답변하는 AI 학습 튜터입니다.
{current_slide_hint}

【슬라이드 인용 규칙】
- 내용을 설명할 때는 반드시 해당 슬라이드 번호를 [슬라이드 N] 형식으로 명시하세요.
  예: "[슬라이드 3]에 따르면 ~", "이 개념은 [슬라이드 5]에서 확인할 수 있습니다."
- 여러 슬라이드에 걸친 내용이면 관련 슬라이드를 모두 명시하세요.
  예: "[슬라이드 2]와 [슬라이드 4]를 종합하면 ~"
- 슬라이드 번호는 반드시 [슬라이드 N] 형식을 정확히 지켜주세요 (클릭 가능한 링크로 변환됩니다).
- 단순 인용에 그치지 말고, 슬라이드 간 흐름과 연결 관계를 함께 설명하세요.

【시각 자료 및 영상 활용 규칙】
- context에 [시각 자료] 항목이 있으면 해당 이미지·차트·다이어그램의 내용을 답변에 반드시 포함하세요.
- context에 [동영상 음성 내용] 항목이 있으면 영상에서 설명하는 내용을 근거로 활용하세요.
- 텍스트와 시각 자료를 함께 종합하여 더 풍부하게 설명하세요.

【답변 원칙】
- PPT 전체 흐름에서 질문이 어느 위치에 해당하는지 간단히 짚고 답변하세요.
- 학생이 특정 슬라이드를 직접 언급하지 않는다면 PPT 전체 내용을 기반으로 답변하는 것을 우선하세요.
- 슬라이드 자료에 없는 내용은 "이 내용은 슬라이드에 없지만," 이라고 먼저 밝힌 뒤 일반 지식으로 보완하세요.
- 슬라이드 내용에만 그치지 말고, 개념의 의미와 실제 활용까지 연결해서 설명하세요.
- 설명이 길어질 경우 절대 중간에 끊지 마세요. 시작한 설명은 반드시 완결하세요.
- "이어서", "계속", "다음" 등의 요청이 오면 이전 대화에서 이미 설명한 내용은 반복하지 말고, 아직 다루지 않은 슬라이드나 개념부터 이어서 설명하세요.
- 전체 요약·정리를 요청받은 경우, 앞 슬라이드에만 치우치지 말고 처음부터 끝 슬라이드까지 균형 있게 커버하세요. 슬라이드별로 핵심 키워드 위주로 간결하게 정리하여 전체를 빠짐없이 다루는 것을 우선하세요.

【추천 질문】
- 답변을 마친 후, 이 대화의 흐름에서 학생이 자연스럽게 궁금해할 만한 질문 2~3개를 제안하세요.
- 반드시 아래 마커 형식으로 답변 맨 끝에 추가하세요 (답변 본문과 분리):
[SUGGESTED_QUESTIONS]
- (질문 1)
- (질문 2)
- (질문 3, 선택)
[/SUGGESTED_QUESTIONS]
- 질문은 지금 답변한 내용과 직접 이어지는 심화·응용 질문으로 구성하세요.

답변 시 {level_hint}

<context>
{context}
</context>"""

    elif is_pdf_doc:
        # PDF 전용 프롬프트
        pdf_name = doc_filenames[0] if doc_filenames else "PDF 문서"
        current_page_hint = (
            f"\n【현재 페이지】\n학생이 현재 [페이지 {current_slide}]를 보고 있습니다. "
            f"질문이 특정 페이지를 명시하지 않는 경우, [페이지 {current_slide}]의 내용을 우선적으로 참조하여 답변하세요.\n"
        ) if current_slide else ""
        system_msg = f"""당신은 '{pdf_name}' PDF 문서를 바탕으로 학생의 질문에 깊이 있게 답변하는 AI 학습 튜터입니다.
{current_page_hint}

【페이지 인용 규칙】
- 내용을 설명할 때는 반드시 해당 페이지 번호를 [페이지 N] 형식으로 명시하세요.
  예: "[페이지 3]에 따르면 ~", "이 내용은 [페이지 5]에서 확인할 수 있습니다."
- 여러 페이지에 걸친 내용이면 관련 페이지를 모두 명시하세요.
  예: "[페이지 2]와 [페이지 4]를 종합하면 ~"
- 페이지 번호는 반드시 [페이지 N] 형식을 정확히 지켜주세요 (클릭 가능한 링크로 변환됩니다).

【표 및 시각 자료 활용 규칙】
- context에 표(마크다운 형식)가 있으면 표의 수치와 항목을 정확히 인용하세요.
- 표의 열·행 관계를 분석해서 의미 있는 비교나 패턴을 설명하세요.
- context에 [시각 자료] 항목이 있으면 해당 이미지·차트·다이어그램의 내용을 답변에 반드시 포함하세요.
- [시각 자료] 안에 코드나 텍스트가 있으면 그 내용을 근거로 설명하세요.
- 시각 자료와 주변 텍스트를 연결해서 더 풍부하게 설명하세요.

【답변 원칙】
- 문서 전체 흐름에서 질문이 어느 위치에 해당하는지 간단히 짚고 답변하세요.
- 문서에 없는 내용은 "이 내용은 문서에 없지만," 이라고 먼저 밝힌 뒤 일반 지식으로 보완하세요.
- 설명이 길어질 경우 절대 중간에 끊지 마세요. 시작한 설명은 반드시 완결하세요.
- "이어서", "계속", "다음" 등의 요청이 오면 이미 설명한 내용은 반복하지 말고 아직 다루지 않은 부분부터 이어서 설명하세요.

【추천 질문 — 필수 출력】
답변 본문 작성 후, 반드시 아래 형식을 답변 맨 끝에 추가하세요. 생략하면 안 됩니다.
[SUGGESTED_QUESTIONS]
- (방금 답변한 내용과 직접 이어지는 심화 질문 1)
- (방금 답변한 내용과 직접 이어지는 심화 질문 2)
- (방금 답변한 내용과 직접 이어지는 심화 질문 3, 선택)
[/SUGGESTED_QUESTIONS]

답변 시 {level_hint}

<context>
{context}
</context>"""

    elif is_hwpx_doc:
        # HWPX 전용 프롬프트 — DOCX와 동일한 [페이지 N] 인용 형식
        hwpx_name = doc_filenames[0] if doc_filenames else "한글 문서"
        system_msg = f"""당신은 '{hwpx_name}' 문서를 바탕으로 학생의 질문에 깊이 있게 답변하는 AI 학습 튜터입니다.

【페이지 인용 규칙】
- 내용을 설명할 때는 반드시 해당 페이지 번호를 [페이지 N] 형식으로 명시하세요.
  예: "[페이지 3]에 따르면 ~", "이 내용은 [페이지 5]에서 확인할 수 있습니다."
- 여러 페이지에 걸친 내용이면 관련 페이지를 모두 명시하세요.
  예: "[페이지 2]와 [페이지 4]를 종합하면 ~"
- 페이지 번호는 반드시 [페이지 N] 형식을 정확히 지켜주세요 (클릭 가능한 링크로 변환됩니다).
- [1], [2] 같은 각주 번호는 절대 사용하지 마세요.

【답변 원칙】
- 문서 전체 흐름에서 질문이 어느 위치에 해당하는지 간단히 짚고 답변하세요.
- 문서에 없는 내용은 "이 내용은 문서에 없지만," 이라고 먼저 밝힌 뒤 일반 지식으로 보완하세요.
- 설명이 길어질 경우 절대 중간에 끊지 마세요. 시작한 설명은 반드시 완결하세요.
- "이어서", "계속", "다음" 등의 요청이 오면 이미 설명한 내용은 반복하지 말고 아직 다루지 않은 부분부터 이어서 설명하세요.
- 전체 요약·정리를 요청받은 경우, 앞 페이지에만 치우치지 말고 처음부터 끝 페이지까지 균형 있게 커버하세요.

【추천 질문 — 필수 출력】
답변 본문 작성 후, 반드시 아래 형식을 답변 맨 끝에 추가하세요. 생략하면 안 됩니다.
[SUGGESTED_QUESTIONS]
- (방금 답변한 내용과 직접 이어지는 심화 질문 1)
- (방금 답변한 내용과 직접 이어지는 심화 질문 2)
- (방금 답변한 내용과 직접 이어지는 심화 질문 3, 선택)
[/SUGGESTED_QUESTIONS]

답변 시 {level_hint}

<context>
{context}
</context>"""

    elif is_docx_doc:
        # DOCX 전용 프롬프트 — PDF와 동일한 [페이지 N] 인용 형식
        docx_name = doc_filenames[0] if doc_filenames else "Word 문서"
        system_msg = f"""당신은 '{docx_name}' 문서를 바탕으로 학생의 질문에 깊이 있게 답변하는 AI 학습 튜터입니다.

【페이지 인용 규칙】
- 내용을 설명할 때는 반드시 해당 페이지 번호를 [페이지 N] 형식으로 명시하세요.
  예: "[페이지 3]에 따르면 ~", "이 내용은 [페이지 5]에서 확인할 수 있습니다."
- 여러 페이지에 걸친 내용이면 관련 페이지를 모두 명시하세요.
  예: "[페이지 2]와 [페이지 4]를 종합하면 ~"
- 페이지 번호는 반드시 [페이지 N] 형식을 정확히 지켜주세요 (클릭 가능한 링크로 변환됩니다).
- [1], [2] 같은 각주 번호는 절대 사용하지 마세요.

【답변 원칙】
- 문서 전체 흐름에서 질문이 어느 위치에 해당하는지 간단히 짚고 답변하세요.
- 문서에 없는 내용은 "이 내용은 문서에 없지만," 이라고 먼저 밝힌 뒤 일반 지식으로 보완하세요.
- 설명이 길어질 경우 절대 중간에 끊지 마세요. 시작한 설명은 반드시 완결하세요.
- "이어서", "계속", "다음" 등의 요청이 오면 이미 설명한 내용은 반복하지 말고 아직 다루지 않은 부분부터 이어서 설명하세요.
- 전체 요약·정리를 요청받은 경우, 앞 페이지에만 치우치지 말고 처음부터 끝 페이지까지 균형 있게 커버하세요.

【추천 질문 — 필수 출력】
답변 본문 작성 후, 반드시 아래 형식을 답변 맨 끝에 추가하세요. 생략하면 안 됩니다.
[SUGGESTED_QUESTIONS]
- (방금 답변한 내용과 직접 이어지는 심화 질문 1)
- (방금 답변한 내용과 직접 이어지는 심화 질문 2)
- (방금 답변한 내용과 직접 이어지는 심화 질문 3, 선택)
[/SUGGESTED_QUESTIONS]

답변 시 {level_hint}

<context>
{context}
</context>"""

    elif is_url_doc:
        # URL 소스 전용 프롬프트
        url_sources = [
            str(doc_meta_map.get(d, {}).get("storage_path", ""))
            for d in doc_ids
            if str(doc_meta_map.get(d, {}).get("storage_path", "")).startswith("http")
        ]
        url_list = "\n".join(f"- {u}" for u in url_sources) if url_sources else ""
        system_msg = f"""당신은 아래 웹사이트에서 추출한 텍스트를 참고하여 질문에 답변하는 AI 어시스턴트입니다.

참조 웹사이트:
{url_list}

【답변 규칙】
1. <context>에 관련 내용이 있으면 웹사이트 내용을 우선적으로 활용하여 답변하세요.
2. <context>에 내용이 없거나 불충분한 경우, 알고 있는 지식으로 답변하되 반드시 "이 웹사이트에는 해당 내용이 없어 일반 지식을 바탕으로 답변합니다." 라고 먼저 밝히세요.
3. 웹사이트 내용과 일반 지식이 섞일 경우 각각 출처를 명확히 구분하세요.
4. 답변은 Markdown 형식으로 작성하고, 핵심 개념·중요 결론·주의할 부분은 **굵게** 표시하세요.
5. 문단 사이에는 빈 줄을 넣으세요.
6. 설명이 여러 포인트로 나뉘면 불릿 목록(*)이나 번호 목록을 사용하되, 각 항목은 반드시 새로운 줄에서 시작하세요.
7. 불릿 목록을 시작하기 전에도 빈 줄을 넣으세요.
8. 핵심 주제가 바뀌면 빈 줄을 넣어 문단을 나누고, 한 문단을 너무 길게 이어 쓰지 마세요.
답변 시 {level_hint}{video_note}

<context>
{context}
</context>"""

    else:
        system_msg = f"""당신은 학습 자료를 분석하는 AI 학습 코치입니다.
아래 문서 내용을 바탕으로 질문에 답변하세요.
답변 시 {level_hint}
제공된 문서 내용에서 최대한 찾아서 답변하세요. 정말로 알 수 없을 때만 "문서에서 찾을 수 없습니다"라고 하세요.{video_note}

<context>
{context}
</context>"""

    # 텍스트 출처가 있을 때만 인용 번호 지시 추가 (이미지 전용은 제외)
    if source_chunks and not (has_images and not has_non_image):
        system_msg += (
            "\n\n답변 시 문서 내용을 인용하는 부분에 [1], [2] 형식으로 출처 번호를 자연스럽게 붙이세요. "
            "예: '...JPA는 1차 캐시를 활용합니다 [1].' "
            "모든 문장에 필요하지 않고, 핵심 근거가 되는 내용에만 표시하세요."
        )

    # ── 메시지 구성 ───────────────────────────────
    messages: list[dict] = [{"role": "system", "content": system_msg}]

    if chat_history:
        for msg in chat_history[-6:]:
            # 히스토리에 이미지가 포함된 경우 텍스트만 유지 (토큰 절약)
            if isinstance(msg.get("content"), list):
                text_only = " ".join(
                    p.get("text", "") for p in msg["content"] if p.get("type") == "text"
                )
                messages.append({"role": msg["role"], "content": text_only})
            else:
                messages.append(msg)

    # PPT 현재 슬라이드 이미지 추가 (detail: low = 85 토큰)
    if is_ppt_doc and current_slide and not image_parts:
        ppt_doc_id = next(
            (d for d in doc_ids if str(doc_meta_map.get(d, {}).get("filename", "")).lower().rsplit(".", 1)[-1] in ("pptx", "ppt")),
            None,
        )
        if ppt_doc_id:
            slide_b64 = _get_slide_image_b64(ppt_doc_id, current_slide)
            if slide_b64:
                image_parts = [{
                    "type": "image_url",
                    "image_url": {"url": f"data:image/webp;base64,{slide_b64}", "detail": "low"},
                }]

    # 이미지가 있으면 multimodal user message
    if image_parts:
        user_content: list[dict] = [{"type": "text", "text": question}] + image_parts
        messages.append({"role": "user", "content": user_content})
    else:
        messages.append({"role": "user", "content": question})

    # 허용 모델 결정
    if _is_claude(model):
        safe_model = model  # Claude 계열은 그대로 허용 (vision 지원)
    elif model.startswith("gpt"):
        safe_model = model
        if image_parts and "gpt-4o" not in safe_model:
            safe_model = "gpt-4o"
    else:
        safe_model = "gpt-4o-mini"

    use_claude = _is_claude(safe_model)
    # youtube_only 재작성 등 팀원 코드는 OpenAI 클라이언트 유지
    client = OpenAI(api_key=settings.OPENAI_API_KEY)

    if stream:
        _stream_meta = {
            "source_chunks": source_chunks,
            "references": references,
            "doc_ids": doc_ids,
            "question": question,
            "asked_questions": asked_questions or [],
        }

        def _token_generator():
            import json as _json
            import threading

            # 스트리밍과 동시에 추천 질문 생성 (별도 스레드)
            _sugg_result: list = []
            def _gen_sugg():
                try:
                    _sugg_result.extend(generate_suggestions(
                        doc_ids=_stream_meta["doc_ids"],
                        asked_questions=_stream_meta["asked_questions"],
                        model=safe_model,
                    ))
                except Exception:
                    pass

            sugg_thread = threading.Thread(target=_gen_sugg, daemon=True)
            sugg_thread.start()

            full_answer_parts: list[str] = []
            _SUGG_MARKER = "[SUGGESTED_QUESTIONS]"
            _sugg_tail = ""   # 마커 시작 감지용 슬라이딩 버퍼
            _stop_stream = False

            def _yield_token(text: str):
                """[SUGGESTED_QUESTIONS] 마커 감지 시 스트리밍 중단."""
                nonlocal _sugg_tail, _stop_stream
                if _stop_stream:
                    return
                combined = _sugg_tail + text
                idx = combined.find(_SUGG_MARKER)
                if idx != -1:
                    _stop_stream = True
                    before = combined[:idx]
                    if before:
                        full_answer_parts.append(before)
                        yield f"data: {_json.dumps({'token': before})}\n\n"
                    return
                # 마커가 걸쳐있을 수 있는 접미사만 버퍼에 유지
                keep = len(_SUGG_MARKER) - 1
                if len(combined) > keep:
                    safe = combined[:-keep]
                    full_answer_parts.append(safe)
                    yield f"data: {_json.dumps({'token': safe})}\n\n"
                    _sugg_tail = combined[-keep:]
                else:
                    _sugg_tail = combined

            try:
                max_tok = 1500 if not is_full_doc_query else 3000
                if use_claude:
                    from anthropic import Anthropic as _Anthropic
                    _claude = _Anthropic(api_key=settings.ANTHROPIC_API_KEY)
                    _sys, _msgs = _convert_messages_for_anthropic(messages)
                    with _claude.messages.stream(
                        model=safe_model,
                        max_tokens=max_tok,
                        system=_sys,
                        messages=_msgs,  # type: ignore
                        temperature=0.3,
                    ) as stream:
                        for text in stream.text_stream:
                            yield from _yield_token(text)
                            if _stop_stream:
                                break
                else:
                    stream_response = client.chat.completions.create(
                        model=safe_model,
                        messages=messages,  # type: ignore
                        temperature=0.3,
                        max_tokens=max_tok,
                        stream=True,
                    )
                    for chunk in stream_response:
                        text = (chunk.choices[0].delta.content
                                if chunk.choices and chunk.choices[0].delta else None)
                        if text:
                            yield from _yield_token(text)
                            if _stop_stream:
                                break
                # 버퍼 잔량 flush
                if _sugg_tail and not _stop_stream:
                    full_answer_parts.append(_sugg_tail)
                    yield f"data: {_json.dumps({'token': _sugg_tail})}\n\n"
            except Exception as e:
                yield f"data: {_json.dumps({'error': str(e)})}\n\n"
                return

            # 유튜브 전용: 스트리밍 완료 후 타임스탬프 재작성
            full_answer = "".join(full_answer_parts)
            final_refs = _stream_meta["references"]
            if youtube_only and full_answer.strip():
                try:
                    aligned_rows = _get_semantic_rows(youtube_doc_ids, full_answer, top_k=max(TOP_K, 6))
                    aligned_references = _build_media_references(aligned_rows)
                    aligned_ref_lines = _format_reference_lines(aligned_references)
                    if aligned_ref_lines:
                        rewrite_msgs = [
                            {
                                "role": "system",
                                "content": (
                                    "너는 유튜브 자막을 근거로 답변을 다듬는 도우미다. "
                                    "초안 답변의 의미와 구조는 유지하되, 아래 자막 근거를 다시 읽고 "
                                    "질문과 직접 관련된 설명 부분의 시간만 답변 본문 안에 자연스럽게 붙여 다시 작성하라. "
                                    "별도의 관련 시점 목록은 만들지 말고, 이어지는 설명이면 범위 표기를 우선하라."
                                ),
                            },
                            {
                                "role": "user",
                                "content": (
                                    f"질문:\n{_stream_meta['question']}\n\n"
                                    f"초안 답변:\n{full_answer}\n\n"
                                    f"자막 근거:\n{aligned_ref_lines}"
                                ),
                            },
                        ]
                        if use_claude:
                            from anthropic import Anthropic as _Anthropic
                            _claude_rw = _Anthropic(api_key=settings.ANTHROPIC_API_KEY)
                            _sys_rw, _msgs_rw = _convert_messages_for_anthropic(rewrite_msgs)
                            _rw_resp = _claude_rw.messages.create(
                                model=safe_model,
                                max_tokens=2000,
                                system=_sys_rw,
                                messages=_msgs_rw,  # type: ignore
                                temperature=0.2,
                            )
                            rewritten = (_rw_resp.content[0].text if _rw_resp.content else "").strip()
                        else:
                            rewrite_resp = client.chat.completions.create(
                                model="gpt-4o-mini",
                                messages=rewrite_msgs,
                                temperature=0.2,
                                max_tokens=2000,
                            )
                            rewritten = (rewrite_resp.choices[0].message.content or "").strip()
                        if rewritten:
                            final_refs = aligned_references
                            yield f"data: {_json.dumps({'rewrite': rewritten})}\n\n"
                except Exception:
                    pass

            # 추천 질문 스레드 대기 (스트리밍 중 이미 완료됐을 가능성 높음)
            sugg_thread.join(timeout=5)

            sc = _stream_meta["source_chunks"]
            for c in sc:
                c.pop("full_text", None)
            srcs: list[Any] = sc if sc else _get_filenames(_stream_meta["doc_ids"])
            yield f"data: {_json.dumps({'done': True, 'sources': srcs, 'references': final_refs, 'suggestions': _sugg_result})}\n\n"

        return _token_generator(), [], []  # type: ignore

    max_tok = 1500 if not is_full_doc_query else 3000
    if use_claude:
        from anthropic import Anthropic as _Anthropic
        _claude = _Anthropic(api_key=settings.ANTHROPIC_API_KEY)
        _sys, _msgs = _convert_messages_for_anthropic(messages)
        _resp = _claude.messages.create(
            model=safe_model,
            max_tokens=max_tok,
            system=_sys,
            messages=_msgs,  # type: ignore
            temperature=0.3,
        )
        answer = _resp.content[0].text if _resp.content else ""
    else:
        response = client.chat.completions.create(
            model=safe_model,
            messages=messages,  # type: ignore
            temperature=0.3,
            max_tokens=max_tok,
        )
        answer = response.choices[0].message.content or ""

    if youtube_only and answer.strip():
        aligned_rows = _get_semantic_rows(youtube_doc_ids, answer, top_k=max(TOP_K, 6))
        aligned_references = _build_media_references(aligned_rows)
        aligned_ref_lines = _format_reference_lines(aligned_references)
        if aligned_ref_lines:
            rewrite_messages = [
                {
                    "role": "system",
                    "content": (
                        "너는 유튜브 자막을 근거로 답변을 다듬는 도우미다. "
                        "초안 답변의 의미와 구조는 유지하되, 아래 자막 근거를 다시 읽고 "
                        "질문과 직접 관련된 설명 부분의 시간만 답변 본문 안에 자연스럽게 붙여 다시 작성하라. "
                        "별도의 관련 시점 목록은 만들지 말고, 이어지는 설명이면 범위 표기를 우선하라."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        f"질문:\n{question}\n\n"
                        f"초안 답변:\n{answer}\n\n"
                        f"자막 근거:\n{aligned_ref_lines}"
                    ),
                },
            ]
            rewrite_response = client.chat.completions.create(
                model=safe_model,
                messages=rewrite_messages,
                temperature=0.2,
                max_tokens=2000,
            )
            rewritten_answer = rewrite_response.choices[0].message.content or ""
            if rewritten_answer.strip():
                answer = rewritten_answer
                references = aligned_references

    # ── 이미지 직접 참조 질문에서 잘못 붙은 📌 제거 ──
    if has_images and _IMAGE_DIRECT_REF_RE.search(question):
        if answer.lstrip().startswith("📌"):
            lines = answer.split("\n")
            filtered = []
            skip_blank = False
            for line in lines:
                if line.lstrip().startswith("📌 이 질문은 업로드된"):
                    skip_blank = True
                    continue
                if skip_blank and line.strip() == "":
                    skip_blank = False
                    continue
                filtered.append(line)
                skip_blank = False
            answer = "\n".join(filtered).lstrip()

    # ── 근거 구절 정밀 추출 (evidence extraction) ──
    if source_chunks:
        _extract_evidence_passages(answer, source_chunks, client, safe_model)
        # full_text 는 내부 처리 전용이므로 응답에서 제거
        for chunk in source_chunks:
            chunk.pop("full_text", None)

    # 추천 질문 파싱: [SUGGESTED_QUESTIONS]...[/SUGGESTED_QUESTIONS] 마커 추출 (PPT·PDF 공통)
    suggested_questions: list[str] = []
    if "[SUGGESTED_QUESTIONS]" in answer:
        try:
            sq_match = re.search(r"\[SUGGESTED_QUESTIONS\](.*?)\[/SUGGESTED_QUESTIONS\]", answer, re.DOTALL)
            if sq_match:
                sq_block = sq_match.group(1)
                for line in sq_block.strip().splitlines():
                    q = re.sub(r"^[-*\d.)\s]+", "", line).strip()
                    if q:
                        suggested_questions.append(q)
                # 답변 본문에서 마커 블록 제거
                answer = answer[:sq_match.start()].rstrip()
                print(f"[RAG] 추천 질문 파싱 완료: {suggested_questions}")
        except Exception as e:
            print(f"[RAG] 추천 질문 파싱 오류: {e}")

    sources: list[Any] = source_chunks if source_chunks else _get_filenames(doc_ids)
    return answer, sources, references, suggested_questions


# ─────────────────────────────────────────────
# 콘텐츠 생성
# ─────────────────────────────────────────────

def generate_content(
    doc_ids: list[str],
    gen_type: str,
    model: str = _DEFAULT_MODEL,
    level: str = "intermediate",
    quiz_count: int = 5,
    topic: str = "",
    difficulty: str = "intermediate",
) -> str:
    """요약 / 퀴즈 / 학습 계획 생성. 결과 문자열 반환."""
    context = _get_context(doc_ids, max_chars=10000)

    level_map = {"beginner": "입문", "intermediate": "중급", "advanced": "심화"}
    difficulty_map = {"easy": "쉬운", "intermediate": "중간", "hard": "어려운"}
    level_ko = level_map.get(level, "중급")
    difficulty_ko = difficulty_map.get(difficulty, "중간")

    if gen_type == "summary":
        prompt = f"""아래는 학습 문서의 내용입니다.
이 내용을 바탕으로 핵심 개념과 주요 내용을 {level_ko} 수준에 맞게 요약해주세요.
마크다운 형식으로 작성해주세요.

문서 내용:
{context}"""

    elif gen_type == "quiz":
        topic_instruction = (
            f"특히 '{topic}' 주제와 관련된 내용으로 퀴즈를 만들어주세요." if topic else ""
        )
        prompt = f"""아래는 학습 문서의 내용입니다.
이 내용을 바탕으로 {difficulty_ko} 난이도의 4지선다 퀴즈 {quiz_count}개를 만들어주세요.
{topic_instruction}

반드시 아래 JSON 형식으로만 응답하세요. 다른 텍스트는 절대 포함하지 마세요:

{{
  "title": "문서 핵심 내용을 반영한 구체적인 퀴즈 제목",
  "questions": [
    {{
      "id": 1,
      "question": "질문 내용",
      "options": ["선택지1", "선택지2", "선택지3", "선택지4"],
      "answerIndex": 0,
      "hint": "정답을 직접 언급하지 않고 방향만 제시하는 힌트",
      "explanation": "정답인 이유를 설명하는 해설"
    }}
  ]
}}

규칙:
- answerIndex는 정답 options의 0부터 시작하는 인덱스
- options는 반드시 4개
- hint는 정답을 직접 언급하지 말 것
- questions 배열에 정확히 {quiz_count}개의 항목 포함
- title은 문서 내용을 구체적으로 반영할 것

문서 내용:
{context}"""

    elif gen_type == "plan":
        prompt = f"""아래는 학습 문서의 내용입니다.
이 내용을 바탕으로 {level_ko} 학습자를 위한 주간 학습 계획을 작성해주세요.
마크다운 형식으로 Day별 목표와 학습 내용을 구체적으로 작성해주세요.

문서 내용:
{context}"""

    else:
        raise ValueError(f"Unknown gen_type: {gen_type}")

    return _call_llm(
        messages=[{"role": "user", "content": prompt}],
        model=model,
        temperature=0.4,
        json_mode=(gen_type == "quiz"),
    )


def generate_audio_overview(
    doc_ids: list[str],
    fmt: str = "deep_analysis",
    language: str = "ko",
    length: str = "default",
    focus: str = "",
    model: str = _DEFAULT_MODEL,
) -> tuple[bytes, str, str]:
    """2인 토크쇼 형식의 오디오 오버뷰 생성."""
    context = _get_context(doc_ids, max_chars=12000)

    length_map = {
        "short": "짧게 (10~14줄 대화)",
        "default": "충분히 (18~26줄 대화)",
    }
    length_desc = length_map.get(length, length_map["default"])

    format_map = {
        "deep_analysis": "두 호스트가 생동감 있게 주고받는 심층 분석 대화. 소스의 주제를 깊이 분석하고 서로 다른 관점을 연결합니다.",
        "summary": "소스의 핵심 아이디어를 빠르게 파악할 수 있도록 간결하게 요약하는 대화.",
        "critique": "소스에 대한 전문가적 시각의 비평 대화. 장단점과 개선점을 건설적으로 논의합니다.",
        "debate": "소스와 관련한 주제로 두 호스트가 서로 다른 입장을 가지고 논쟁하는 토론.",
    }
    format_desc = format_map.get(fmt, format_map["deep_analysis"])

    lang_map = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}
    lang_label = lang_map.get(language, "한국어")
    focus_instruction = f"\n특히 다음 부분에 집중해주세요: {focus}" if focus.strip() else ""

    prompt = f"""아래 문서 내용을 바탕으로 팟캐스트 스타일의 오디오 오버뷰 스크립트를 {lang_label}로 작성해주세요.

형식: {format_desc}
길이: {length_desc}{focus_instruction}

규칙:
- Host A와 Host B 두 사람이 자연스럽게 대화합니다
- 각 발화는 짧고 자연스럽게 (1~3문장)
- 인사 없이 바로 핵심 주제로 시작
- "그렇군요", "맞아요", "흥미롭네요" 같은 자연스러운 리액션 포함
- 청취자가 이해하기 쉽게 전문 용어는 설명 포함

문서 내용:
{context}

반드시 아래 JSON 형식으로만 응답하세요:
{{
  "title": "오디오 오버뷰 제목",
  "lines": [
    {{"speaker": "A", "text": "말할 내용"}},
    {{"speaker": "B", "text": "응답 내용"}}
  ]
}}"""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt}],
        model=model,
        temperature=0.7,
        json_mode=True,
    )
    result = json.loads(raw)
    title = result.get("title", "오디오 오버뷰")
    lines: list[dict] = result.get("lines", [])

    script = "\n".join(
        f"{'Host A' if l.get('speaker') == 'A' else 'Host B'}: {l.get('text', '')}"
        for l in lines
    )

    # Host A → sarah(여성), Host B → josh(남성) ElevenLabs 음성
    voice_map = {"A": "sarah", "B": "josh"}
    audio_parts: list[bytes] = []
    for line in lines:
        text = line.get("text", "").strip()
        if not text:
            continue
        voice_name = voice_map.get(line.get("speaker", "A"), "sarah")
        try:
            audio_parts.append(_elevenlabs_tts_simple(text, voice_name=voice_name))
        except Exception as e:
            print(f"[audio_overview] TTS 실패 (무시): {e}")

    audio_bytes = b"".join(audio_parts)
    return audio_bytes, script, title


def generate_mindmap(
    doc_ids: list[str],
    language: str = "ko",
    focus: str = "",
    model: str = _DEFAULT_MODEL,
) -> tuple[list, str]:
    """문서 내용을 평면 노드 배열 형식의 마인드맵 JSON으로 변환."""
    context = _get_context(doc_ids, max_chars=10000)

    lang_map = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}
    lang_label = lang_map.get(language, "한국어")
    focus_instruction = f"\n특히 다음 주제에 집중해주세요: {focus}" if focus.strip() else ""

    prompt = f"""아래 문서 내용을 바탕으로 학습용 마인드맵을 {lang_label}로 작성해주세요.{focus_instruction}

규칙:
- 루트 노드는 문서의 핵심 주제
- 2~4개의 1단계 브랜치 (주요 개념)
- 각 브랜치에 2~4개의 2단계 노드 (세부 내용)
- 필요시 3단계 노드 추가 가능
- 각 노드는 짧고 핵심적인 텍스트 (10단어 이내)

문서 내용:
{context}

반드시 아래 JSON 형식으로만 응답하세요:
{{
  "title": "마인드맵 제목",
  "nodes": [
    {{"id": "root", "text": "중심 주제"}},
    {{"id": "1", "text": "브랜치 1", "parent": "root"}},
    {{"id": "1-1", "text": "세부 내용 1", "parent": "1"}},
    {{"id": "1-2", "text": "세부 내용 2", "parent": "1"}},
    {{"id": "2", "text": "브랜치 2", "parent": "root"}},
    {{"id": "2-1", "text": "세부 내용 1", "parent": "2"}}
  ]
}}"""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt}],
        model=model,
        temperature=0.5,
        json_mode=True,
    )
    result = json.loads(raw)
    title = result.get("title", "마인드맵")
    nodes = result.get("nodes", [{"id": "root", "text": title}])
    return nodes, title


def generate_flashcards(
    doc_ids: list[str],
    count: str = "standard",
    difficulty: str = "intermediate",
    topic: str = "",
    language: str = "ko",
    model: str = _DEFAULT_MODEL,
) -> tuple[list, str]:
    """문서 내용을 바탕으로 플래시카드(앞면/뒷면) 배열을 생성.

    Returns:
        (cards_list, title)
        cards_list: [{"front": "용어/개념", "back": "설명", "hint": "힌트(선택)"}]
    """
    count_map = {"fewer": 10, "standard": 20, "more": 40}
    card_count = count_map.get(count, 20)

    difficulty_map = {"easy": "쉬운", "intermediate": "보통", "hard": "어려운"}
    difficulty_ko = difficulty_map.get(difficulty, "보통")

    lang_map = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}
    lang_label = lang_map.get(language, "한국어")

    topic_instruction = f"\n특히 '{topic}' 주제와 관련된 카드를 만들어주세요." if topic.strip() else ""
    context = _get_context(doc_ids, max_chars=10000)

    prompt = f"""아래 문서 내용을 바탕으로 {difficulty_ko} 난이도의 학습용 플래시카드 {card_count}개를 {lang_label}로 만들어주세요.{topic_instruction}

플래시카드 규칙:
- 앞면(front): 핵심 용어, 개념 또는 질문 (짧고 명확하게)
- 뒷면(back): 앞면에 대한 설명, 정의 또는 답변 (이해하기 쉽게)
- 힌트(hint): 정답을 직접 언급하지 않고 방향만 제시 (선택사항)
- 내용이 점점 어려워지도록 순서를 정렬해주세요

문서 내용:
{context}

반드시 아래 JSON 형식으로만 응답하세요:
{{
  "title": "플래시카드 세트 제목",
  "cards": [
    {{
      "front": "앞면 내용 (질문 또는 용어)",
      "back": "뒷면 내용 (답변 또는 설명)",
      "hint": "힌트 (선택, 없으면 빈 문자열)"
    }}
  ]
}}"""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt}],
        model=model,
        temperature=0.4,
        json_mode=True,
    )
    result = json.loads(raw)
    title = result.get("title", "플래시카드")
    cards = result.get("cards", [])
    return cards, title


def generate_slides(
    doc_ids: list[str],
    format: str = "presenter",
    length: str = "default",
    language: str = "ko",
    prompt: str = "",
    model: str = _DEFAULT_MODEL,
) -> tuple[list, str, str]:
    """문서 내용을 바탕으로 슬라이드 자료(JSON)를 생성.

    Returns:
        (slides_list, title, cover_image_b64)
        slides_list: [{"title": "슬라이드 제목", "bullets": ["내용1", "내용2"], "speaker_notes": "발표자 노트", "layout": "title|content|two_column"}]
        cover_image_b64: 표지 이미지 base64 문자열 (DALL-E 3)
    """
    length_map = {"short": (5, 8), "default": (8, 12), "long": (12, 18)}
    min_slides, max_slides = length_map.get(length, (8, 12))

    lang_map = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}
    lang_label = lang_map.get(language, "한국어")

    if format == "detailed":
        format_instruction = "자세한 자료: 전체 텍스트와 세부정보를 포함한 포괄적인 슬라이드로, 이메일로 공유하거나 단독으로 읽기에 적합합니다. 각 슬라이드에 충분한 텍스트를 포함해 주세요."
    else:
        format_instruction = "발표자 슬라이드: 발표 중 참고할 핵심 키워드와 간결한 bullet point만 포함한 시각적인 슬라이드입니다. 각 bullet은 5단어 이내로 짧게 작성해 주세요."

    custom_instruction = f"\n추가 요구사항: {prompt.strip()}" if prompt.strip() else ""
    context = _get_context(doc_ids, max_chars=12000)

    prompt_text = f"""아래 문서 내용을 바탕으로 {lang_label} 슬라이드 자료를 만들어주세요.

형식: {format_instruction}{custom_instruction}

슬라이드 수: {min_slides}~{max_slides}장 (첫 슬라이드는 표지, 마지막은 정리/결론)

규칙:
- 첫 슬라이드: layout = "title", bullets는 비우고 subtitle 필드에 부제목만 기입
- 중간 슬라이드: layout = "content" (단일 컬럼) 또는 "two_column" (양쪽 비교)
- 마지막 슬라이드: layout = "summary", 핵심 내용 3~5개 정리
- bullets: 각 슬라이드의 본문 내용 (문자열 배열)
- speaker_notes: 발표자가 구두로 설명할 내용 (발표자 슬라이드 형식일 때 더 자세히)

문서 내용:
{context}

반드시 아래 JSON 형식으로만 응답하세요:
{{
  "title": "슬라이드 전체 제목",
  "slides": [
    {{
      "title": "슬라이드 제목",
      "subtitle": "부제목 (표지에만 사용, 나머지는 빈 문자열)",
      "bullets": ["내용1", "내용2", "내용3"],
      "speaker_notes": "발표자 노트",
      "layout": "title"
    }}
  ]
}}"""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt_text}],
        model=model,
        temperature=0.5,
        json_mode=True,
    )
    result = json.loads(raw)
    title = result.get("title", "슬라이드 자료")
    slides = result.get("slides", [])

    # DALL-E 3 표지 이미지 생성 (OpenAI 전용 기능, 실패 시 빈 문자열)
    cover_image_b64 = ""
    try:
        oai_client = OpenAI(api_key=settings.OPENAI_API_KEY)
        img_prompt = _call_llm(
            messages=[{
                "role": "user",
                "content": f"""다음 슬라이드 제목을 위한 DALL-E 이미지 생성 프롬프트를 영어로 한 문장으로 작성해주세요.
제목: {title}

요구사항:
- 학술/교육 자료에 어울리는 전문적인 flat illustration 스타일
- 파란색 계열 색상 팔레트 (#1e3a5f, #2563eb 계열)
- 깔끔하고 미니멀한 디자인, 흰색 배경
- 텍스트 없이 아이콘/도형/다이어그램으로만 구성
- 16:9 비율에 맞는 가로형 구도
프롬프트만 출력하세요.""",
            }],
            model=model,
            max_tokens=200,
        )
        dalle_resp = oai_client.images.generate(
            model="dall-e-3",
            prompt=img_prompt,
            size="1792x1024",
            quality="standard",
            response_format="b64_json",
            n=1,
        )
        cover_image_b64 = dalle_resp.data[0].b64_json or ""
    except Exception as e:
        print(f"[slides] DALL-E image generation failed: {e}")

    return slides, title, cover_image_b64


def generate_report(
    doc_ids: list[str],
    format: str = "briefing",
    language: str = "ko",
    length: str = "default",
    tone: str = "formal",
    instructions: str = "",
    model: str = _DEFAULT_MODEL,
) -> tuple[list, str]:
    """문서 내용을 바탕으로 구조화된 보고서를 생성.

    Returns:
        (sections, title)
        sections: [{"heading": "섹션 제목", "content": "내용 텍스트"}]
    """
    lang_map = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}
    lang_label = lang_map.get(language, "한국어")

    tone_map = {
        "formal": "격식체(공식적이고 전문적인 문체)",
        "casual": "구어체(친근하고 읽기 쉬운 문체)",
        "academic": "학술체(논문 수준의 엄밀하고 인용 중심의 문체)",
    }
    tone_label = tone_map.get(tone, "격식체")

    length_map = {"short": 3, "default": 5, "long": 8}
    section_count = length_map.get(length, 5)

    format_prompts = {
        "briefing": f"""브리핑 문서: 핵심 인사이트를 간결하게 정리한 요약 보고서입니다.
구성: 개요, 주요 발견사항 {section_count-2}개, 결론
각 섹션은 bullet point 중심으로 간결하게 작성하세요.""",

        "study_guide": f"""학습 가이드: 학습자가 내용을 효과적으로 복습하고 이해할 수 있도록 돕는 구조입니다.
구성: 핵심 개념 요약, 단답형 퀴즈 (질문 {max(section_count-2, 3)}개와 정답), 추천 에세이 질문 2개, 핵심 용어집
학습자 친화적이고 교육적인 톤으로 작성하세요.""",

        "blog": f"""블로그 게시물: 일반 독자가 쉽게 이해할 수 있는 기사 형식입니다.
구성: 흥미로운 서론, 본론 {section_count-2}개의 소제목(각 2~3문단), 결론
스토리텔링 방식으로 읽기 쉽게 작성하세요.""",

        "prd": f"""제품 요구사항 정의서(PRD): 제품/서비스의 요구사항을 체계적으로 정리한 문서입니다.
구성: 개요 및 목적, 사용자 페르소나, 핵심 기능 요구사항, 기술적 제약 및 고려사항, 성공 지표
명확하고 구체적인 요구사항 형식으로 작성하세요.""",

        "architecture": f"""시스템 아키텍처 설계서: 시스템의 구조와 구성요소를 설명하는 기술 문서입니다.
구성: 시스템 개요, 핵심 컴포넌트 설명, 데이터 흐름, 기술 스택, 확장성/보안 고려사항
기술적이고 상세하게 작성하세요.""",

        "tech_explainer": f"""기술 개념 설명서: 복잡한 기술 개념을 이해하기 쉽게 설명하는 문서입니다.
구성: 개념 소개(비유 포함), 원리 설명, 실제 동작 방식, 활용 사례, 장단점 및 한계
단계적으로 개념을 쌓아가는 방식으로 설명하세요.""",

        "learning_guide": f"""학습 활용 가이드: 서비스나 도구를 효과적으로 활용하는 방법을 안내하는 입문용 자료입니다.
구성: 시작하기, 주요 기능 소개 (각 {max(section_count-3, 2)}가지), 활용 팁 및 모범 사례, 자주 묻는 질문
초보자를 위한 친절한 안내서 형식으로 작성하세요.""",

        "custom": "사용자가 아래 추가 지시사항에 명시한 형식으로 보고서를 작성하세요.",
    }

    format_instruction = format_prompts.get(format, format_prompts["briefing"])
    custom_instruction = f"\n\n추가 지시사항: {instructions.strip()}" if instructions.strip() else ""
    context = _get_context(doc_ids, max_chars=12000)

    prompt_text = f"""아래 문서 내용을 바탕으로 {lang_label} 보고서를 생성해주세요.

보고서 형식:
{format_instruction}{custom_instruction}

문체: {tone_label}

문서 내용:
{context}

반드시 아래 JSON 형식으로만 응답하세요:
{{
  "title": "보고서 제목",
  "sections": [
    {{
      "heading": "섹션 제목",
      "content": "섹션 내용 (마크다운 bullet 사용 가능. \\n으로 줄바꿈)"
    }}
  ]
}}"""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt_text}],
        model=model,
        temperature=0.5,
        json_mode=True,
    )
    result = json.loads(raw)
    title = result.get("title", "보고서")
    sections = result.get("sections", [])
    return sections, title


def generate_data_table(
    doc_ids: list[str],
    format: str = "summary_table",
    language: str = "ko",
    instructions: str = "",
    model: str = _DEFAULT_MODEL,
) -> tuple[str, str, list, list]:
    """문서 내용을 바탕으로 구조화된 데이터 표를 생성.

    Returns:
        (title, description, columns, rows)
        columns: [{"id": "col_id", "title": "Column Title", "type": "text|number|date"}]
        rows: [{"col_id": value, ...}, ...]
    """
    lang_map = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}
    lang_label = lang_map.get(language, "한국어")

    context = _get_context(doc_ids, max_chars=6000)
    custom_instruction = f"\n\n추가 지시사항: {instructions.strip()}" if instructions.strip() else ""

    # 형식별 상세 프롬프트 정의
    format_specs = {
        "summary_table": {
            "description": "핵심 내용 정리표",
            "instruction": """다음 구조의 요약 표를 생성하세요:
- 컬럼: 주제(text) | 핵심 내용(text) | 중요도(text: 상/중/하) | 관련 개념(text)
- 5~8개의 행: 문서의 핵심 내용을 주제별로 체계적으로 정리
- 각 행은 하나의 학습 포인트를 명확하게 표현

예시:
- 주제: "데이터베이스의 정의", 핵심 내용: "구조화된 데이터 모음", 중요도: "상", 관련 개념: "SQL, 스키마"
- 주제: "인덱싱", 핵심 내용: "검색 성능 향상 기법", 중요도: "중", 관련 개념: "B-tree, 쿼리"
""",
        },
        "comparison_table": {
            "description": "비교 분석 표",
            "instruction": """문서에서 비교할 2~4개의 개념/항목을 추출하여 다음 구조로 표를 생성하세요:
- 첫 번째 컬럼: 비교 항목(text)
- 나머지 컬럼: 각 개념/항목별 특성(text)
- 5~7개의 행: 명확한 비교 항목 (정의, 특징, 장점, 단점, 사용 사례 등)
- 각 셀에는 간결하고 구체적인 설명 작성

예시:
- 비교항목: "정의", [개념A]: "...", [개념B]: "..."
- 비교항목: "장점", [개념A]: "...", [개념B]: "..."
""",
        },
        "concept_definition": {
            "description": "개념 정의 표",
            "instruction": """문서의 핵심 개념/용어를 추출하여 다음 구조의 완전한 정의 표를 생성하세요:
- 컬럼: 개념(text) | 정의(text) | 예시(text) | 유사어/반대어(text) | 관련 분야(text)
- 8~12개의 행: 문서에서 중요한 개념들을 모두 포함
- 정의는 명확하고 이해하기 쉽게
- 예시는 실제로 발생할 수 있는 사례로
- 유사어와 반대어를 모두 포함

예시:
- 개념: "인덱싱", 정의: "데이터 검색 성능 향상 기법", 예시: "데이터베이스 컬럼에 인덱스 생성", 유사어: "색인화", 관련분야: "데이터베이스"
""",
        },
        "learning_checklist": {
            "description": "학습 점검표",
            "instruction": """학습자가 자신의 이해도를 점검할 수 있는 다음 구조의 체크리스트를 생성하세요:
- 컬럼: 학습 항목(text) | 상세 설명(text) | 중요도(text: 상/중/하) | 자가 평가 기준(text)
- 8~15개의 행: 문서의 학습 목표를 구체적인 항목으로 분해
- 자가 평가 기준: "설명할 수 있다", "예시를 들 수 있다", "응용할 수 있다" 등 명확한 기준 제시
- 중요도에 따라 학습 우선순위 명시

예시:
- 학습항목: "SELECT 문장", 상세설명: "데이터베이스에서 데이터를 조회하는 명령어", 중요도: "상", 평가기준: "SELECT * FROM 문법을 이해하고 사용할 수 있다"
""",
        },
        "progress_tracking": {
            "description": "진도 추적표",
            "instruction": """주차/단계별 학습 계획을 추적할 수 있는 다음 구조의 표를 생성하세요:
- 컬럼: 주차(text) | 학습 내용(text) | 예상 소요 시간(text) | 학습 목표(text) | 완료 여부(text: 미완료/진행중/완료)
- 6~10개의 행: 문서 내용을 주차/단계별로 분할하여 구성
- 소요 시간: 현실적인 학습 시간 제시 (예: "2시간", "1주일")
- 학습 목표: 각 주차의 구체적인 달성 목표 명시
- 완료 여부: 초기값은 모두 "미완료"로 설정

예시:
- 주차: "1주차", 내용: "데이터베이스 기초 개념", 소요시간: "3시간", 목표: "DBMS의 정의와 종류 이해", 완료여부: "미완료"
""",
        },
    }

    spec = format_specs.get(format, format_specs["summary_table"])
    format_instruction = spec["instruction"]

    prompt_text = f"""당신은 학습 자료를 체계적으로 정리하는 교육 전문가입니다.
아래 문서 내용을 바탕으로 {spec["description"]}을 생성해주세요.{custom_instruction}

생성 규칙:
{format_instruction}

문서 내용:
{context}

반드시 아래 JSON 형식으로만 응답하세요. 필드를 생략하지 마세요:
{{
  "title": "표의 제목 (문서 내용을 반영한 구체적인 제목)",
  "description": "표의 목적과 사용 방법을 설명하는 한두 문장",
  "columns": [
    {{
      "id": "col_1",
      "title": "첫 번째 열 제목",
      "type": "text"
    }},
    {{
      "id": "col_2",
      "title": "두 번째 열 제목",
      "type": "text"
    }}
  ],
  "rows": [
    {{
      "col_1": "첫 번째 셀 값",
      "col_2": "두 번째 셀 값"
    }},
    {{
      "col_1": "첫 번째 셀 값",
      "col_2": "두 번째 셀 값"
    }}
  ]
}}"""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt_text}],
        model=model,
        temperature=0.7,
        json_mode=True,
    )
    result = json.loads(raw)
    title = result.get("title", "데이터표")
    description = result.get("description", "")
    columns = result.get("columns", [])
    rows = result.get("rows", [])
    return title, description, columns, rows


# ─────────────────────────────────────────────
# 영상 생성 (Remotion + TTS)
# ─────────────────────────────────────────────

def _extract_marker_timings(notes_raw: str, n_bullets: int) -> tuple[str, list[float]]:
    """[N] 마커를 파싱하고 문자 위치 기반 타이밍을 계산한다.

    핵심 원리:
    - TTS 발화 시간은 텍스트 길이에 비례 → 마커 이전 문자 수 / 전체 문자 수 = 타이밍 비율
    - 스프링 애니메이션 지연을 보정하기 위해 전체 길이의 6%를 앞서 트리거

    Returns: (마커 제거된 clean text, 타이밍 리스트 0.0~1.0)
    """
    pattern = re.compile(r'\[(\d+)\]')

    clean = pattern.sub('', notes_raw)
    clean = re.sub(r'  +', ' ', clean).strip()
    total_len = max(1, len(clean))

    if n_bullets <= 0:
        return clean, []

    marker_list = [(m.start(), int(m.group(1)), len(m.group(0))) for m in pattern.finditer(notes_raw)]

    timing_map: dict[int, float] = {}
    removed_chars = 0
    for raw_pos, num, mlen in marker_list:
        clean_pos = raw_pos - removed_chars
        frac = clean_pos / total_len
        frac = max(0.03, frac - 0.06)
        timing_map[num] = round(min(0.90, frac), 3)
        removed_chars += mlen

    if not timing_map:
        return clean, [round((i + 1) / (n_bullets + 1), 3) for i in range(n_bullets)]

    timings: list[float] = []
    sorted_map = sorted(timing_map.items())
    for i in range(n_bullets):
        if (i + 1) in timing_map:
            timings.append(timing_map[i + 1])
        else:
            prev = max((t for k, t in sorted_map if k <= i),     default=0.03)
            nxt  = min((t for k, t in sorted_map if k >= i + 2), default=0.90)
            timings.append(round((prev + nxt) / 2, 3))

    return clean, timings


def generate_video(
    doc_ids: list[str],
    language: str = "ko",
    length: str = "default",
    model: str = _DEFAULT_MODEL,
) -> tuple[list[dict], str]:
    """문서 내용을 바탕으로 Remotion 비디오용 슬라이드 데이터(TTS 오디오 포함)를 생성."""
    import base64 as _base64

    length_map = {"short": (4, 5), "default": (5, 7), "long": (7, 10)}
    min_slides, max_slides = length_map.get(length, (5, 7))

    lang_map = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}
    lang_label = lang_map.get(language, "한국어")

    context = _sanitize(_get_context(doc_ids, max_chars=14000))

    prompt_text = f"""당신은 대학 강의를 진행하는 강사입니다.
아래 문서를 바탕으로, 실제 강의실에서 학생들에게 설명하듯 {lang_label} 강의 슬라이드 스크립트를 작성해주세요.

슬라이드 수: {min_slides}~{max_slides}장 (첫 장 타이틀, 마지막 장 정리)

━━━ 핵심 원칙 ━━━
각 슬라이드는 하나의 주제만 다룹니다.
강의 전체가 자연스럽게 이어지도록, 앞 내용과 연결하거나 흐름을 만들어주세요.

━━━ speaker_notes 작성 규칙 ━━━
어조: 강사가 학생들에게 말하는 자연스러운 구어체 강의 말투
  - "~습니다", "~입니다" 로 끝나는 부드럽고 친근한 경어체
  - 딱딱하게 정의를 나열하지 말고, 강사가 실제로 입으로 말하듯 풀어서 설명
  - 청중이 따라오기 쉽도록 짧고 명확하게, 그러나 자연스러운 호흡으로
  - "~라고 볼 수 있습니다", "~인 셈입니다", "쉽게 말하면" 같은 표현도 자연스럽게 사용 가능

문장 수: 정확히 1 + bullets 개수 문장
  - 1번째 문장: 이 슬라이드 주제를 강사답게 도입 (앞 내용과 연결하거나, 왜 중요한지 한 문장으로)
  - 2번째~ 문장: bullets[0], bullets[1]... 순서대로 각각 1문장씩, 학생 눈높이에 맞게 설명

슬라이드 흐름: 강의 전체가 하나의 이야기처럼 이어지도록 작성합니다.
  - 앞 슬라이드 내용을 짧게 받아서 이 슬라이드로 자연스럽게 넘어오는 도입도 좋습니다.
  - 단, 이 슬라이드의 본론(bullets 설명)이 흐려지지 않도록 도입은 1문장 이내로 짧게.
연결어 적극 활용: "또한", "그리고", "반면에", "특히", "따라서", "이를 통해" 등으로 문장 흐름을 자연스럽게 연결
[N] 마커: 각 bullet 설명 문장 맨 앞에 [1], [2], [3]... 삽입 (렌더링 시 자동 제거)

예시 (bullets 3개 → 정확히 4문장):
"TCP가 왜 중요한지 이해하려면, 우선 데이터가 어떻게 전달되는지를 알아야 합니다. [1]연결을 시작할 때는 3-way handshake라는 과정을 거쳐서, 서로 준비됐는지 확인하고 통신을 시작합니다. [2]전송 중 패킷이 사라지면 자동으로 다시 보내줘서, 데이터가 빠짐없이 도착하도록 보장합니다. [3]통신이 끝날 때도 4-way handshake로 안전하게 연결을 끊어줍니다."

━━━ layout 선택 ━━━
"steps"  → 순서·절차·단계 (3~5개)
"cards"  → 독립된 개념·기능 (3~4개)
"table"  → 비교·대조
"list"   → 그 외 일반 내용

━━━ 나머지 필드 ━━━
[title] 슬라이드 주제 (10자 이내, 명사형)
[summary] 이 슬라이드의 핵심 결론 한 줄 (25자 이내, "~이다" 형식)
[bullets] "핵심어 | 한 줄 설명" 형식, 3~4개
  - 핵심어: 2~5자의 키워드
  - 설명: speaker_notes 해당 문장과 동일한 내용을 압축한 한 문장

━━━━━━━━━━━━━━━━━━

문서 내용:
{context}

반드시 아래 JSON 형식으로만 응답하세요:
{{
  "title": "동영상 제목",
  "slides": [
    {{
      "layout": "list",
      "title": "슬라이드 제목",
      "summary": "핵심 결론 한 줄",
      "bullets": ["핵심어 | 설명", "핵심어2 | 설명2"],
      "speaker_notes": "핵심 도입 문장. [1]bullet1 설명 문장. [2]bullet2 설명 문장."
    }}
  ]
}}"""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt_text}],
        model=model,
        temperature=0.6,
        max_tokens=4000,
        json_mode=True,
    )
    result = json.loads(raw)
    video_title = result.get("title", "동영상 개요")
    raw_slides = result.get("slides", [])

    FPS = 30
    MIN_SLIDE_FRAMES = 150
    CHARS_PER_SEC = 5.0

    slides_with_audio: list[dict] = []
    for slide in raw_slides:
        notes_raw = slide.get("speaker_notes", "").strip()
        bullets   = slide.get("bullets", [])

        notes, bullet_timings = _extract_marker_timings(notes_raw, len(bullets))
        notes = _sanitize(notes)

        try:
            tts_audio = _elevenlabs_tts_simple(notes or _sanitize(slide.get("title", "")), voice_name="sarah")
            audio_b64 = _base64.b64encode(tts_audio).decode()
            audio_sec = max(4.0, len(notes) / CHARS_PER_SEC)
            estimated_sec = audio_sec + 1.0
        except Exception as e:
            print(f"[video] TTS 실패, 무음으로 대체: {e}")
            audio_b64 = ""
            audio_sec = max(4.0, len(notes) / CHARS_PER_SEC)
            estimated_sec = audio_sec + 1.0

        duration_frames = max(MIN_SLIDE_FRAMES, int(estimated_sec * FPS))

        slides_with_audio.append({
            "title": slide.get("title", ""),
            "layout": slide.get("layout", "list"),
            "summary": slide.get("summary", ""),
            "bullets": bullets,
            "bulletTimings": bullet_timings,
            "audioSec": round(audio_sec, 2),
            "speakerNotes": notes,
            "audioBase64": audio_b64,
            "durationInFrames": duration_frames,
        })

    return slides_with_audio, video_title


def generate_infographic(
    doc_ids: list[str],
    format: str = "overview",
    language: str = "ko",
    instructions: str = "",
    model: str = _DEFAULT_MODEL,
) -> tuple[str, str, list[dict]]:
    """문서 내용을 바탕으로 인포그래픽을 생성.

    반환: (title, description, sections)
    - sections: 각각 title, icon, color, stat, stat_label, points를 가진 섹션 객체 리스트
    """
    context = _get_context(doc_ids, max_chars=6000)

    colors = ["#FF6B6B", "#4ECDC4", "#45B7D1", "#FFA07A", "#98D8C8", "#F7DC6F", "#BB8FCE"]

    format_specs = {
        "overview": {
            "label": "개요형",
            "instruction": "문서의 주요 내용을 5개의 주제별 섹션으로 요약합니다. 각 섹션은 제목, 이모지 아이콘, 색상, 핵심 포인트 3~4개를 포함합니다.",
            "section_count": 5,
            "stat_rule": "stat은 선택사항입니다. 문서에 명확한 수치나 통계가 있을 때만 포함하세요. 없으면 stat과 stat_label을 생략하세요.",
        },
        "process": {
            "label": "프로세스형",
            "instruction": "문서에서 설명하는 단계적 프로세스나 절차를 순서대로 4~5개의 섹션으로 나타냅니다.",
            "section_count": 5,
            "stat_rule": "stat은 필수입니다. 각 섹션마다 '1단계', '2단계', '3단계' 등의 형식으로 순서 번호를 포함하세요. stat_label에는 해당 단계의 핵심 행동을 한 줄로 적으세요.",
        },
        "comparison": {
            "label": "비교형",
            "instruction": "문서의 주요 개념이나 대상들의 특징과 차이를 비교 분석합니다. 4개의 섹션으로 나누어 각각의 장점, 특징, 사용 시기를 설명합니다.",
            "section_count": 4,
            "stat_rule": "stat은 선택사항입니다. 비교 대상의 수치 차이나 비율 데이터가 있을 때만 포함하세요. 명확한 수치가 없으면 생략해도 괜찮습니다.",
        },
        "statistics": {
            "label": "통계형",
            "instruction": "문서에서 강조하는 수치, 통계, 비율을 중심으로 4~5개의 섹션을 만듭니다.",
            "section_count": 5,
            "stat_rule": "stat은 필수입니다. 각 섹션마다 문서에서 가장 임팩트 있는 숫자나 비율을 포함하세요 (예: '85%', '3배', '1억명'). stat_label에는 그 수치의 맥락을 설명하세요. 모든 섹션에 stat이 있어야 합니다.",
        },
        "timeline": {
            "label": "타임라인형",
            "instruction": "문서가 다루는 역사적, 발전적, 또는 시간 순서의 흐름을 4~5개의 섹션으로 표현합니다.",
            "section_count": 5,
            "stat_rule": "stat은 필수입니다. 각 섹션의 시간을 명확히 표시하세요 (예: '2024년', '1990년대', '3월'). stat_label에는 그 시기의 핵심 사건이나 변화를 적으세요.",
        },
    }

    spec = format_specs.get(format, format_specs["overview"])
    language_label = {"ko": "한국어", "en": "English", "ja": "日本語", "zh": "中文"}.get(language, "한국어")
    color_palette = ", ".join(colors)

    prompt_text = f"""당신은 전문적인 교육 콘텐츠 크리에이터입니다.

다음 문서 내용을 바탕으로 '{spec["label"]}' 형태의 인포그래픽을 생성하세요.

【문서 내용】
{context}

【요청사항】
- 형식: {spec["label"]}
- {spec["instruction"]}
- 섹션 개수: 약 {spec["section_count"]}개
- 언어: {language_label}
{f'- 추가 지시사항: {instructions}' if instructions else ''}

【stat 필드 규칙】
{spec["stat_rule"]}

【JSON 응답 형식】
각 섹션은 다음 구조를 가져야 합니다:
- title: 섹션 제목 (문자열)
- icon: 이모지 아이콘 (1개 이모지 문자)
- color: 섹션 배경색 (hex 컬러 코드, 권장 팔레트: {color_palette})
- stat: 위의 stat 필드 규칙에 따라 포함 또는 생략
- stat_label: stat이 있을 때만 포함 (문자열)
- points: 핵심 포인트 배열 (3~5개의 문자열)

전체 응답:
{{
  "title": "전체 제목",
  "description": "간단한 설명 (1~2문장)",
  "sections": [
    {{
      "title": "섹션 제목",
      "icon": "🎯",
      "color": "#FF6B6B",
      "stat": "85%",
      "stat_label": "통계 설명",
      "points": ["포인트 1", "포인트 2", "포인트 3"]
    }}
  ]
}}

올바른 JSON 형식으로만 응답하세요."""

    raw = _call_llm(
        messages=[{"role": "user", "content": prompt_text}],
        model=model,
        temperature=0.7,
        json_mode=True,
    )
    result = json.loads(raw)
    title = result.get("title", "인포그래픽")
    description = result.get("description", "")
    sections = result.get("sections", [])

    for idx, section in enumerate(sections):
        if not section.get("color"):
            section["color"] = colors[idx % len(colors)]

    return title, description, sections


# ─────────────────────────────────────────────
# 추천 질문 생성
# ─────────────────────────────────────────────

def generate_suggestions(
    doc_ids: list[str],
    asked_questions: list[str] = [],
    model: str = _DEFAULT_MODEL,
    last_answer: str = "",
) -> list[dict]:
    """문서 타입과 내용 기반으로 카테고리별 추천 질문 3개를 생성한다.
    이미지 소스는 실제 이미지를 Vision API에 직접 전달.
    반환: [{"text": "...", "category": "이해|분석|적용"}]
    """
    import base64

    filenames = _get_filenames(doc_ids)
    image_docs = _get_image_docs(doc_ids)
    non_image_ids = [d for d in doc_ids if d not in {img["doc_id"] for img in image_docs}]

    # 파일 타입 감지
    file_types: set[str] = set()
    for fn in filenames:
        ext = fn.lower().rsplit(".", 1)[-1] if "." in fn else ""
        if ext in IMAGE_EXTENSIONS:
            file_types.add("image")
        elif ext in VIDEO_AUDIO_EXTENSIONS:
            file_types.add("video_audio")
        elif ext in ("pptx", "ppt"):
            file_types.add("ppt")
        else:
            file_types.add("document")
    if any(_is_youtube_url(fn) for fn in filenames):
        file_types.add("video_audio")

    # 텍스트 컨텍스트 (비이미지 문서)
    context_ids = non_image_ids if non_image_ids else doc_ids
    context = _get_context(context_ids, max_chars=4000) if context_ids else ""

    # 이미지 메타데이터 조회
    img_meta_map = _get_image_metadata([img["doc_id"] for img in image_docs])
    all_subjects: list[str] = []
    all_img_types: list[str] = []
    for img in image_docs:
        meta = img_meta_map.get(img["doc_id"], {})
        all_subjects.extend(meta.get("subjects", []))
        if meta.get("image_type"):
            all_img_types.append(meta["image_type"])
    subjects_str = ", ".join(dict.fromkeys(all_subjects))
    primary_img_type = all_img_types[0] if all_img_types else "photo"

    # 이미지 base64 수집 (최대 1568px 리사이즈 + JPEG 압축 후 전송)
    image_parts: list[dict] = []
    for img in image_docs:
        raw = _download_image_b64(img["storage_path"])
        if raw:
            b64, mime = _compress_image_to_base64(raw, max_width=1568, quality=85)
            image_parts.append({
                "type": "image_url",
                "image_url": {"url": f"data:{mime};base64,{b64}", "detail": "high"},
            })

    asked_block = ""
    if asked_questions:
        asked_block = "\n\n[이미 질문된 항목 - 반드시 제외]\n" + "\n".join(
            f"- {q}" for q in asked_questions[-6:]
        )

    # 이미지 전용 프롬프트
    if image_parts and not non_image_ids:
        image_names = ", ".join(f"'{img['filename']}'" for img in image_docs)
        # 이미지 타입별 질문 가이드
        img_type_hint = {
            "photo":        "피사체([{subjects}])의 외형·행동·습성·생태 등 다양한 각도의 질문 포함.",
            "chart":        "축 값, 데이터 수치, 추세, 비교 포인트를 구체적으로 언급하는 질문 포함.",
            "diagram":      "구성 요소 간 관계, 흐름 순서, 핵심 노드에 관한 질문 포함.",
            "screenshot":   "코드 로직, UI 요소, 오류 원인, 설정 의미에 관한 기술적 질문 포함.",
            "document":     "문서에 적힌 특정 항목, 조항, 수치를 직접 인용하는 질문 포함.",
            "illustration": "그림이 표현하는 개념, 메시지, 상징적 의미에 관한 질문 포함.",
        }.get(primary_img_type, "이미지에서 보이는 구체적 요소를 언급하는 질문 포함.").format(subjects=subjects_str)

        meta_hint = f"\n[사전 분석 결과] 이미지 유형: {primary_img_type} / 등장 대상: {subjects_str}" if subjects_str else ""

        prompt_text = f"""당신은 이미지 기반 학습을 돕는 전문 비주얼 튜터입니다.
첨부된 이미지({image_names})를 직접 분석하여, 학습자가 이 이미지를 완전히 이해하기 위해 반드시 물어봐야 할 핵심 질문 3개를 생성하세요.
{meta_hint}
{asked_block}

[질문 생성 규칙]
{img_type_hint}

1. 이미지에서 실제로 보이는 구체적 요소(수치, 대상, 색상, 구조, 텍스트)를 질문에 직접 언급할 것
   ✗ "이미지가 무엇을 나타내나요?" (너무 일반적)
   ✓ "그래프에서 2023년 매출이 급등한 원인은 무엇인가요?" (구체적)

2. 카테고리별 1개씩:
   - "이해": 이미지의 특정 요소·대상·수치를 확인하는 질문
   - "분석": 이미지 전체 의미, 패턴, 시사점을 파악하는 질문
   - "적용": 이미지 내용을 실제 상황이나 다른 개념에 연결하는 질문

3. 각 질문은 완전한 문장, 40자 이내

JSON 형식으로만 응답:
{{"questions": [{{"text": "질문", "category": "이해"}}, {{"text": "질문", "category": "분석"}}, {{"text": "질문", "category": "적용"}}]}}"""

        user_content: list[dict] = [{"type": "text", "text": prompt_text}] + image_parts

    else:
        # 텍스트 문서 (+ 선택적으로 이미지 혼합) 프롬프트
        type_guidance = ""
        if "image" in file_types:
            type_guidance += "\n[이미지 포함] 이미지에서 보이는 시각적 요소, 차트 수치, 다이어그램을 직접 언급하는 질문 포함."
        if "video_audio" in file_types:
            type_guidance += "\n[영상/음성 포함] 영상에서 언급된 발언, 수치, 설명된 개념에 관한 질문 포함."
        if "ppt" in file_types:
            type_guidance += "\n[PPT 포함] 슬라이드 흐름, 핵심 주장, 데이터에 관한 질문 포함."

        answer_block = f"\n\n[방금 AI가 답변한 내용 — 이 내용을 기반으로 질문을 생성하세요]\n{last_answer[:2000]}" if last_answer else ""
        context_block = f"\n══════════════════════════════\n[문서 내용]\n{context[:4000]}\n══════════════════════════════" if not last_answer else ""

        prompt_text = f"""당신은 학습 내용을 깊이 이해하도록 돕는 전문 튜터입니다.
{"방금 AI가 답변한 내용을 바탕으로, 학습자가 자연스럽게 이어서 궁금해할 심화 질문 3개를 생성하세요." if last_answer else "아래 문서 내용을 분석하여, 학습자가 반드시 짚고 넘어가야 할 핵심 질문 3개를 생성하세요."}
{type_guidance}{answer_block}{context_block}
{asked_block}

[생성 규칙 - 반드시 준수]
1. 각 질문은 문서에 실제로 등장하는 용어, 수치, 개념을 반드시 포함할 것
   ✗ "핵심 개념이 무엇인가요?" → ✓ "SHA-256이 MD5보다 안전한 이유는?"

2. 카테고리별 1개씩:
   - "이해": 개념·사실 확인 → "~란?", "~는 어떻게?"
   - "분석": 비교·원인·관계 → "~와 ~의 차이는?", "왜 ~?"
   - "적용": 실제 활용·확장 → "~을 어떻게 활용?", "만약 ~라면?"

3. 각 질문은 완전한 문장, 40자 이내

JSON 형식으로만 응답:
{{"questions": [{{"text": "질문", "category": "이해"}}, {{"text": "질문", "category": "분석"}}, {{"text": "질문", "category": "적용"}}]}}"""

        if image_parts:
            user_content = [{"type": "text", "text": prompt_text}] + image_parts
        else:
            user_content = [{"type": "text", "text": prompt_text}]

    try:
        _CLAUDE_MODELS = ("claude-opus", "claude-sonnet", "claude-haiku")
        is_claude = any(model.startswith(m) for m in _CLAUDE_MODELS)

        if is_claude and settings.ANTHROPIC_API_KEY:
            # Claude 모델: 이미지는 base64 형식 그대로 전달, JSON 출력 유도
            import anthropic
            claude_client = anthropic.Anthropic(api_key=settings.ANTHROPIC_API_KEY)

            # Anthropic 이미지 블록으로 변환
            anthropic_content: list[dict] = []
            for part in user_content if isinstance(user_content, list) else [{"type": "text", "text": user_content}]:
                if part.get("type") == "image_url":
                    url: str = part["image_url"]["url"]
                    if url.startswith("data:"):
                        mime, b64data = url[5:].split(";base64,", 1)
                        anthropic_content.append({
                            "type": "image",
                            "source": {"type": "base64", "media_type": mime, "data": b64data},
                        })
                else:
                    anthropic_content.append({"type": "text", "text": part.get("text", "")})

            claude_resp = claude_client.messages.create(
                model=model,
                max_tokens=512,
                system="당신은 학습 내용을 깊이 이해하도록 돕는 전문 교육 튜터입니다. 자료의 실제 내용에 기반한 구체적이고 통찰력 있는 질문을 생성합니다. 반드시 JSON 형식으로만 응답하세요.",
                messages=[{"role": "user", "content": anthropic_content}],
            )
            raw = claude_resp.content[0].text if claude_resp.content else "{}"
            # JSON 블록 추출 (```json ... ``` 감싸진 경우 대비)
            json_match = re.search(r'\{[\s\S]*\}', raw)
            result = json.loads(json_match.group(0) if json_match else raw)
        else:
            raw = _call_llm(
                messages=[
                    {
                        "role": "system",
                        "content": "당신은 학습 내용을 깊이 이해하도록 돕는 전문 교육 튜터입니다. 자료의 실제 내용에 기반한 구체적이고 통찰력 있는 질문을 생성합니다.",
                    },
                    {"role": "user", "content": user_content},  # type: ignore
                ],
                model=model,
                temperature=0.75,
                json_mode=True,
            )
            result = json.loads(raw)

        questions = result.get("questions", [])
        if isinstance(questions, list):
            return [q for q in questions if isinstance(q, dict) and "text" in q][:3]
        return []
    except Exception:
        return []
